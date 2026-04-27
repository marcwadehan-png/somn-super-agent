"""
__all__ = [
    'parse_excel',
    'parse_powerpoint',
    'parse_word',
]

Office文档解析模块 (Word/Excel/PowerPoint)
"""

from pathlib import Path
from typing import Dict
import zipfile
import re
from ._dv_types import DocumentContent, DocumentMetadata, DocumentType

def parse_word(viewer, path: Path, options: Dict) -> DocumentContent:
    """解析Word文档"""
    extract_tables = options.get('extract_tables', True)
    
    content = DocumentContent(metadata=None)
    
    try:
        from docx import Document
        
        doc = Document(str(path))
        
        # 提取段落
        for para in doc.paragraphs:
            if para.text.strip():
                content.text_content += para.text + "\n"
        
        # 提取表格
        if extract_tables:
            for table_idx, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                content.tables.append(table_data)
        
        # 提取图片
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                content.images.append({
                    'name': rel.target_ref,
                    'type': rel.reltype.split('/')[-1]
                })
        
        content.metadata = DocumentMetadata(
            file_path=str(path),
            file_name=path.name,
            file_size=path.stat().st_size,
            file_type=DocumentType.WORD,
            mime_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            page_count=len(doc.paragraphs) // 30 + 1  # 估算
        )
        
    except ImportError:
        content.warnings.append("未安装python-docx库")
        content = _parse_docx_fallback(path)
    except Exception as e:
        content.warnings.append(f"Word解析错误: {e}")
    
    return content

def _parse_docx_fallback(path: Path) -> DocumentContent:
    """使用zipfile解析docx(无python-docx时)"""
    content = DocumentContent(metadata=None)
    
    try:
        with zipfile.ZipFile(path, 'r') as zf:
            # 读取document.xml
            xml_content = zf.read('word/document.xml').decode('utf-8')
            content.text_content = _extract_text_from_xml(xml_content)
    except Exception as e:
        content.warnings.append(f"Docx回退解析失败: {e}")
    
    return content

def _extract_text_from_xml(xml_content: str) -> str:
    """从XML中提取文本"""
    # 移除XML标签,保留文本
    text = re.sub(r'<[^>]+>', ' ', xml_content)
    # 清理多余空格
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def parse_excel(viewer, path: Path, options: Dict) -> DocumentContent:
    """解析Excel文档"""
    content = DocumentContent(metadata=None)
    
    try:
        import pandas as pd
        
        # 根据扩展名选择读取方式
        extension = path.suffix.lower()
        
        if extension == '.csv':
            df = pd.read_csv(path, encoding='utf-8', on_bad_lines='skip')
        else:
            # xlsx, xls
            df = pd.read_excel(path, sheet_name=None)  # 读取所有sheet
        
        if isinstance(df, dict):
            # 多sheet
            for sheet_name, sheet_df in df.items():
                content.text_content += f"\n[Sheet: {sheet_name}]\n"
                content.text_content += sheet_df.to_string() + "\n"
                content.tables.append(sheet_df.values.tolist())
        else:
            content.text_content = df.to_string()
            content.tables.append(df.values.tolist())
        
    except ImportError:
        content.warnings.append("未安装pandas库,尝试基础解析")
        content = _parse_excel_basic(path)
    except Exception as e:
        content.warnings.append(f"Excel解析错误: {e}")
    
    return content

def _parse_excel_basic(path: Path) -> DocumentContent:
    """基础Excel解析(无pandas时)"""
    from loguru import logger
    
    content = DocumentContent(metadata=None)
    
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True)
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            content.text_content += f"\n[Sheet: {sheet_name}]\n"
            
            rows = []
            for row in sheet.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else '' for cell in row]
                rows.append(row_data)
                content.text_content += '\t'.join(row_data) + '\n'
            
            content.tables.append(rows)
        
    except ImportError:
        content.warnings.append("未安装openpyxl库")
    except Exception as e:
        content.warnings.append(f"Excel基础解析失败: {e}")
    
    return content

def parse_powerpoint(viewer, path: Path, options: Dict) -> DocumentContent:
    """解析PowerPoint文档"""
    content = DocumentContent(metadata=None)
    
    try:
        from pptx import Presentation
        
        prs = Presentation(str(path))
        slide_count = len(prs.slides)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = f"\n[第{slide_num}页]\n"
            
            # 提取标题
            if slide.shapes.title:
                slide_text += f"标题: {slide.shapes.title.text}\n"
            
            # 提取文本框
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    if not (slide.shapes.title and shape == slide.shapes.title):
                        slide_text += shape.text + "\n"
            
            content.pages.append(slide_text)
            content.text_content += slide_text
            
            # 提取图片
            for shape in slide.shapes:
                if hasattr(shape, 'image'):
                    content.images.append({
                        'slide': slide_num,
                        'type': shape.shape_type
                    })
        
        # 更新元数据
        content.metadata = DocumentMetadata(
            file_path=str(path),
            file_name=path.name,
            file_size=path.stat().st_size,
            file_type=DocumentType.POWERPOINT,
            mime_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            page_count=slide_count
        )
        
    except ImportError:
        content.warnings.append("未安装python-pptx库,请运行: pip install python-pptx")
    except Exception as e:
        content.warnings.append(f"PowerPoint解析错误: {e}")
    
    return content
