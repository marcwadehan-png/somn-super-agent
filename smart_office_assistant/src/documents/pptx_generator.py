"""
__all__ = [
    'add_bullet_slide',
    'add_image_slide',
    'add_section_divider',
    'add_slide',
    'add_table_slide',
    'add_title_slide',
    'create_business_presentation',
    'create_presentation_from_outline',
    'create_training_presentation',
    'save',
]

PPT 演示文稿generate器 - PPTX Generator
基于 python-pptx 的演示文稿generate功能
"""

from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple, Union
from dataclasses import dataclass
from enum import Enum

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap

from loguru import logger

class SlideLayout(Enum):
    """幻灯片布局类型"""
    TITLE = "title"
    TITLE_AND_CONTENT = "title_and_content"
    BLANK = "blank"
    TITLE_ONLY = "title_only"
    TWO_CONTENT = "two_content"
    COMPARISON = "comparison"
    CONTENT_WITH_CAPTION = "content_with_caption"

@dataclass
class SlideContent:
    """幻灯片内容"""
    title: str
    content: Union[str, List[str], Dict[str, Any]]
    layout: SlideLayout = SlideLayout.TITLE_AND_CONTENT
    notes: Optional[str] = None

@dataclass
class ChartData:
    """图表数据"""
    chart_type: str  # 'bar', 'line', 'pie', etc.
    categories: List[str]
    series: List[Dict[str, Any]]  # [{'name': 'Series1', 'values': [1, 2, 3]}]
    title: Optional[str] = None

class PPTXGenerator:
    """
    PPT 演示文稿generate器
    
    功能:
    - 创建和编辑 PPT 演示文稿
    - 支持多种幻灯片布局
    - 插入图表和图片
    - 设置主题和样式
    - 添加演讲者备注
    """
    
    def __init__(self, template_path: Optional[str] = None):
        """
        initgenerate器
        
        Args:
            template_path: 可选的模板文件路径
        """
        if template_path and Path(template_path).exists():
            self.prs = Presentation(template_path)
            logger.info(f"加载模板: {template_path}")
        else:
            self.prs = Presentation()
            self._setup_default_theme()
        
        self.template_path = template_path
        
        # 布局mapping
        self._layout_map = {
            SlideLayout.TITLE: 0,
            SlideLayout.TITLE_AND_CONTENT: 1,
            SlideLayout.BLANK: 6,
            SlideLayout.TITLE_ONLY: 5,
            SlideLayout.TWO_CONTENT: 3,
            SlideLayout.COMPARISON: 4,
            SlideLayout.CONTENT_WITH_CAPTION: 2
        }
    
    def _setup_default_theme(self):
        """设置默认主题"""
        # 设置幻灯片尺寸 (16:9)
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
    
    def add_slide(
        self,
        content: SlideContent
    ) -> Any:
        """
        添加幻灯片
        
        Args:
            content: 幻灯片内容
        
        Returns:
            幻灯片对象
        """
        # get布局
        layout_idx = self._layout_map.get(content.layout, 1)
        slide_layout = self.prs.slide_layouts[layout_idx]
        
        # 添加幻灯片
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 设置标题
        if slide.shapes.title:
            title = slide.shapes.title
            title.text = content.title
            
            # 设置标题格式
            for paragraph in title.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = 'Microsoft YaHei'
                    run.font.size = Pt(32)
                    run.font.bold = True
        
        # 设置内容
        if content.layout == SlideLayout.TITLE_AND_CONTENT:
            self._set_content_placeholder(slide, content.content)
        elif content.layout == SlideLayout.TITLE:
            pass  # 标题幻灯片只有标题
        elif content.layout == SlideLayout.TWO_CONTENT:
            self._set_two_content(slide, content.content)
        
        # 添加备注
        if content.notes:
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = content.notes
        
        return slide
    
    def _set_content_placeholder(self, slide: Any, content: Union[str, List[str]]):
        """设置内容占位符"""
        if len(slide.placeholders) < 2:
            return
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        if isinstance(content, str):
            p = tf.paragraphs[0]
            p.text = content
            p.level = 0
            for run in p.runs:
                run.font.name = 'Microsoft YaHei'
                run.font.size = Pt(18)
        
        elif isinstance(content, list):
            for i, item in enumerate(content):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = item
                p.level = 0
                for run in p.runs:
                    run.font.name = 'Microsoft YaHei'
                    run.font.size = Pt(18)
    
    def _set_two_content(self, slide: Any, content: Dict[str, Any]):
        """设置双栏内容"""
        # 左栏
        if len(slide.placeholders) > 1:
            left = slide.placeholders[1]
            left_tf = left.text_frame
            left_tf.clear()
            
            left_content = content.get('left', [])
            if isinstance(left_content, list):
                for i, item in enumerate(left_content):
                    if i == 0:
                        p = left_tf.paragraphs[0]
                    else:
                        p = left_tf.add_paragraph()
                    p.text = item
                    p.level = 0
                    for run in p.runs:
                        run.font.name = 'Microsoft YaHei'
        
        # 右栏
        if len(slide.placeholders) > 2:
            right = slide.placeholders[2]
            right_tf = right.text_frame
            right_tf.clear()
            
            right_content = content.get('right', [])
            if isinstance(right_content, list):
                for i, item in enumerate(right_content):
                    if i == 0:
                        p = right_tf.paragraphs[0]
                    else:
                        p = right_tf.add_paragraph()
                    p.text = item
                    p.level = 0
                    for run in p.runs:
                        run.font.name = 'Microsoft YaHei'
    
    def add_title_slide(
        self,
        title: str,
        subtitle: str = "",
        presenter: str = ""
    ) -> Any:
        """
        添加标题幻灯片
        
        Args:
            title: 主标题
            subtitle: 副标题
            presenter: 演讲者
        """
        slide_layout = self.prs.slide_layouts[0]  # 标题幻灯片布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 设置标题
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = title
            
            for paragraph in title_shape.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = 'Microsoft YaHei'
                    run.font.size = Pt(44)
                    run.font.bold = True
        
        # 设置副标题
        if len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_text = subtitle
            if presenter:
                subtitle_text += f"\n\n演讲者: {presenter}"
            
            subtitle_shape.text = subtitle_text
            
            for paragraph in subtitle_shape.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = 'Microsoft YaHei'
                    run.font.size = Pt(24)
        
        return slide
    
    def add_section_divider(self, title: str, subtitle: str = "") -> Any:
        """
        添加章节分隔页
        
        Args:
            title: 章节标题
            subtitle: 章节副标题
        """
        content = {
            'left': [title],
            'right': [subtitle] if subtitle else []
        }
        
        slide_content = SlideContent(
            title="",
            content=content,
            layout=SlideLayout.TWO_CONTENT
        )
        
        return self.add_slide(slide_content)
    
    def add_bullet_slide(
        self,
        title: str,
        bullets: List[str],
        notes: str = ""
    ) -> Any:
        """
        添加项目符号列表幻灯片
        
        Args:
            title: 标题
            bullets: 列表项
            notes: 演讲者备注
        """
        slide_content = SlideContent(
            title=title,
            content=bullets,
            layout=SlideLayout.TITLE_AND_CONTENT,
            notes=notes
        )
        
        return self.add_slide(slide_content)
    
    def add_image_slide(
        self,
        title: str,
        image_path: str,
        caption: str = "",
        notes: str = ""
    ) -> Any:
        """
        添加图片幻灯片
        
        Args:
            title: 标题
            image_path: 图片路径
            caption: 图片说明
            notes: 演讲者备注
        """
        if not Path(image_path).exists():
            logger.warning(f"图片不存在: {image_path}")
            return None
        
        slide_content = SlideContent(
            title=title,
            content=caption,
            layout=SlideLayout.TITLE_AND_CONTENT,
            notes=notes
        )
        
        slide = self.add_slide(slide_content)
        if slide is None:
            logger.error("幻灯片创建失败")
            return None
        
        # 添加图片
        left = Inches(1)
        top = Inches(2)
        height = Inches(4.5)
        
        slide.shapes.add_picture(image_path, left, top, height=height)
        
        return slide
    
    def add_table_slide(
        self,
        title: str,
        headers: List[str],
        rows: List[List[str]],
        notes: str = ""
    ) -> Any:
        """
        添加表格幻灯片
        
        Args:
            title: 标题
            headers: 表头
            rows: 数据行
            notes: 演讲者备注
        """
        slide_content = SlideContent(
            title=title,
            content="",
            layout=SlideLayout.TITLE_AND_CONTENT,
            notes=notes
        )
        
        slide = self.add_slide(slide_content)
        
        # 添加表格
        rows_count = len(rows) + 1  # +1 for header
        cols_count = len(headers)
        
        left = Inches(1)
        top = Inches(2)
        width = Inches(11)
        height = Inches(0.5 * rows_count)
        
        table = slide.shapes.add_table(
            rows_count, cols_count, left, top, width, height
        ).table
        
        # 设置表头
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.name = 'Microsoft YaHei'
                    run.font.size = Pt(14)
        
        # 设置数据
        for row_idx, row_data in enumerate(rows, start=1):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_data)
                
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Microsoft YaHei'
                        run.font.size = Pt(12)
        
        return slide
    
    def create_presentation_from_outline(
        self,
        title: str,
        subtitle: str,
        outline: List[Dict[str, Any]],
        presenter: str = ""
    ):
        """
        根据大纲创建完整演示文稿
        
        Args:
            title: 演示文稿标题
            subtitle: 副标题
            outline: 大纲内容
            presenter: 演讲者
        """
        # 标题页
        self.add_title_slide(title, subtitle, presenter)
        
        # 目录页
        toc_items = [item.get('title', '') for item in outline]
        self.add_bullet_slide("目录", toc_items)
        
        # 内容页
        for item in outline:
            slide_title = item.get('title', '')
            content = item.get('content', [])
            notes = item.get('notes', '')
            
            if isinstance(content, list):
                self.add_bullet_slide(slide_title, content, notes)
            else:
                self.add_bullet_slide(slide_title, [content], notes)
        
        # 结束页
        self.add_title_slide("谢谢", "Q & A", presenter)
    
    def save(self, output_path: str) -> str:
        """
        保存演示文稿
        
        Args:
            output_path: 输出路径
        
        Returns:
            保存的文件路径
        """
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        self.prs.save(output_path)
        logger.info(f"演示文稿已保存: {output_path}")
        
        return str(output_path)

# 预设模板
class PresentationTemplates:
    """演示文稿模板"""
    
    @staticmethod
    def create_business_presentation(
        company_name: str,
        presentation_title: str,
        key_points: List[str],
        data_slides: List[Dict[str, Any]]
    ) -> PPTXGenerator:
        """创建商业演示模板"""
        gen = PPTXGenerator()
        
        # 标题页
        gen.add_title_slide(
            presentation_title,
            f"{company_name} 商业计划",
            ""
        )
        
        # 要点页
        gen.add_bullet_slide("核心要点", key_points)
        
        # 数据页
        for data_slide in data_slides:
            gen.add_bullet_slide(
                data_slide.get('title', ''),
                data_slide.get('points', [])
            )
        
        # 总结页
        gen.add_title_slide("感谢聆听", "期待合作")
        
        return gen
    
    @staticmethod
    def create_training_presentation(
        course_title: str,
        instructor: str,
        modules: List[Dict[str, Any]]
    ) -> PPTXGenerator:
        """创建培训课件模板"""
        gen = PPTXGenerator()
        
        # 标题页
        gen.add_title_slide(course_title, "培训课程", instructor)
        
        # 课程大纲
        outline = [m.get('title', '') for m in modules]
        gen.add_bullet_slide("课程大纲", outline)
        
        # 各模块内容
        for module in modules:
            # 模块分隔页
            gen.add_section_divider(module.get('title', ''))
            
            # 模块内容
            for slide in module.get('slides', []):
                gen.add_bullet_slide(
                    slide.get('title', ''),
                    slide.get('content', [])
                )
        
        # 结束页
        gen.add_title_slide("课程结束", "谢谢参与", instructor)
        
        return gen
