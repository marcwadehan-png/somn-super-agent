# -*- coding: utf-8 -*-
"""
文档生成集成测试 [v22.3]
Document Generation Integration Tests

测试文档模块与底层库的完整集成。
需要安装: python-pptx, python-docx, openpyxl, pandas

版本: v22.3.0
日期: 2026-04-25
"""

from __future__ import annotations

import sys
import os
from pathlib import Path
from typing import Dict, Any, List, Optional

# 添加项目路径
PROJECT_ROOT = Path(__file__).parent.parent.parent
sys.path.insert(0, str(PROJECT_ROOT / 'smart_office_assistant' / 'src'))


class DependencyChecker:
    """依赖检查器"""
    
    DEPENDENCIES = {
        'python-pptx': {
            'package': 'pptx',
            'import_name': 'pptx',
            'purpose': 'PPT生成',
            'required': True
        },
        'python-docx': {
            'package': 'python-docx',
            'import_name': 'docx',
            'purpose': 'Word生成',
            'required': True
        },
        'openpyxl': {
            'package': 'openpyxl',
            'import_name': 'openpyxl',
            'purpose': 'Excel生成',
            'required': True
        },
        'pandas': {
            'package': 'pandas',
            'import_name': 'pandas',
            'purpose': '数据分析',
            'required': False
        },
        'Pillow': {
            'package': 'Pillow',
            'import_name': 'PIL',
            'purpose': '图片处理',
            'required': False
        },
        'matplotlib': {
            'package': 'matplotlib',
            'import_name': 'matplotlib',
            'purpose': '图表生成',
            'required': False
        },
        'numpy': {
            'package': 'numpy',
            'import_name': 'numpy',
            'purpose': '数值计算',
            'required': False
        },
    }
    
    @classmethod
    def check_all(cls) -> Dict[str, Dict[str, Any]]:
        """检查所有依赖"""
        results = {}
        
        for name, info in cls.DEPENDENCIES.items():
            try:
                __import__(info['import_name'])
                results[name] = {
                    'installed': True,
                    'import_name': info['import_name'],
                    'purpose': info['purpose'],
                    'required': info['required']
                }
            except ImportError:
                results[name] = {
                    'installed': False,
                    'package': info['package'],
                    'import_name': info['import_name'],
                    'purpose': info['purpose'],
                    'required': info['required']
                }
        
        return results
    
    @classmethod
    def print_report(cls) -> None:
        """打印依赖报告"""
        print("\n" + "=" * 60)
        print("文档模块依赖检查报告")
        print("=" * 60)
        
        results = cls.check_all()
        
        installed = []
        missing = []
        
        for name, info in results.items():
            if info['installed']:
                installed.append(name)
                status = f"[{chr(0x2713)}] {name}"
            else:
                missing.append(name)
                status = f"[{chr(0x2717)}] {name}"
            
            purpose = info.get('purpose', '')
            pkg = info.get('package', info['import_name'])
            
            if not info['installed']:
                status += f" (安装: pip install {pkg})"
            
            print(f"  {status}")
        
        print("\n" + "-" * 60)
        print(f"已安装: {len(installed)}/{len(results)}")
        
        if missing:
            required = [m for m in missing if cls.DEPENDENCIES[m]['required']]
            optional = [m for m in missing if not cls.DEPENDENCIES[m]['required']]
            
            print(f"缺失必需: {len(required)}")
            print(f"缺失可选: {len(optional)}")
            
            if required:
                print("\n请安装必需依赖:")
                for pkg in required:
                    print(f"  pip install {cls.DEPENDENCIES[pkg]['package']}")
        
        return results


class DocumentGeneratorTester:
    """文档生成测试器"""
    
    def __init__(self, output_dir: Optional[Path] = None):
        self.output_dir = output_dir or Path(__file__).parent / 'test_output'
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.results = []
    
    def test_pptx_generation(self) -> bool:
        """测试PPT生成"""
        print("\n" + "=" * 60)
        print("测试: PPT生成")
        print("=" * 60)
        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RgbColor
        except ImportError:
            print("  [SKIP] python-pptx 未安装")
            return False
        
        try:
            # 创建演示文稿
            prs = Presentation()
            
            # 添加标题幻灯片
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = title_slide.shapes.title
            subtitle = title_slide.placeholders[1]
            title.text = "测试演示文稿"
            subtitle.text = "v22.3 集成测试"
            
            # 添加内容幻灯片
            content_slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = content_slide.shapes.title
            body = content_slide.placeholders[1]
            title.text = "测试内容"
            
            # 添加文本框
            text_frame = body.text_frame
            text_frame.text = "这是一个测试段落"
            p = text_frame.add_paragraph()
            p.text = "这是第二个段落"
            
            # 保存
            output_path = self.output_dir / 'test_presentation.pptx'
            prs.save(str(output_path))
            
            print(f"  [PASS] PPT生成成功: {output_path}")
            print(f"  文件大小: {output_path.stat().st_size} bytes")
            self.results.append(('PPT', True, str(output_path)))
            return True
            
        except Exception as e:
            print(f"  [FAIL] PPT生成失败: {e}")
            self.results.append(('PPT', False, str(e)))
            return False
    
    def test_docx_generation(self) -> bool:
        """测试Word生成"""
        print("\n" + "=" * 60)
        print("测试: Word生成")
        print("=" * 60)
        
        try:
            from docx import Document
            from docx.shared import Inches, Pt, RGBColor
        except ImportError:
            print("  [SKIP] python-docx 未安装")
            return False
        
        try:
            # 创建文档
            doc = Document()
            
            # 添加标题
            doc.add_heading('测试文档', 0)
            
            # 添加段落
            doc.add_paragraph('这是一个测试段落')
            
            # 添加加粗文本
            para = doc.add_paragraph()
            para.add_run('加粗文本').bold = True
            para.add_run(' 和普通文本')
            
            # 添加列表
            doc.add_heading('测试列表', 1)
            doc.add_paragraph('列表项1', style='List Bullet')
            doc.add_paragraph('列表项2', style='List Bullet')
            
            # 添加表格
            doc.add_heading('测试表格', 1)
            table = doc.add_table(rows=3, cols=3)
            table.style = 'Light Grid Accent 1'
            
            # 填充表格
            headers = table.rows[0].cells
            headers[0].text = '名称'
            headers[1].text = '数量'
            headers[2].text = '价格'
            
            for i in range(1, 3):
                row = table.rows[i].cells
                row[0].text = f'产品{i}'
                row[1].text = str(i * 10)
                row[2].text = f'{i * 100}元'
            
            # 保存
            output_path = self.output_dir / 'test_document.docx'
            doc.save(str(output_path))
            
            print(f"  [PASS] Word生成成功: {output_path}")
            print(f"  文件大小: {output_path.stat().st_size} bytes")
            self.results.append(('Word', True, str(output_path)))
            return True
            
        except Exception as e:
            print(f"  [FAIL] Word生成失败: {e}")
            self.results.append(('Word', False, str(e)))
            return False
    
    def test_xlsx_generation(self) -> bool:
        """测试Excel生成"""
        print("\n" + "=" * 60)
        print("测试: Excel生成")
        print("=" * 60)
        
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment
            from openpyxl.chart import BarChart, Reference
        except ImportError:
            print("  [SKIP] openpyxl 未安装")
            return False
        
        try:
            # 创建工作簿
            wb = Workbook()
            
            # 第一个工作表 - 数据表
            ws = wb.active
            ws.title = "销售数据"
            
            # 添加标题
            ws['A1'] = '产品名称'
            ws['B1'] = '销量'
            ws['C1'] = '收入'
            
            # 标题样式
            for cell in ws[1]:
                cell.font = Font(bold=True)
                cell.fill = PatternFill(start_color='CCCCCC', end_color='CCCCCC', fill_type='solid')
            
            # 添加数据
            data = [
                ('产品A', 100, 5000),
                ('产品B', 150, 7500),
                ('产品C', 80, 4000),
                ('产品D', 200, 10000),
            ]
            
            for i, (name, qty, revenue) in enumerate(data, start=2):
                ws[f'A{i}'] = name
                ws[f'B{i}'] = qty
                ws[f'C{i}'] = revenue
            
            # 添加汇总
            ws['A6'] = '合计'
            ws['B6'] = '=SUM(B2:B5)'
            ws['C6'] = '=SUM(C2:C5)'
            ws['A6'].font = Font(bold=True)
            
            # 第二个工作表 - 图表
            chart_ws = wb.create_sheet("图表")
            chart_ws['A1'] = '月份'
            chart_ws['B1'] = '销售额'
            months = ['1月', '2月', '3月', '4月']
            values = [1000, 1500, 1200, 1800]
            for i, (m, v) in enumerate(zip(months, values), start=2):
                chart_ws[f'A{i}'] = m
                chart_ws[f'B{i}'] = v
            
            # 创建图表
            chart = BarChart()
            chart.title = "月度销售额"
            data_ref = Reference(chart_ws, min_col=2, min_row=1, max_row=5)
            cats = Reference(chart_ws, min_col=1, min_row=2, max_row=5)
            chart.add_data(data_ref, titles_from_data=True)
            chart.set_categories(cats)
            chart_ws.add_chart(chart, "E2")
            
            # 保存
            output_path = self.output_dir / 'test_workbook.xlsx'
            wb.save(str(output_path))
            
            print(f"  [PASS] Excel生成成功: {output_path}")
            print(f"  文件大小: {output_path.stat().st_size} bytes")
            self.results.append(('Excel', True, str(output_path)))
            return True
            
        except Exception as e:
            print(f"  [FAIL] Excel生成失败: {e}")
            self.results.append(('Excel', False, str(e)))
            return False
    
    def test_template_integration(self) -> bool:
        """测试模板填充集成"""
        print("\n" + "=" * 60)
        print("测试: 模板填充集成")
        print("=" * 60)
        
        try:
            from documents.template_filler import TemplateFiller, quick_fill
        except ImportError:
            print("  [SKIP] 模板填充模块导入失败")
            return False
        
        try:
            # 测试简单填充
            template = "您好，{{name}}！今天是{{date}}。"
            result = quick_fill(template, name="张三", date="2026-04-25")
            assert "张三" in result
            assert "2026-04-25" in result
            print(f"  [PASS] 简单模板填充: {result}")
            
            # 测试列表填充
            filler = TemplateFiller()
            template = "{{#loop items}}<li>{{name}}</li>{{/loop}}"
            variables = {"items": [{"name": "项目1"}, {"name": "项目2"}]}
            result = filler.fill_text(template, variables)
            assert "项目1" in result
            assert "项目2" in result
            print(f"  [PASS] 列表模板填充: {result[:50]}...")
            
            self.results.append(('模板填充', True, '模块正常'))
            return True
            
        except Exception as e:
            print(f"  [FAIL] 模板填充失败: {e}")
            self.results.append(('模板填充', False, str(e)))
            return False
    
    def test_intelligence_integration(self) -> bool:
        """测试智能分析集成"""
        print("\n" + "=" * 60)
        print("测试: 智能分析集成")
        print("=" * 60)
        
        try:
            from documents.document_intelligence import (
                DocumentIntelligenceService, extract_entities, summarize_document
            )
        except ImportError:
            print("  [SKIP] 智能分析模块导入失败")
            return False
        
        try:
            # 测试实体识别
            text = "会议定于2024年1月15日在北京市召开，张三和李四参加。"
            entities = extract_entities(text)
            print(f"  实体识别: 发现 {len(entities)} 个实体")
            for e in entities[:5]:
                print(f"    - {e.text} ({e.entity_type.name})")
            
            # 测试文档摘要
            content = """# 测试报告

## 概述

这是一个测试文档，用于验证智能分析功能。

## 内容

会议于2024年1月15日在北京召开，预算100万元。
"""
            summary = summarize_document(content)
            print(f"  摘要生成: {summary.title}")
            print(f"  章节数: {len(summary.sections)}")
            print(f"  统计: {summary.statistics}")
            
            self.results.append(('智能分析', True, '模块正常'))
            return True
            
        except Exception as e:
            print(f"  [FAIL] 智能分析失败: {e}")
            self.results.append(('智能分析', False, str(e)))
            return False
    
    def run_all_tests(self) -> Dict[str, Any]:
        """运行所有测试"""
        print("\n" + "=" * 60)
        print("文档模块 v22.3 集成测试")
        print("=" * 60)
        
        # 检查依赖
        dep_results = DependencyChecker.print_report()
        
        # 运行测试
        self.test_pptx_generation()
        self.test_docx_generation()
        self.test_xlsx_generation()
        self.test_template_integration()
        self.test_intelligence_integration()
        
        # 打印结果汇总
        print("\n" + "=" * 60)
        print("测试结果汇总")
        print("=" * 60)
        
        passed = sum(1 for _, success, _ in self.results if success)
        total = len(self.results)
        
        for name, success, detail in self.results:
            status = f"[{chr(0x2713)}]" if success else f"[{chr(0x2717)}]"
            print(f"  {status} {name}: {detail}")
        
        print(f"\n通过: {passed}/{total}")
        
        # 检查是否需要安装依赖
        missing_required = [k for k, v in dep_results.items() 
                           if not v['installed'] and v['required']]
        
        if missing_required:
            print("\n" + "-" * 60)
            print("安装缺失的必需依赖:")
            for pkg in missing_required:
                print(f"  pip install {dep_results[pkg]['package']}")
        
        return {
            'total': total,
            'passed': passed,
            'failed': total - passed,
            'results': self.results,
            'dependencies': dep_results
        }


def main():
    """主函数"""
    import argparse
    
    parser = argparse.ArgumentParser(description='文档模块集成测试')
    parser.add_argument('--check-deps', action='store_true', help='仅检查依赖')
    parser.add_argument('--output', type=str, help='输出目录')
    args = parser.parse_args()
    
    if args.check_deps:
        DependencyChecker.print_report()
    else:
        output_dir = Path(args.output) if args.output else None
        tester = DocumentGeneratorTester(output_dir)
        results = tester.run_all_tests()
        
        # 返回退出码
        sys.exit(0 if results['failed'] == 0 else 1)


if __name__ == "__main__":
    main()
