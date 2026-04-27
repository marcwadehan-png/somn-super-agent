"""
统一文档处理服务 - Natural Language Document Service
整合PPT/Word/Excel的智能处理能力

功能:
- 自然语言驱动的文档生成/读取/修改
- 多格式转换和批处理
- 模板填充和智能分析
"""

from pathlib import Path
from typing import Dict, Any, List, Optional, Union
from enum import Enum
import json
import logging
from dataclasses import dataclass

logger = logging.getLogger(__name__)

# Skill路径配置（使用动态路径）
_workbuddy_base = Path.home() / ".workbuddy"
SKILL_BASE = _workbuddy_base / "plugins" / "marketplaces" / "cb_teams_marketplace" / "plugins" / "document-skills" / "skills"
SCRIPTS_BASE = SKILL_BASE.parent / "scripts"


class DocumentIntent(Enum):
    """文档处理意图"""
    GENERATE_PPT = "生成PPT"
    GENERATE_WORD = "生成Word"
    GENERATE_EXCEL = "生成Excel"
    READ_DOCUMENT = "读取文档"
    EDIT_DOCUMENT = "编辑文档"
    CONVERT_FORMAT = "格式转换"
    MERGE_DOCUMENTS = "合并文档"
    ANALYZE_DATA = "数据分析"
    TEMPLATE_FILL = "模板填充"
    UNKNOWN = "未知"


class DocumentFormat(Enum):
    """文档格式"""
    PPTX = "pptx"
    DOCX = "docx"
    XLSX = "xlsx"
    PDF = "pdf"
    CSV = "csv"
    MARKDOWN = "md"
    HTML = "html"


@dataclass
class DocumentResult:
    """文档处理结果"""
    success: bool
    message: str
    file_path: Optional[str] = None
    data: Optional[Any] = None
    errors: Optional[List[str]] = None


class NaturalLanguageDocumentService:
    """
    自然语言文档处理服务

    支持的指令模式:
    - "帮我做一个PPT/演示文稿/幻灯片"
    - "生成项目报告Word文档"
    - "制作销售数据Excel表格"
    - "读取这份文档的内容"
    - "给Word添加目录"
    - "把PDF转成Word"
    - "合并这三个文档"
    - "分析这份Excel数据"
    """

    def __init__(self):
        self._init_services()
        self.intent_keywords = self._build_intent_keywords()

    def _init_services(self):
        """初始化各文档服务"""
        # PPT服务
        try:
            # 尝试多种导入方式
            try:
                from src.ppt.ppt_service import PPTService
            except ImportError:
                try:
                    from ppt.ppt_service import PPTService
                except ImportError:
                    try:
                        from smart_office_assistant.src.ppt.ppt_service import PPTService
                    except ImportError:
                        # 最后尝试相对导入
                        import sys
                        import os
                        # 添加项目根目录到路径
                        current_file = Path(__file__).resolve()
                        project_root = current_file.parent.parent.parent
                        if str(project_root) not in sys.path:
                            sys.path.insert(0, str(project_root))
                        from src.ppt.ppt_service import PPTService
            self.ppt_service = PPTService(theme="business", enable_charts=True)
            self._ppt_available = True
        except ImportError as e:
            logger.warning(f"PPT服务初始化失败: {e}")
            self.ppt_service = None
            self._ppt_available = False

        # Word生成器
        try:
            try:
                from src.documents.docx_generator import DOCXGenerator
            except ImportError:
                from documents.docx_generator import DOCXGenerator
            self.docx_generator = DOCXGenerator()
            self._docx_available = True
        except ImportError as e:
            logger.warning(f"Word服务初始化失败: {e}")
            self.docx_generator = None
            self._docx_available = False

        # Excel生成器
        try:
            try:
                from src.documents.excel_generator import ExcelGenerator
            except ImportError:
                from documents.excel_generator import ExcelGenerator
            self.excel_generator = ExcelGenerator()
            self._xlsx_available = True
        except ImportError as e:
            logger.warning(f"Excel服务初始化失败: {e}")
            self.excel_generator = None
            self._xlsx_available = False

        # 文档解析器
        self._parser_available = True

    def _build_intent_keywords(self) -> Dict[DocumentIntent, List[str]]:
        """构建意图识别关键词"""
        return {
            DocumentIntent.GENERATE_PPT: [
                "ppt", "PPT", "演示", "幻灯片", "演示文稿",
                "做PPT", "生成PPT", "创建PPT", "制作PPT",
                "做演示", "生成演示"
            ],
            DocumentIntent.GENERATE_WORD: [
                "word", "Word", "文档", "报告", "文章",
                "生成Word", "写Word", "创建Word", "制作Word",
                "生成文档", "写文档", "创建文档"
            ],
            DocumentIntent.GENERATE_EXCEL: [
                "excel", "Excel", "表格", "报表", "数据表",
                "生成Excel", "做表格", "制作Excel", "做报表",
                "生成表格", "制作表格", "数据分析"
            ],
            DocumentIntent.READ_DOCUMENT: [
                "读取", "读取文档", "查看内容", "打开", "分析",
                "提取内容", "查看这个", "读取这份"
            ],
            DocumentIntent.EDIT_DOCUMENT: [
                "编辑", "修改", "添加", "更新", "删除",
                "给...添加", "在...中加入"
            ],
            DocumentIntent.CONVERT_FORMAT: [
                "转换", "转成", "转PDF", "转Word", "转Excel",
                "导出为", "转换成"
            ],
            DocumentIntent.MERGE_DOCUMENTS: [
                "合并", "整合", "拼接", "汇总"
            ],
            DocumentIntent.ANALYZE_DATA: [
                "分析", "统计", "汇总", "透视",
                "分析这份数据", "统计一下"
            ],
            DocumentIntent.TEMPLATE_FILL: [
                "模板", "填充", "套用模板", "按模板"
            ]
        }

    def process_command(self, command: str, context: Optional[Dict] = None) -> DocumentResult:
        """
        处理自然语言文档指令

        Args:
            command: 自然语言指令
            context: 可选的上下文信息

        Returns:
            DocumentResult: 处理结果
        """
        # 1. 意图识别
        intent = self._classify_intent(command)

        # 2. 参数提取
        params = self._extract_params(command, intent, context)

        # 3. 执行处理
        try:
            result = self._execute_intent(intent, params, context)
            return result
        except Exception as e:
            logger.error(f"文档处理失败: {e}")
            return DocumentResult(
                success=False,
                message="处理失败"
            )

    def _classify_intent(self, command: str) -> DocumentIntent:
        """识别处理意图"""
        # 优先级匹配（生成类优先）
        intent_priority = [
            DocumentIntent.GENERATE_PPT,
            DocumentIntent.GENERATE_WORD,
            DocumentIntent.GENERATE_EXCEL,
            DocumentIntent.CONVERT_FORMAT,
            DocumentIntent.READ_DOCUMENT,
            DocumentIntent.EDIT_DOCUMENT,
            DocumentIntent.MERGE_DOCUMENTS,
            DocumentIntent.ANALYZE_DATA,
            DocumentIntent.TEMPLATE_FILL
        ]

        for intent in intent_priority:
            keywords = self.intent_keywords.get(intent, [])
            if any(kw in command for kw in keywords):
                return intent

        return DocumentIntent.UNKNOWN

    def _extract_params(
        self,
        command: str,
        intent: DocumentIntent,
        context: Optional[Dict]
    ) -> Dict[str, Any]:
        """提取处理参数"""
        params = {}

        # 提取文件路径
        path_keywords = ["这个", "这份", "这个文件", "那份", "该文件"]
        for kw in path_keywords:
            if kw in command:
                idx = command.find(kw)
                # 简化处理：假设上下文中有文件路径
                if context and context.get("file_path"):
                    params["file_path"] = context["file_path"]
                break

        # 提取格式要求
        format_keywords = {
            ".pdf": DocumentFormat.PDF,
            ".docx": DocumentFormat.DOCX,
            ".xlsx": DocumentFormat.XLSX,
            ".pptx": DocumentFormat.PPTX
        }
        for kw, fmt in format_keywords.items():
            if kw in command:
                params["target_format"] = fmt
                break

        # 提取内容主题（简化处理）
        # 实际应用中应该用NLP解析
        params["content"] = command

        return params

    def _execute_intent(
        self,
        intent: DocumentIntent,
        params: Dict,
        context: Optional[Dict]
    ) -> DocumentResult:
        """执行文档处理"""
        handlers = {
            DocumentIntent.GENERATE_PPT: self._handle_generate_ppt,
            DocumentIntent.GENERATE_WORD: self._handle_generate_word,
            DocumentIntent.GENERATE_EXCEL: self._handle_generate_excel,
            DocumentIntent.READ_DOCUMENT: self._handle_read_document,
            DocumentIntent.EDIT_DOCUMENT: self._handle_edit_document,
            DocumentIntent.CONVERT_FORMAT: self._handle_convert_format,
            DocumentIntent.MERGE_DOCUMENTS: self._handle_merge_documents,
            DocumentIntent.ANALYZE_DATA: self._handle_analyze_data,
            DocumentIntent.TEMPLATE_FILL: self._handle_template_fill
        }

        handler = handlers.get(intent, self._handle_unknown)
        return handler(params, context)

    def _handle_generate_ppt(
        self,
        params: Dict,
        context: Optional[Dict]
    ) -> DocumentResult:
        """处理PPT生成"""
        if not self._ppt_available:
            return DocumentResult(
                success=False,
                message="PPT服务暂不可用"
            )

        try:
            from src.ppt.ppt_generator import ContentFormat

            content = params.get("content", "")
            output_path = params.get("output_path")

            # 使用PPT服务生成
            result_path = self.ppt_service.generate_ppt(
                content=content,
                format=ContentFormat.MARKDOWN,
                output_path=output_path,
                beautify=True,
                auto_charts=True
            )

            return DocumentResult(
                success=True,
                message="PPT生成成功",
                file_path=result_path
            )

        except Exception as e:
            return DocumentResult(
                success=False,
                message="PPT生成失败"
            )

    def _handle_generate_word(
        self,
        params: Dict,
        context: Optional[Dict]
    ) -> DocumentResult:
        """处理Word文档生成"""
        if not self._docx_available:
            return DocumentResult(
                success=False,
                message="Word服务暂不可用"
            )

        try:
            content = params.get("content", "")
            output_path = params.get("output_path", "output/document.docx")

            # 创建Word文档
            doc = self.docx_generator

            # 简单的标题-内容解析
            lines = content.split("\n")
            for line in lines:
                line = line.strip()
                if not line:
                    continue

                if line.startswith("#"):
                    level = line.count("#")
                    text = line.lstrip("#").strip()
                    doc.add_heading(text, level=min(level, 6))
                elif line.startswith("-") or line.startswith("*"):
                    doc.add_bullet_list([line.lstrip("-* ").strip()])
                else:
                    doc.add_paragraph(line)

            result_path = doc.save(output_path)

            return DocumentResult(
                success=True,
                message="Word文档生成成功",
                file_path=result_path
            )

        except Exception as e:
            return DocumentResult(
                success=False,
                message="Word文档生成失败"
            )

    def _handle_generate_excel(
        self,
        params: Dict,
        context: Optional[Dict]
    ) -> DocumentResult:
        """处理Excel表格生成"""
        if not self._xlsx_available:
            return DocumentResult(
                success=False,
                message="Excel服务暂不可用"
            )

        try:
            content = params.get("content", "")
            output_path = params.get("output_path", "output/data.xlsx")

            # 创建Excel
            excel = self.excel_generator
            excel.create_sheet("数据表")

            # 解析简单的表格数据（CSV格式）
            lines = [l for l in content.split("\n") if l.strip()]
            if lines:
                # 第一行作为表头
                headers = [h.strip() for h in lines[0].split(",")]
                data = []

                for line in lines[1:]:
                    row = [c.strip() for c in line.split(",")]
                    data.append(row)

                excel.add_data(data, headers)

            result_path = excel.save(output_path)

            return DocumentResult(
                success=True,
                message="Excel表格生成成功",
                file_path=result_path
            )

        except Exception as e:
            return DocumentResult(
                success=False,
                message="Excel表格生成失败"
            )

    def _handle_read_document(
        self,
        params: Dict,
        context: Optional[Dict]
    ) -> DocumentResult:
        """处理文档读取"""
        file_path = params.get("file_path") or (context.get("file_path") if context else None)

        if not file_path:
            return DocumentResult(
                success=False,
                message="未指定文件路径"
            )

        file_path = Path(file_path)
        if not file_path.exists():
            return DocumentResult(
                success=False,
                message=f"文件不存在: {file_path}"
            )

        try:
            content = self._parse_document(file_path)
            return DocumentResult(
                success=True,
                message="文档读取成功",
                data=content
            )
        except Exception as e:
            return DocumentResult(
                success=False,
                message="文档读取失败"
            )

    def _parse_document(self, file_path: Path) -> Dict[str, Any]:
        """解析文档内容"""
        suffix = file_path.suffix.lower()

        if suffix == ".docx":
            return self._parse_word(file_path)
        elif suffix in [".xlsx", ".xls", ".csv"]:
            return self._parse_excel(file_path)
        elif suffix == ".pptx":
            return self._parse_powerpoint(file_path)
        elif suffix == ".pdf":
            return self._parse_pdf(file_path)
        else:
            # 尝试作为文本文件读取
            with open(file_path, "r", encoding="utf-8") as f:
                return {"text": f.read(), "type": "text"}

    def _parse_word(self, file_path: Path) -> Dict[str, Any]:
        """解析Word文档"""
        try:
            from docx import Document
            doc = Document(str(file_path))

            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)

            return {
                "type": "word",
                "paragraphs": paragraphs,
                "tables": tables,
                "text": "\n".join(paragraphs)
            }
        except ImportError:
            return self._parse_docx_fallback(file_path)

    def _parse_docx_fallback(self, file_path: Path) -> Dict[str, Any]:
        """使用zipfile解析docx（无python-docx时）"""
        import zipfile
        import re

        try:
            with zipfile.ZipFile(file_path, 'r') as zf:
                xml_content = zf.read('word/document.xml').decode('utf-8')
                # 移除XML标签
                text = re.sub(r'<[^>]+>', ' ', xml_content)
                text = re.sub(r'\s+', ' ', text)
                return {"type": "word", "text": text.strip()}
        except Exception as e:
            return {"type": "word", "error": "Word文档解析失败"}

    def _parse_excel(self, file_path: Path) -> Dict[str, Any]:
        """解析Excel文档"""
        try:
            import pandas as pd

            suffix = file_path.suffix.lower()
            if suffix == ".csv":
                df = pd.read_csv(file_path, encoding="utf-8")
            else:
                df = pd.read_excel(file_path)

            return {
                "type": "excel",
                "data": df.to_dict("records"),
                "columns": list(df.columns),
                "row_count": len(df),
                "preview": df.head(10).to_dict("records")
            }
        except ImportError:
            return self._parse_excel_basic(file_path)

    def _parse_excel_basic(self, file_path: Path) -> Dict[str, Any]:
        """基础Excel解析（无pandas时）"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)

            sheets = {}
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    row_data = [str(cell) if cell else "" for cell in row]
                    if any(row_data):
                        rows.append(row_data)
                sheets[sheet_name] = rows

            return {
                "type": "excel",
                "sheets": sheets,
                "sheet_count": len(sheets)
            }
        except Exception as e:
            return {"type": "excel", "error": "Excel文档解析失败"}

    def _parse_powerpoint(self, file_path: Path) -> Dict[str, Any]:
        """解析PowerPoint文档"""
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))

            slides = []
            for i, slide in enumerate(prs.slides, 1):
                slide_content = {"slide": i, "text": []}

                if slide.shapes.title:
                    slide_content["title"] = slide.shapes.title.text

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        if not (slide.shapes.title and shape == slide.shapes.title):
                            slide_content["text"].append(shape.text)

                slides.append(slide_content)

            return {
                "type": "powerpoint",
                "slides": slides,
                "slide_count": len(slides)
            }
        except ImportError:
            return {"type": "powerpoint", "error": "python-pptx未安装"}

    def _parse_pdf(self, file_path: Path) -> Dict[str, Any]:
        """解析PDF文档"""
        # 优先使用PDF skill的脚本
        script_path = SCRIPTS_BASE / "office" / "unpack.py"
        if script_path.exists():
            # 实际应用中应该调用脚本
            pass

        return {
            "type": "pdf",
            "message": "PDF解析需要使用PDF skill"
        }

    def _handle_edit_document(
        self,
        params: Dict,
        context: Optional[Dict]
    ) -> DocumentResult:
        """处理文档编辑"""
        return DocumentResult(
            success=False,
            message="文档编辑功能开发中"
        )

    def _handle_convert_format(
        self,
        params: Dict,
        context: Optional[Dict]
    ) -> DocumentResult:
        """处理格式转换"""
        return DocumentResult(
            success=False,
            message="格式转换功能开发中，可使用PDF skill"
        )

    def _handle_merge_documents(
        self,
        params: Dict,
        context: Optional[Dict]
    ) -> DocumentResult:
        """处理文档合并"""
        return DocumentResult(
            success=False,
            message="文档合并功能开发中"
        )

    def _handle_analyze_data(
        self,
        params: Dict,
        context: Optional[Dict]
    ) -> DocumentResult:
        """处理数据分析"""
        file_path = params.get("file_path") or (context.get("file_path") if context else None)

        if not file_path:
            return DocumentResult(
                success=False,
                message="未指定数据文件"
            )

        try:
            data = self._parse_document(Path(file_path))

            if data.get("type") == "excel":
                # 简单的数据分析
                df_data = data.get("data", [])
                analysis = {
                    "record_count": len(df_data),
                    "columns": data.get("columns", []),
                    "summary": f"共{len(df_data)}条记录，{len(data.get('columns', []))}个字段"
                }
                return DocumentResult(
                    success=True,
                    message="数据分析完成",
                    data=analysis
                )

            return DocumentResult(
                success=True,
                message="数据提取完成",
                data=data
            )

        except Exception as e:
            return DocumentResult(
                success=False,
                message="数据分析失败"
            )

    def _handle_template_fill(
        self,
        params: Dict,
        context: Optional[Dict]
    ) -> DocumentResult:
        """处理模板填充"""
        return DocumentResult(
            success=False,
            message="模板填充功能开发中"
        )

    def _handle_unknown(
        self,
        params: Dict,
        context: Optional[Dict]
    ) -> DocumentResult:
        """处理未知意图"""
        return DocumentResult(
            success=False,
            message="无法识别处理意图，请明确说明要做什么"
        )


# 便捷函数
def quick_process(command: str, **kwargs) -> DocumentResult:
    """快速处理文档指令"""
    service = NaturalLanguageDocumentService()
    return service.process_command(command, kwargs)


if __name__ == "__main__":
    # 示例用法
    service = NaturalLanguageDocumentService()

    # 测试指令
    test_commands = [
        "帮我生成一个项目汇报PPT",
        "创建一个会议纪要Word文档",
        "制作销售数据分析Excel表格"
    ]

    for cmd in test_commands:
        print(f"\n指令: {cmd}")
        result = service.process_command(cmd)
        print(f"结果: {result.message}")
