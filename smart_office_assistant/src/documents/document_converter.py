# -*- coding: utf-8 -*-
"""
文档转换工具 [v22.0 Phase 3]
Document Converter

功能:
- Word/Excel/PowerPoint互转
- Markdown/HTML互转
- PDF与其他格式转换
- 支持批量转换

版本: v22.0.0
日期: 2026-04-25
"""

from __future__ import annotations

import os
import logging
from pathlib import Path
from typing import Dict, List, Optional, Any, Union, Tuple
from dataclasses import dataclass
from enum import Enum, auto
import subprocess
import tempfile

logger = logging.getLogger(__name__)


class Format(Enum):
    """文档格式"""
    DOCX = "docx"
    DOC = "doc"
    XLSX = "xlsx"
    XLS = "xls"
    CSV = "csv"
    PPTX = "pptx"
    PPT = "ppt"
    PDF = "pdf"
    TXT = "txt"
    MD = "md"
    HTML = "html"
    RTF = "rtf"
    ODT = "odt"
    ODS = "ods"
    ODP = "odp"
    XML = "xml"
    JSON = "json"


@dataclass
class ConversionResult:
    """转换结果"""
    success: bool
    input_path: Optional[Path] = None
    output_path: Optional[Path] = None
    error: Optional[str] = None
    warnings: List[str] = None

    def __post_init__(self):
        if self.warnings is None:
            self.warnings = []


class DocumentConverter:
    """
    文档格式转换器
    
    支持的转换:
    - Word: docx, doc, rtf, odt, txt
    - Excel: xlsx, xls, csv, ods
    - PowerPoint: pptx, ppt, odp
    - 通用: pdf, html, md, txt
    
    示例:
    ```python
    converter = DocumentConverter()
    
    # 转换单个文件
    result = converter.convert("input.docx", "output.pdf")
    
    # 批量转换
    results = converter.batch_convert(
        input_dir="docs/",
        output_dir="pdfs/",
        from_format="docx",
        to_format="pdf"
    )
    ```
    """

    # 转换映射表 (源格式 -> 目标格式 -> 方法)
    CONVERSION_MAP: Dict[str, Dict[str, str]] = {
        'docx': {
            'pdf': '_docx_to_pdf',
            'html': '_docx_to_html',
            'md': '_docx_to_md',
            'txt': '_docx_to_txt',
            'doc': '_copy',
            'rtf': '_docx_to_rtf',
        },
        'xlsx': {
            'csv': '_xlsx_to_csv',
            'pdf': '_xlsx_to_pdf',
            'html': '_xlsx_to_html',
            'xls': '_xlsx_to_xls',
        },
        'pptx': {
            'pdf': '_pptx_to_pdf',
            'html': '_pptx_to_html',
            'ppt': '_copy',
            'odp': '_pptx_to_odp',
        },
        'md': {
            'html': '_md_to_html',
            'docx': '_md_to_docx',
            'pdf': '_md_to_pdf',
            'txt': '_copy',
        },
        'html': {
            'md': '_html_to_md',
            'docx': '_html_to_docx',
            'pdf': '_html_to_pdf',
            'txt': '_html_to_txt',
        },
        'txt': {
            'docx': '_txt_to_docx',
            'md': '_copy',
            'html': '_txt_to_html',
        },
        'csv': {
            'xlsx': '_csv_to_xlsx',
            'json': '_csv_to_json',
        },
        'pdf': {
            'txt': '_pdf_to_txt',
            'html': '_pdf_to_html',
            'md': '_pdf_to_md',
        }
    }

    def __init__(self, temp_dir: Optional[Union[str, Path]] = None):
        """
        初始化转换器
        
        Args:
            temp_dir: 临时目录
        """
        self.temp_dir = Path(temp_dir) if temp_dir else Path(tempfile.gettempdir())
        self._check_dependencies()

    def _check_dependencies(self) -> None:
        """检查依赖"""
        self._libreoffice_available = self._check_libreoffice()
        self._pandoc_available = self._check_pandoc()

        if not self._libreoffice_available and not self._pandoc_available:
            logger.warning(
                "建议安装 LibreOffice 或 Pandoc 以支持更多格式转换\n"
                "LibreOffice: https://www.libreoffice.org/download/\n"
                "Pandoc: https://pandoc.org/installing.html"
            )

    def _check_libreoffice(self) -> bool:
        """检查LibreOffice"""
        try:
            result = subprocess.run(
                ['soffice', '--version'],
                capture_output=True,
                timeout=5
            )
            return result.returncode == 0
        except (subprocess.SubprocessError, FileNotFoundError):
            return False

    def _check_pandoc(self) -> bool:
        """检查Pandoc"""
        try:
            result = subprocess.run(
                ['pandoc', '--version'],
                capture_output=True,
                timeout=5
            )
            return result.returncode == 0
        except (subprocess.SubprocessError, FileNotFoundError):
            return False

    def convert(
        self,
        input_path: Union[str, Path],
        output_path: Union[str, Path],
        options: Optional[Dict[str, Any]] = None
    ) -> ConversionResult:
        """
        转换单个文件
        
        Args:
            input_path: 输入文件路径
            output_path: 输出文件路径
            options: 转换选项
        
        Returns:
            ConversionResult: 转换结果
        """
        input_path = Path(input_path)
        output_path = Path(output_path)
        options = options or {}

        if not input_path.exists():
            return ConversionResult(
                success=False,
                input_path=input_path,
                error=f"输入文件不存在: {input_path}"
            )

        # 获取格式
        from_format = input_path.suffix.lstrip('.').lower()
        to_format = output_path.suffix.lstrip('.').lower()

        # 相同格式直接复制
        if from_format == to_format:
            return self._copy(input_path, output_path, options)

        # 查找转换方法
        if from_format in self.CONVERSION_MAP and to_format in self.CONVERSION_MAP[from_format]:
            method_name = self.CONVERSION_MAP[from_format][to_format]
        else:
            # 尝试通过中间格式转换
            return self._convert_via_intermediate(input_path, output_path, from_format, to_format, options)

        # 调用转换方法
        method = getattr(self, method_name, None)
        if method:
            return method(input_path, output_path, options)
        else:
            return ConversionResult(
                success=False,
                input_path=input_path,
                error=f"不支持的转换: {from_format} -> {to_format}"
            )

    def batch_convert(
        self,
        input_dir: Union[str, Path],
        output_dir: Union[str, Path],
        from_format: str,
        to_format: str,
        recursive: bool = False,
        options: Optional[Dict[str, Any]] = None
    ) -> List[ConversionResult]:
        """
        批量转换文件
        
        Args:
            input_dir: 输入目录
            output_dir: 输出目录
            from_format: 源格式
            to_format: 目标格式
            recursive: 是否递归子目录
            options: 转换选项
        
        Returns:
            List[ConversionResult]: 转换结果列表
        """
        input_dir = Path(input_dir)
        output_dir = Path(output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)

        # 收集文件
        pattern = f"**/*.{from_format}" if recursive else f"*.{from_format}"
        files = list(input_dir.glob(pattern))

        results = []
        for file_path in files:
            # 计算输出路径
            if recursive:
                rel_path = file_path.relative_to(input_dir)
                out_path = output_dir / rel_path.with_suffix(f".{to_format}")
            else:
                out_path = output_dir / file_path.with_suffix(f".{to_format}").name

            # 确保输出子目录存在
            out_path.parent.mkdir(parents=True, exist_ok=True)

            # 执行转换
            result = self.convert(file_path, out_path, options)
            results.append(result)

        return results

    def _convert_via_intermediate(
        self,
        input_path: Path,
        output_path: Path,
        from_format: str,
        to_format: str,
        options: Dict
    ) -> ConversionResult:
        """通过中间格式转换"""
        # 定义转换路径
        conversion_paths = {
            ('docx', 'pdf'): [('docx', 'html'), ('html', 'pdf')] if not self._libreoffice_available else [],
            ('xlsx', 'pdf'): [('xlsx', 'html'), ('html', 'pdf')] if not self._libreoffice_available else [],
            ('pptx', 'pdf'): [('pptx', 'html'), ('html', 'pdf')] if not self._libreoffice_available else [],
            ('md', 'docx'): [('md', 'html'), ('html', 'docx')] if not self._pandoc_available else [],
        }

        key = (from_format, to_format)
        if key in conversion_paths:
            path = conversion_paths[key]
            if path:
                # 需要多次转换
                current = input_path
                for i, (step_from, step_to) in enumerate(path):
                    temp_out = self.temp_dir / f"temp_{i}_{current.stem}.{step_to}"
                    result = self.convert(current, temp_out, options)
                    if not result.success:
                        return result
                    current = temp_out
                
                # 最后一步
                result = self.convert(current, output_path, options)
                # 清理临时文件
                for f in self.temp_dir.glob(f"temp_*"):
                    f.unlink()
                return result

        return ConversionResult(
            success=False,
            input_path=input_path,
            error=f"不支持的转换: {from_format} -> {to_format}"
        )

    # ==================== 转换方法 ====================

    def _copy(self, input_path: Path, output_path: Path, options: Dict) -> ConversionResult:
        """复制文件"""
        import shutil
        try:
            output_path.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(input_path, output_path)
            return ConversionResult(success=True, input_path=input_path, output_path=output_path)
        except Exception as e:
            logger.error(f"文件复制失败: {e}")
            return ConversionResult(success=False, input_path=input_path, error="转换失败")

    def _docx_to_pdf(self, input_path: Path, output_path: Path, options: Dict) -> ConversionResult:
        """Word转PDF"""
        if self._libreoffice_available:
            return self._libreoffice_convert(input_path, output_path, 'pdf')
        return ConversionResult(
            success=False,
            input_path=input_path,
            error="需要安装LibreOffice才能转换Word为PDF"
        )

    def _docx_to_html(self, input_path: Path, output_path: Path, options: Dict) -> ConversionResult:
        """Word转HTML"""
        try:
            from docx import Document
            doc = Document(input_path)
            
            html_parts = ['<!DOCTYPE html>', '<html>', '<head>', 
                         f'<meta charset="utf-8">',
                         f'<title>{input_path.stem}</title>',
                         '</head>', '<body>']
            
            for para in doc.paragraphs:
                style = para.style.name.lower()
                if 'heading' in style:
                    level = style.replace('heading', '').strip() or '1'
                    html_parts.append(f'<h{level}>{para.text}</h{level}>')
                else:
                    html_parts.append(f'<p>{para.text}</p>')
            
            # 处理表格
            for table in doc.tables:
                html_parts.append('<table>')
                for row in table.rows:
                    html_parts.append('<tr>')
                    for cell in row.cells:
                        html_parts.append(f'<td>{cell.text}</td>')
                    html_parts.append('</tr>')
                html_parts.append('</table>')
            
            html_parts.extend(['</body>', '</html>'])
            
            output_path.parent.mkdir(parents=True, exist_ok=True)
            output_path.write_text('\n'.join(html_parts), encoding='utf-8')
            return ConversionResult(success=True, input_path=input_path, output_path=output_path)
        except ImportError:
            return ConversionResult(success=False, input_path=input_path, error="python-docx未安装")
        except Exception as e:
            logger.error(f"Word转HTML失败: {e}")
            return ConversionResult(success=False, input_path=input_path, error="转换失败")

    def _docx_to_md(self, input_path: Path, output_path: Path, options: Dict) -> ConversionResult:
        """Word转Markdown"""
        if self._pandoc_available:
            return self._pandoc_convert(input_path, output_path, 'markdown')
        return self._docx_to_html(input_path, output_path.with_suffix('.html'), options)

    def _docx_to_txt(self, input_path: Path, output_path: Path, options: Dict) -> ConversionResult:
        """Word转纯文本"""
        try:
            from docx import Document
            doc = Document(input_path)
            text = '\n\n'.join(para.text for para in doc.paragraphs)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            output_path.write_text(text, encoding='utf-8')
            return ConversionResult(success=True, input_path=input_path, output_path=output_path)
        except ImportError:
            return ConversionResult(success=False, input_path=input_path, error="python-docx未安装")
        except Exception as e:
            logger.error(f"Word转纯文本失败: {e}")
            return ConversionResult(success=False, input_path=input_path, error="转换失败")

    def _xlsx_to_csv(self, input_path: Path, output_path: Path, options: Dict) -> ConversionResult:
        """Excel转CSV"""
        try:
            import pandas as pd
            df = pd.read_excel(input_path, **options.get('read_options', {}))
            output_path.parent.mkdir(parents=True, exist_ok=True)
            df.to_csv(output_path, index=False, encoding='utf-8')
            return ConversionResult(success=True, input_path=input_path, output_path=output_path)
        except ImportError:
            return ConversionResult(success=False, input_path=input_path, error="pandas未安装")
        except Exception as e:
            logger.error(f"Excel转CSV失败: {e}")
            return ConversionResult(success=False, input_path=input_path, error="转换失败")

    def _xlsx_to_pdf(self, input_path: Path, output_path: Path, options: Dict) -> ConversionResult:
        """Excel转PDF"""
        if self._libreoffice_available:
            return self._libreoffice_convert(input_path, output_path, 'pdf')
        return ConversionResult(
            success=False,
            input_path=input_path,
            error="需要安装LibreOffice才能转换Excel为PDF"
        )

    def _xlsx_to_html(self, input_path: Path, output_path: Path, options: Dict) -> ConversionResult:
        """Excel转HTML"""
        try:
            import pandas as pd
            df = pd.read_excel(input_path, **options.get('read_options', {}))
            html = df.to_html()
            
            full_html = f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>{input_path.stem}</title>
    <style>
        table {{ border-collapse: collapse; width: 100%; }}
        th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
        th {{ background-color: #f2f2f2; }}
    </style>
</head>
<body>
{html}
</body>
</html>'''
            output_path.parent.mkdir(parents=True, exist_ok=True)
            output_path.write_text(full_html, encoding='utf-8')
            return ConversionResult(success=True, input_path=input_path, output_path=output_path)
        except ImportError:
            return ConversionResult(success=False, input_path=input_path, error="pandas未安装")
        except Exception as e:
            logger.error(f"Excel转HTML失败: {e}")
            return ConversionResult(success=False, input_path=input_path, error="转换失败")

    def _pptx_to_pdf(self, input_path: Path, output_path: Path, options: Dict) -> ConversionResult:
        """PowerPoint转PDF"""
        if self._libreoffice_available:
            return self._libreoffice_convert(input_path, output_path, 'pdf')
        return ConversionResult(
            success=False,
            input_path=input_path,
            error="需要安装LibreOffice才能转换PowerPoint为PDF"
        )

    def _pptx_to_html(self, input_path: Path, output_path: Path, options: Dict) -> ConversionResult:
        """PowerPoint转HTML"""
        try:
            from pptx import Presentation
            prs = Presentation(input_path)
            
            slides_html = []
            for i, slide in enumerate(prs.slides, 1):
                slides_html.append(f'<div class="slide" id="slide-{i}">')
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides_html.append(f'<p>{shape.text}</p>')
                slides_html.append('</div>')
            
            html = f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>{input_path.stem}</title>
    <style>
        .slide {{ page-break-after: always; margin: 20px; }}
        p {{ margin: 10px 0; }}
    </style>
</head>
<body>
{''.join(slides_html)}
</body>
</html>'''
            output_path.parent.mkdir(parents=True, exist_ok=True)
            output_path.write_text(html, encoding='utf-8')
            return ConversionResult(success=True, input_path=input_path, output_path=output_path)
        except ImportError:
            return ConversionResult(success=False, input_path=input_path, error="python-pptx未安装")
        except Exception as e:
            logger.error(f"PPT转HTML失败: {e}")
            return ConversionResult(success=False, input_path=input_path, error="转换失败")

    def _md_to_html(self, input_path: Path, output_path: Path, options: Dict) -> ConversionResult:
        """Markdown转HTML"""
        if self._pandoc_available:
            return self._pandoc_convert(input_path, output_path, 'html')
        
        # 简单转换
        content = input_path.read_text(encoding='utf-8')
        html_content = content.replace('# ', '<h1>').replace('\n## ', '</h1>\n<h2>')
        html_content = html_content.replace('\n# ', '</h1>\n<h1>').replace('\n### ', '</h2>\n<h3>')
        html_content = html_content.replace('\n- ', '</h3>\n<ul>\n<li>').replace('\n\n', '</li>\n</ul>\n<p>')
        
        html = f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>{input_path.stem}</title>
</head>
<body>
{html_content}
</body>
</html>'''
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(html, encoding='utf-8')
        return ConversionResult(success=True, input_path=input_path, output_path=output_path)

    def _md_to_docx(self, input_path: Path, output_path: Path, options: Dict) -> ConversionResult:
        """Markdown转Word"""
        if self._pandoc_available:
            return self._pandoc_convert(input_path, output_path, 'docx')
        return ConversionResult(
            success=False,
            input_path=input_path,
            error="需要安装Pandoc才能转换Markdown为Word"
        )

    def _html_to_md(self, input_path: Path, output_path: Path, options: Dict) -> ConversionResult:
        """HTML转Markdown"""
        if self._pandoc_available:
            return self._pandoc_convert(input_path, output_path, 'markdown')
        return ConversionResult(
            success=False,
            input_path=input_path,
            error="需要安装Pandoc才能转换HTML为Markdown"
        )

    def _html_to_txt(self, input_path: Path, output_path: Path, options: Dict) -> ConversionResult:
        """HTML转纯文本"""
        try:
            from html.parser import HTMLParser
            
            class TextExtractor(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.text = []
                    self.skip = False
                
                def handle_starttag(self, tag, attrs):
                    if tag in ('script', 'style'):
                        self.skip = True
                    elif tag in ('p', 'div', 'br'):
                        self.text.append('\n')
                
                def handle_endtag(self, tag):
                    if tag in ('script', 'style'):
                        self.skip = False
                
                def handle_data(self, data):
                    if not self.skip:
                        self.text.append(data)
            
            content = input_path.read_text(encoding='utf-8')
            parser = TextExtractor()
            parser.feed(content)
            
            output_path.parent.mkdir(parents=True, exist_ok=True)
            output_path.write_text(''.join(parser.text).strip(), encoding='utf-8')
            return ConversionResult(success=True, input_path=input_path, output_path=output_path)
        except Exception as e:
            logger.error(f"Markdown转HTML失败: {e}")
            return ConversionResult(success=False, input_path=input_path, error="转换失败")

    def _csv_to_json(self, input_path: Path, output_path: Path, options: Dict) -> ConversionResult:
        """CSV转JSON"""
        try:
            import pandas as pd
            df = pd.read_csv(input_path)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            output_path.write_text(df.to_json(orient='records', force_ascii=False, indent=2), encoding='utf-8')
            return ConversionResult(success=True, input_path=input_path, output_path=output_path)
        except ImportError:
            return ConversionResult(success=False, input_path=input_path, error="pandas未安装")
        except Exception as e:
            logger.error(f"CSV转JSON失败: {e}")
            return ConversionResult(success=False, input_path=input_path, error="转换失败")

    def _csv_to_xlsx(self, input_path: Path, output_path: Path, options: Dict) -> ConversionResult:
        """CSV转Excel"""
        try:
            import pandas as pd
            df = pd.read_csv(input_path)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            df.to_excel(output_path, index=False, **options.get('write_options', {}))
            return ConversionResult(success=True, input_path=input_path, output_path=output_path)
        except ImportError:
            return ConversionResult(success=False, input_path=input_path, error="pandas未安装")
        except Exception as e:
            logger.error(f"CSV转Excel失败: {e}")
            return ConversionResult(success=False, input_path=input_path, error="转换失败")

    def _txt_to_docx(self, input_path: Path, output_path: Path, options: Dict) -> ConversionResult:
        """纯文本转Word"""
        try:
            from docx import Document
            doc = Document()
            text = input_path.read_text(encoding='utf-8')
            doc.add_paragraph(text)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            doc.save(output_path)
            return ConversionResult(success=True, input_path=input_path, output_path=output_path)
        except ImportError:
            return ConversionResult(success=False, input_path=input_path, error="python-docx未安装")
        except Exception as e:
            logger.error(f"文本转Word失败: {e}")
            return ConversionResult(success=False, input_path=input_path, error="转换失败")

    def _txt_to_html(self, input_path: Path, output_path: Path, options: Dict) -> ConversionResult:
        """纯文本转HTML"""
        text = input_path.read_text(encoding='utf-8')
        html = f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>{input_path.stem}</title>
</head>
<body>
<pre>{text}</pre>
</body>
</html>'''
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(html, encoding='utf-8')
        return ConversionResult(success=True, input_path=input_path, output_path=output_path)

    def _libreoffice_convert(self, input_path: Path, output_path: Path, format: str) -> ConversionResult:
        """使用LibreOffice转换"""
        try:
            output_path.parent.mkdir(parents=True, exist_ok=True)
            # LibreOffice转换命令
            cmd = [
                'soffice',
                '--headless',
                '--convert-to', format,
                '--outdir', str(output_path.parent),
                str(input_path)
            ]
            result = subprocess.run(cmd, capture_output=True, timeout=60)
            
            if result.returncode == 0:
                # LibreOffice可能改变文件名
                expected = output_path.parent / f"{input_path.stem}.{format}"
                if expected.exists():
                    if expected != output_path:
                        expected.rename(output_path)
                    return ConversionResult(success=True, input_path=input_path, output_path=output_path)
                return ConversionResult(success=True, input_path=input_path, output_path=expected)
            else:
                return ConversionResult(
                    success=False,
                    input_path=input_path,
                    error=result.stderr.decode('utf-8', errors='ignore')
                )
        except subprocess.TimeoutExpired:
            return ConversionResult(success=False, input_path=input_path, error="转换超时")
        except Exception as e:
            logger.error(f"LibreOffice转换失败: {e}")
            return ConversionResult(success=False, input_path=input_path, error="转换失败")

    def _pandoc_convert(self, input_path: Path, output_path: Path, format: str) -> ConversionResult:
        """使用Pandoc转换"""
        try:
            output_path.parent.mkdir(parents=True, exist_ok=True)
            cmd = ['pandoc', str(input_path), '-o', str(output_path)]
            result = subprocess.run(cmd, capture_output=True, timeout=30)
            
            if result.returncode == 0:
                return ConversionResult(success=True, input_path=input_path, output_path=output_path)
            else:
                return ConversionResult(
                    success=False,
                    input_path=input_path,
                    error=result.stderr.decode('utf-8', errors='ignore')
                )
        except subprocess.TimeoutExpired:
            return ConversionResult(success=False, input_path=input_path, error="转换超时")
        except Exception as e:
            logger.error(f"Pandoc转换失败: {e}")
            return ConversionResult(success=False, input_path=input_path, error="转换失败")


# 便捷函数
def quick_convert(
    input_path: Union[str, Path],
    output_path: Union[str, Path],
    **options
) -> ConversionResult:
    """快速转换"""
    converter = DocumentConverter()
    return converter.convert(input_path, output_path, options)
