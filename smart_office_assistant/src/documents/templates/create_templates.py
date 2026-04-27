# -*- coding: utf-8 -*-
"""创建内置模板文件"""

import sys
from pathlib import Path

# 添加项目路径
sys.path.insert(0, str(Path(__file__).parent.parent.parent.parent))

from docx import Document
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt as PPt
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment

TEMPLATES_DIR = Path(__file__).parent


def create_word_template():
    """创建Word模板"""
    doc = Document()

    # 标题
    title = doc.add_heading('', 0)
    run = title.add_run('{{title}}')
    run.font.size = Pt(24)
    run.font.bold = True

    # 副标题
    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub_run = subtitle.add_run('{{subtitle}}')
    sub_run.font.size = Pt(14)
    sub_run.font.color.rgb = RGBColor(128, 128, 128)

    # 基本信息
    doc.add_paragraph()
    info_table = doc.add_table(rows=4, cols=2)
    info_table.style = 'Light Grid Accent 1'
    info_data = [
        ('作者', '{{author}}'),
        ('部门', '{{department}}'),
        ('日期', '{{date}}'),
        ('版本', 'V1.0')
    ]
    for i, (key, value) in enumerate(info_data):
        info_table.rows[i].cells[0].text = key
        info_table.rows[i].cells[1].text = value

    # 分隔线
    doc.add_paragraph('_' * 50)

    # 摘要
    doc.add_heading('摘要', level=1)
    doc.add_paragraph('{{summary}}')

    # 目录占位
    doc.add_heading('目录', level=1)
    doc.add_paragraph('[请在生成后手动更新目录]')

    # 正文占位
    doc.add_heading('正文内容', level=1)
    doc.add_paragraph('{{content}}')

    output_path = TEMPLATES_DIR / 'report_template.docx'
    doc.save(output_path)
    print(f'✓ Word模板创建成功: {output_path}')
    return output_path


def create_meeting_template():
    """创建会议纪要模板"""
    doc = Document()

    # 标题
    title = doc.add_heading('', 0)
    run = title.add_run('{{meeting_title}}')
    run.font.size = Pt(22)

    # 会议信息
    doc.add_heading('会议信息', level=1)
    info_table = doc.add_table(rows=5, cols=2)
    info_table.style = 'Light Grid Accent 1'
    info_data = [
        ('会议名称', '{{meeting_title}}'),
        ('日期时间', '{{date}}'),
        ('主持人', '{{chairman}}'),
        ('记录人', '{{recorder}}'),
        ('参会人员', '{{attendees}}')
    ]
    for i, (key, value) in enumerate(info_data):
        info_table.rows[i].cells[0].text = key
        info_table.rows[i].cells[1].text = value

    # 会议议程
    doc.add_heading('会议议程', level=1)
    doc.add_paragraph('{{agenda}}')

    # 讨论要点
    doc.add_heading('讨论要点', level=1)
    doc.add_paragraph('{{discussion_points}}')

    # 决议事项
    doc.add_heading('决议事项', level=1)
    doc.add_paragraph('{{decisions}}')

    # 行动项
    doc.add_heading('行动项', level=1)
    action_table = doc.add_table(rows=4, cols=3)
    action_table.style = 'Light Grid Accent 1'
    headers = ['任务', '负责人', '截止日期']
    for i, h in enumerate(headers):
        action_table.rows[0].cells[i].text = h
    action_table.rows[1].cells[0].text = '{{task_1}}'
    action_table.rows[2].cells[0].text = '{{task_2}}'
    action_table.rows[3].cells[0].text = '{{task_3}}'

    output_path = TEMPLATES_DIR / 'meeting_template.docx'
    doc.save(output_path)
    print(f'✓ 会议模板创建成功: {output_path}')
    return output_path


def create_excel_data_template():
    """创建Excel数据模板"""
    wb = Workbook()

    # Sheet1 - 数据表
    ws = wb.active
    ws.title = "数据表"

    # 标题
    ws.merge_cells('A1:F1')
    ws['A1'] = '{{title}}'
    ws['A1'].font = Font(size=16, bold=True)
    ws['A1'].alignment = Alignment(horizontal='center')

    # 日期
    ws['A2'] = '日期:'
    ws['B2'] = '{{date}}'

    # 表头
    headers = ['序号', '项目', '分类', '数值', '备注', '状态']
    header_fill = PatternFill(start_color='4472C4', end_color='4472C4', fill_type='solid')
    header_font = Font(color='FFFFFF', bold=True)

    for col, header in enumerate(headers, 1):
        cell = ws.cell(row=4, column=col, value=header)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal='center')

    # 数据行占位
    for row in range(5, 15):
        ws.cell(row=row, column=1, value=row-4)

    # Sheet2 - 统计
    stats = wb.create_sheet("统计")
    stats['A1'] = '统计分析'
    stats['A1'].font = Font(size=14, bold=True)

    output_path = TEMPLATES_DIR / 'data_table.xlsx'
    wb.save(output_path)
    print(f'✓ Excel模板创建成功: {output_path}')
    return output_path


def create_excel_budget_template():
    """创建预算表模板"""
    wb = Workbook()
    ws = wb.active
    ws.title = "预算表"

    # 标题
    ws.merge_cells('A1:G1')
    ws['A1'] = '{{project_name}} - 预算表'
    ws['A1'].font = Font(size=16, bold=True)
    ws['A1'].alignment = Alignment(horizontal='center')

    # 项目信息
    ws['A3'] = '项目名称:'
    ws['B3'] = '{{project_name}}'
    ws['A4'] = '开始日期:'
    ws['B4'] = '{{start_date}}'
    ws['A5'] = '结束日期:'
    ws['B5'] = '{{end_date}}'
    ws['A6'] = '项目经理:'
    ws['B6'] = '{{manager}}'

    # 预算表头
    headers = ['序号', '费用项', '预算金额', '实际金额', '差异', '说明', '状态']
    header_fill = PatternFill(start_color='70AD47', end_color='70AD47', fill_type='solid')
    header_font = Font(color='FFFFFF', bold=True)

    for col, header in enumerate(headers, 1):
        cell = ws.cell(row=8, column=col, value=header)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal='center')

    # 数据行
    categories = ['人员成本', '设备采购', '差旅费用', '外包服务', '其他费用']
    for row, cat in enumerate(categories, 9):
        ws.cell(row=row, column=1, value=row-8)
        ws.cell(row=row, column=2, value=cat)
        ws.cell(row=row, column=3, value=0)
        ws.cell(row=row, column=4, value=0)
        ws.cell(row=row, column=5, value='=D{}-C{}'.format(row, row))

    # 合计行
    ws['A14'] = '合计'
    ws['C14'] = '=SUM(C9:C13)'
    ws['D14'] = '=SUM(D9:D13)'
    ws['E14'] = '=D14-C14'

    output_path = TEMPLATES_DIR / 'budget.xlsx'
    wb.save(output_path)
    print(f'✓ 预算表创建成功: {output_path}')
    return output_path


def create_ppt_template():
    """创建PPT模板"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 标题幻灯片
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = '{{title}}'
    subtitle.text = '{{subtitle}}'

    # 内容幻灯片
    content_layout = prs.slide_layouts[1]
    for i in range(3):
        slide = prs.slides.add_slide(content_layout)
        title = slide.shapes.title
        title.text = f'第{i+1}页 - {{slide_title_{i+1}}}'

    # 结束幻灯片
    end_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(end_layout)

    output_path = TEMPLATES_DIR / 'presentation.pptx'
    prs.save(output_path)
    print(f'✓ PPT模板创建成功: {output_path}')
    return output_path


if __name__ == '__main__':
    print('开始创建模板文件...')
    print()

    try:
        create_word_template()
        create_meeting_template()
        create_excel_data_template()
        create_excel_budget_template()
        create_ppt_template()

        print()
        print('✓ 所有模板创建完成!')
    except ImportError as e:
        print(f'依赖未安装: {e}')
        print('请运行: pip install python-docx python-pptx openpyxl')
    except Exception as e:
        print(f'创建失败: {e}')
