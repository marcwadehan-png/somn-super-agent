# -*- coding: utf-8 -*-
"""主窗口 - 文件操作模块

__all__ = [
    'extract_file_content',
    'index_directory',
    'open_docx',
    'open_excel',
    'open_file',
    'open_pdf',
    'open_ppt',
    'open_with_system_default',
]

处理文件打开、保存、索引等
"""

from pathlib import Path
from PySide6.QtWidgets import QFileDialog, QMessageBox, QProgressDialog
from PySide6.QtCore import Qt
import hashlib

def open_file(main_window, file_path: str):
    """
    打开指定文件

    Args:
        main_window: MainWindow 实例
        file_path: 文件路径
    """
    path = Path(file_path)
    if not path.exists():
        QMessageBox.warning(main_window, "错误", f"文件不存在: {file_path}")
        return

    try:
        suffix = path.suffix.lower()

        # 文本文件 - 直接在编辑器中打开
        if suffix in ['.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.yaml', '.yml', '.xml']:
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            main_window.doc_title_input.setText(path.stem)
            main_window.editor.setPlainText(content)
            main_window.status_bar.showMessage(f"已打开: {path.name}")

        # Word文档
        elif suffix == '.docx':
            open_docx(main_window, path)

        # PDF文档
        elif suffix == '.pdf':
            open_pdf(main_window, path)

        # PPT文档
        elif suffix in ['.pptx', '.ppt']:
            open_ppt(main_window, path)

        # Excel文档
        elif suffix in ['.xlsx', '.xls', '.csv']:
            open_excel(main_window, path)

        # 其他文件 - 使用系统默认程序打开
        else:
            open_with_system_default(main_window, path)

    except Exception as e:
        from loguru import logger
        logger.error(f"打开文件失败: {e}")
        QMessageBox.warning(main_window, "错误", f"无法打开文件: {e}")

def open_docx(main_window, path: Path):
    """打开Word文档"""
    try:
        from docx import Document
        doc = Document(path)

        # 提取文本内容
        content_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                content_parts.append(para.text)

        # 提取表格内容
        for table in doc.tables:
            content_parts.append("\n[表格]")
            for row in table.rows:
                row_text = " | ".join(cell.text for cell in row.cells)
                content_parts.append(row_text)

        content = "\n\n".join(content_parts)
        main_window.doc_title_input.setText(path.stem)
        main_window.editor.setPlainText(content)
        main_window.status_bar.showMessage(f"已打开Word文档: {path.name}")

    except ImportError:
        # 未安装python-docx,使用系统默认程序
        open_with_system_default(main_window, path)
    except Exception as e:
        from loguru import logger
        logger.error(f"打开Word文档失败: {e}")
        open_with_system_default(main_window, path)

def open_pdf(main_window, path: Path):
    """打开PDF文档"""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(path)

        content_parts = []
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            content_parts.append(f"--- 第 {page_num + 1} 页 ---\n{text}")

        doc.close()

        content = "\n\n".join(content_parts)
        main_window.doc_title_input.setText(path.stem)
        main_window.editor.setPlainText(content)
        main_window.status_bar.showMessage(f"已打开PDF文档: {path.name} ({len(doc)} 页)")

    except ImportError:
        # 未安装PyMuPDF,使用系统默认程序
        open_with_system_default(main_window, path)
    except Exception as e:
        from loguru import logger
        logger.error(f"打开PDF文档失败: {e}")
        open_with_system_default(main_window, path)

def open_ppt(main_window, path: Path):
    """打开PPT文档"""
    try:
        from pptx import Presentation
        prs = Presentation(path)

        content_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            content_parts.append(f"=== 幻灯片 {slide_num} ===")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content_parts.append(shape.text)

            content_parts.append("")  # 空行分隔

        content = "\n".join(content_parts)
        main_window.doc_title_input.setText(path.stem)
        main_window.editor.setPlainText(content)
        main_window.status_bar.showMessage(f"已打开PPT: {path.name} ({len(prs.slides)} 页)")

    except ImportError:
        # 未安装python-pptx,使用系统默认程序
        open_with_system_default(main_window, path)
    except Exception as e:
        from loguru import logger
        logger.error(f"打开PPT失败: {e}")
        open_with_system_default(main_window, path)

def open_excel(main_window, path: Path):
    """打开Excel文档"""
    try:
        import pandas as pd

        if path.suffix.lower() == '.csv':
            df = pd.read_csv(path)
            content = df.to_string(index=False)
        else:
            # 读取所有sheet
            xl = pd.ExcelFile(path)
            content_parts = []

            for sheet_name in xl.sheet_names:
                df = pd.read_excel(path, sheet_name=sheet_name)
                content_parts.append(f"=== 工作表: {sheet_name} ===")
                content_parts.append(df.to_string(index=False))
                content_parts.append("")

            content = "\n".join(content_parts)

        main_window.doc_title_input.setText(path.stem)
        main_window.editor.setPlainText(content)
        main_window.status_bar.showMessage(f"已打开Excel: {path.name}")

    except ImportError:
        # 未安装pandas,使用系统默认程序
        open_with_system_default(main_window, path)
    except Exception as e:
        from loguru import logger
        logger.error(f"打开Excel失败: {e}")
        open_with_system_default(main_window, path)

def open_with_system_default(main_window, path: Path):
    """使用系统默认程序打开文件"""
    import sys
    import os
    import subprocess

    try:
        if sys.platform == 'win32':
            os.startfile(path)
        elif sys.platform == 'darwin':  # macOS
            subprocess.run(['open', str(path)], check=True)
        else:  # Linux
            subprocess.run(['xdg-open', str(path)], check=True)

        main_window.status_bar.showMessage(f"已使用系统程序打开: {path.name}")
    except Exception as e:
        QMessageBox.warning(main_window, "错误", f"无法打开文件: {e}")

def extract_file_content(file_path: Path) -> str:
    """
    提取文件内容

    Args:
        file_path: 文件路径

    Returns:
        文件内容文本
    """
    suffix = file_path.suffix.lower()

    try:
        # 文本文件
        if suffix in ['.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.yaml', '.yml', '.xml']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()

        # Word文档
        elif suffix == '.docx':
            try:
                from docx import Document
                doc = Document(file_path)
                return '\n'.join(para.text for para in doc.paragraphs if para.text.strip())
            except ImportError:
                return f"[Word文档: {file_path.name}]"

        # PDF文档
        elif suffix == '.pdf':
            try:
                import fitz
                doc = fitz.open(file_path)
                content = []
                for page in doc:
                    content.append(page.get_text())
                doc.close()
                return '\n'.join(content)
            except ImportError:
                return f"[PDF文档: {file_path.name}]"

        # PPT文档
        elif suffix == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                content = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            content.append(shape.text)
                return '\n'.join(content)
            except ImportError:
                return f"[PPT演示文稿: {file_path.name}]"

        # Excel文档
        elif suffix in ['.xlsx', '.xls', '.csv']:
            try:
                import pandas as pd
                if suffix == '.csv':
                    df = pd.read_csv(file_path)
                else:
                    df = pd.read_excel(file_path)
                return df.to_string(index=False)
            except ImportError:
                return f"[Excel表格: {file_path.name}]"

        else:
            return f"[文件: {file_path.name}]"

    except Exception as e:
        from loguru import logger
        logger.error(f"提取文件内容失败 {file_path}: {e}")
        return f"[无法读取: {file_path.name}]"

def index_directory(main_window, dir_path: str):
    """
    索引目录中的文件到知识库

    Args:
        main_window: MainWindow 实例
        dir_path: 目录路径
    """
    from pathlib import Path
    import pandas as pd

    path = Path(dir_path)
    if not path.exists() or not path.is_dir():
        QMessageBox.warning(main_window, "错误", "无效的目录路径")
        return

    # 支持的文件类型
    supported_extensions = {
        '.txt': '文本文件',
        '.md': 'Markdown文档',
        '.docx': 'Word文档',
        '.pdf': 'PDF文档',
        '.pptx': 'PPT演示文稿',
        '.xlsx': 'Excel表格',
        '.csv': 'CSV数据文件',
        '.py': 'Python代码',
        '.js': 'JavaScript代码',
        '.html': 'HTML文档',
        '.css': 'CSS样式',
        '.json': 'JSON数据',
        '.yaml': 'YAML配置',
        '.yml': 'YAML配置',
        '.xml': 'XML文档'
    }

    # 收集文件
    files_to_index = []
    for ext in supported_extensions.keys():
        files_to_index.extend(path.rglob(f"*{ext}"))

    if not files_to_index:
        QMessageBox.information(main_window, "索引完成", "目录中没有支持的文件类型")
        return

    # 创建进度对话框
    progress = QProgressDialog("正在索引文件...", "取消", 0, len(files_to_index), main_window)
    progress.setWindowTitle("文件索引")
    progress.setWindowModality(Qt.WindowModal)
    progress.setMinimumDuration(0)

    indexed_count = 0
    failed_count = 0

    for i, file_path in enumerate(files_to_index):
        if progress.wasCanceled():
            break

        progress.setValue(i)
        progress.setLabelText(f"正在索引: {file_path.name}")

        try:
            # 提取文件内容
            content = extract_file_content(file_path)
            if content:
                # 生成文件ID
                file_id = hashlib.md5(str(file_path).encode()).hexdigest()[:12]

                # 添加到知识库
                main_window.knowledge_base.add_knowledge(
                    id=f"file_{file_id}",
                    title=file_path.stem,
                    content=content[:2000],  # 限制内容长度
                    category=supported_extensions.get(file_path.suffix.lower(), '其他'),
                    source=str(file_path),
                    metadata={
                        'file_path': str(file_path),
                        'file_name': file_path.name,
                        'file_size': file_path.stat().st_size,
                        'modified_time': file_path.stat().st_mtime,
                        'indexed_time': pd.Timestamp.now().isoformat()
                    }
                )
                indexed_count += 1

        except Exception as e:
            from loguru import logger
            logger.error(f"索引文件失败 {file_path}: {e}")
            failed_count += 1

    progress.setValue(len(files_to_index))

    # 显示结果
    QMessageBox.information(
        main_window,
        "索引完成",
        f"索引完成!\n\n"
        f"成功: {indexed_count} 个文件\n"
        f"失败: {failed_count} 个文件\n"
        f"总计: {len(files_to_index)} 个文件"
    )

    main_window.status_bar.showMessage(f"索引完成: {indexed_count} 个文件已添加到知识库")
    main_window._refresh_kb_list()
