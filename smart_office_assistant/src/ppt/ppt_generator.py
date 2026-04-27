"""
__all__ = [
    'design_visual_hierarchy',
    'distribute_slides',
    'generate_from_content',
    'generate_ppt',
    'hex_to_rgb',
    'infer_slide_type',
    'main',
    'parse_content',
]

PPT内容generate引擎 - 将结构化内容转换为专业PPT
"""

from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, field
from enum import Enum
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import yaml
import json
from datetime import datetime
import logging
from src.core.paths import PROJECT_ROOT

logger = logging.getLogger(__name__)

class SlideType(Enum):
    """幻灯片类型"""
    TITLE = "title"           # 标题页
    TOC = "toc"               # 目录页
    CONTENT = "content"       # 内容页
    CHART = "chart"           # 图表页
    IMAGE = "image"           # 图片页
    CODE = "code"             # 代码页
    QUOTE = "quote"           # 引用页
    SUMMARY = "summary"       # 总结页
    QA = "qa"                 # Q&A页

class ContentFormat(Enum):
    """内容格式"""
    MARKDOWN = "markdown"
    YAML = "yaml"
    JSON = "json"
    DOCX = "docx"

@dataclass
class ContentNode:
    """内容节点"""
    level: int                    # 层级(0=标题, 1=章节, 2=小节, 3=要点)
    title: str                    # 标题
    content: Optional[str]        # 内容
    children: List['ContentNode'] = field(default_factory=list)
    slide_type: Optional[SlideType] = None
    metadata: Dict[str, Any] = field(default_factory=dict)

@dataclass
class SlideSpec:
    """幻灯片规范"""
    title: str
    slide_type: SlideType
    content: List[str]
    visual_elements: List[Dict[str, Any]] = field(default_factory=list)
    layout: Optional[str] = None
    metadata: Dict[str, Any] = field(default_factory=dict)

class PPTContentGenerator:
    """PPT内容generate器"""

    def __init__(self, theme: str = "business"):
        self.theme = theme
        self.theme_config = self._load_theme(theme)

    def _load_theme(self, theme: str) -> Dict[str, Any]:
        """加载主题配置"""
        config = {
            "business": {
                "colors": {
                    "primary": "#2C3E50",
                    "secondary": "#3498DB",
                    "accent": "#E74C3C",
                    "background": "#FFFFFF",
                    "text": "#2C3E50"
                },
                "fonts": {
                    "title": "Arial Bold",
                    "body": "Arial"
                },
                "layouts": {
                    "title": "Title Only",
                    "toc": "Two Content",
                    "content": "Two Content",
                    "chart": "Title and Content"
                }
            },
            "tech": {
                "colors": {
                    "primary": "#1A1A1A",
                    "secondary": "#00D4FF",
                    "accent": "#FF006E",
                    "background": "#0A0A0A",
                    "text": "#FFFFFF"
                },
                "fonts": {
                    "title": "Roboto Mono Bold",
                    "body": "Inter"
                },
                "layouts": {
                    "title": "Title Only",
                    "toc": "Two Content",
                    "content": "Two Content",
                    "chart": "Title and Content"
                }
            }
        }
        return config.get(theme, config["business"])

    def parse_content(self, content: str, format: ContentFormat = ContentFormat.MARKDOWN) -> ContentNode:
        """
        解析内容为结构化节点
        """
        if format == ContentFormat.MARKDOWN:
            return self._parse_markdown(content)
        elif format == ContentFormat.YAML:
            return self._parse_yaml(content)
        elif format == ContentFormat.JSON:
            return self._parse_json(content)
        else:
            raise ValueError(f"不支持的格式: {format}")

    def _parse_markdown(self, markdown: str) -> ContentNode:
        """
        解析Markdown内容
        """
        lines = markdown.strip().split('\n')
        root = ContentNode(level=0, title="Root", content="")

        stack = [root]

        for line in lines:
            line = line.rstrip()

            # 跳过空行
            if not line:
                continue

            # 检测标题级别
            if line.startswith('#'):
                level = line.count('#')
                title = line.lstrip('#').strip()
                node = ContentNode(level=level, title=title, content=None)

                # 找到父节点
                while stack and stack[-1].level >= level:
                    stack.pop()

                stack[-1].children.append(node)
                stack.append(node)

            # 检测列表项
            elif line.strip().startswith(('-', '*', '+')):
                content = line.strip()[1:].strip()
                node = ContentNode(level=3, title="", content=content)
                stack[-1].children.append(node)

            # 普通文本
            else:
                if stack[-1].content:
                    stack[-1].content += '\n' + line
                else:
                    stack[-1].content = line

        return root

    def _parse_yaml(self, yaml_content: str) -> ContentNode:
        """
        解析YAML内容
        """
        data = yaml.safe_load(yaml_content)
        return self._parse_dict(data)

    def _parse_json(self, json_content: str) -> ContentNode:
        """
        解析JSON内容
        """
        data = json.loads(json_content)
        return self._parse_dict(data)

    def _parse_dict(self, data: Dict[str, Any]) -> ContentNode:
        """
        解析字典结构
        """
        if "title" in data:
            root = ContentNode(level=0, title=data["title"], content=data.get("content"))
        else:
            root = ContentNode(level=0, title="Untitled", content=None)

        if "sections" in data:
            for section in data["sections"]:
                section_node = ContentNode(level=1, title=section.get("title", ""), content=section.get("content"))

                if "items" in section:
                    for item in section["items"]:
                        item_node = ContentNode(level=2, title=item.get("title", ""), content=item.get("content"))
                        section_node.children.append(item_node)

                root.children.append(section_node)

        return root

    def infer_slide_type(self, node: ContentNode) -> SlideType:
        """
        推断幻灯片类型
        """
        # 检查元数据
        if node.metadata.get("slide_type"):
            return SlideType(node.metadata["slide_type"])

        # 根据内容推断
        content_lower = node.content.lower() if node.content else ""

        # 目录页
        if "目录" in node.title or "index" in node.title.lower():
            return SlideType.TOC

        # 图表相关
        if "chart" in content_lower or "graph" in content_lower or "图" in node.title:
            return SlideType.CHART

        # 图片相关
        if "image" in content_lower or "图片" in node.title:
            return SlideType.IMAGE

        # 代码相关
        if "code" in content_lower or "```" in (node.content or ""):
            return SlideType.CODE

        # 引用相关
        if "quote" in content_lower or "引述" in node.title:
            return SlideType.QUOTE

        # 总结相关
        if "summary" in node.title.lower() or "总结" in node.title or "conclusion" in node.title.lower():
            return SlideType.SUMMARY

        # Q&A相关
        if "qa" in node.title.lower() or "问答" in node.title or "q&a" in node.title.lower():
            return SlideType.QA

        # 默认内容页
        return SlideType.CONTENT

    def distribute_slides(self, root: ContentNode) -> List[SlideSpec]:
        """
        将内容分配到幻灯片
        """
        slides = []

        # 封面页
        if root.children:
            first_section = root.children[0]
            slides.append(SlideSpec(
                title=root.title if root.title != "Root" else first_section.title,
                slide_type=SlideType.TITLE,
                content=[first_section.content] if first_section.content else [],
                layout=self.theme_config["layouts"]["title"]
            ))

        # 目录页(如果有多个章节)
        if len(root.children) > 2:
            toc_items = [child.title for child in root.children if child.title]
            slides.append(SlideSpec(
                title="目录",
                slide_type=SlideType.TOC,
                content=toc_items,
                layout=self.theme_config["layouts"]["toc"]
            ))

        # 内容页
        for section in root.children:
            # 跳过已处理的封面
            if section == root.children[0]:
                continue

            # 章节标题页
            if section.level == 1 and section.title:
                slide_type = self.infer_slide_type(section)
                slides.append(SlideSpec(
                    title=section.title,
                    slide_type=slide_type,
                    content=[section.content] if section.content else [],
                    layout=self._get_layout_for_type(slide_type)
                ))

            # 子章节/要点
            for item in section.children:
                if item.content or item.children:
                    slide_type = self.infer_slide_type(item)

                    # 收集内容
                    contents = []
                    if item.content:
                        contents.append(item.content)
                    for child in item.children:
                        if child.content:
                            contents.append(child.content)

                    if contents:
                        slides.append(SlideSpec(
                            title=item.title if item.title else section.title,
                            slide_type=slide_type,
                            content=contents,
                            layout=self._get_layout_for_type(slide_type)
                        ))

        # 总结页
        if slides and slides[-1].slide_type != SlideType.SUMMARY:
            summary_content = slides[-1].content[:3]  # 取最后3个要点
            slides.append(SlideSpec(
                title="总结",
                slide_type=SlideType.SUMMARY,
                content=summary_content,
                layout=self.theme_config["layouts"]["content"]
            ))

        # Q&A页
        slides.append(SlideSpec(
            title="Q&A",
            slide_type=SlideType.QA,
            content=[],
            layout=self.theme_config["layouts"]["title"]
        ))

        return slides

    def _get_layout_for_type(self, slide_type: SlideType) -> str:
        """get幻灯片类型对应的布局"""
        layout_map = {
            SlideType.TITLE: "Title Only",
            SlideType.TOC: "Two Content",
            SlideType.CONTENT: "Two Content",
            SlideType.CHART: "Title and Content",
            SlideType.IMAGE: "Two Content",
            SlideType.CODE: "Title and Content",
            SlideType.QUOTE: "Title Only",
            SlideType.SUMMARY: "Two Content",
            SlideType.QA: "Title Only"
        }
        return layout_map.get(slide_type, "Title and Content")

    def design_visual_hierarchy(self, slide: SlideSpec) -> Dict[str, Any]:
        """
        设计视觉层级
        """
        hierarchy = {
            "title": {
                "font_size": 32,
                "font_name": self.theme_config["fonts"]["title"],
                "color": self.theme_config["colors"]["primary"],
                "bold": True,
                "alignment": PP_ALIGN.LEFT
            },
            "body": {
                "font_size": 18,
                "font_name": self.theme_config["fonts"]["body"],
                "color": self.theme_config["colors"]["text"],
                "bold": False,
                "alignment": PP_ALIGN.LEFT
            },
            "subtitle": {
                "font_size": 24,
                "font_name": self.theme_config["fonts"]["title"],
                "color": self.theme_config["colors"]["secondary"],
                "bold": True,
                "alignment": PP_ALIGN.LEFT
            },
            "caption": {
                "font_size": 14,
                "font_name": self.theme_config["fonts"]["body"],
                "color": self.theme_config["colors"]["text"],
                "bold": False,
                "alignment": PP_ALIGN.CENTER
            }
        }

        # 根据幻灯片类型调整
        if slide.slide_type == SlideType.TITLE:
            hierarchy["title"]["font_size"] = 44
            hierarchy["title"]["alignment"] = PP_ALIGN.CENTER

        return hierarchy

    def hex_to_rgb(self, hex_color: str) -> RGBColor:
        """将HEX颜色转换为RGB"""
        hex_color = hex_color.lstrip('#')
        rgb = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        return RGBColor(rgb[0], rgb[1], rgb[2])

    def generate_ppt(self, slides: List[SlideSpec], output_path: str):
        """
        generatePPT文件
        """
        prs = Presentation()

        # 设置幻灯片尺寸(16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for slide_spec in slides:
            # 选择布局
            if slide_spec.layout == "Title Only":
                slide_layout = prs.slide_layouts[5]
            elif slide_spec.layout == "Two Content":
                slide_layout = prs.slide_layouts[3]
            else:  # Title and Content
                slide_layout = prs.slide_layouts[1]

            slide = prs.slides.add_slide(slide_layout)

            # get视觉层级
            hierarchy = self.design_visual_hierarchy(slide_spec)

            # 设置标题
            if slide.shapes.title:
                title = slide.shapes.title
                title.text = slide_spec.title
                title_para = title.text_frame.paragraphs[0]
                title_para.font.size = Pt(hierarchy["title"]["font_size"])
                title_para.font.name = hierarchy["title"]["font_name"]
                title_para.font.color.rgb = self.hex_to_rgb(hierarchy["title"]["color"])
                title_para.font.bold = hierarchy["title"]["bold"]
                title_para.alignment = hierarchy["title"]["alignment"]

            # 设置内容
            if slide_spec.content:
                # 查找内容占位符
                content_placeholder = None
                for shape in slide.placeholders:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        content_placeholder = shape
                        break

                if content_placeholder:
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()

                    for i, content in enumerate(slide_spec.content):
                        if i > 0:
                            text_frame.add_paragraph()

                        p = text_frame.paragraphs[i] if i < len(text_frame.paragraphs) else text_frame.add_paragraph()
                        p.text = content
                        p.font.size = Pt(hierarchy["body"]["font_size"])
                        p.font.name = hierarchy["body"]["font_name"]
                        p.font.color.rgb = self.hex_to_rgb(hierarchy["body"]["color"])
                        p.font.bold = hierarchy["body"]["bold"]
                        p.alignment = hierarchy["body"]["alignment"]
                        p.space_before = Pt(10)
                        p.space_after = Pt(10)

        # 保存PPT
        prs.save(output_path)
        logger.info(f"PPT已generate: {output_path}")

    def generate_from_content(self, content: str, format: ContentFormat = ContentFormat.MARKDOWN,
                             output_path: str = None) -> str:
        """
        从内容generatePPT(完整流程)
        """
        # 解析内容
        root = self.parse_content(content, format)

        # 分配幻灯片
        slides = self.distribute_slides(root)

        # generate文件名
        if output_path is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = str(PROJECT_ROOT / "outputs" / f"presentation_{timestamp}.pptx")

        # generatePPT
        self.generate_ppt(slides, output_path)

        return output_path

def main():
    """主函数 - 示例"""
    generator = PPTContentGenerator(theme="business")

    # 示例内容(Markdown格式)
    markdown_content = """
# PPTgenerate引擎演示

## 简介
- 自动将结构化内容转换为PPT
- 支持多种内容格式
- 智能幻灯片分配
- 专业主题应用

## 核心功能
- 内容解析
- 类型推断
- 视觉设计
- 文件generate

## 技术特点
- 基于python-pptx
- 支持主题定制
- 自动排版
- 格式导出

## 总结
- 高效自动化
- 专业输出
- 易于扩展
"""

    # generatePPT
    output_path = generator.generate_from_content(
        content=markdown_content,
        format=ContentFormat.MARKDOWN,
        output_path=str(PROJECT_ROOT / "outputs/demo_presentation.pptx")
    )

    print(f"PPTgenerate成功: {output_path}")

# if __name__ == "__main__":
# #     raise RuntimeError("此入口已禁用 - 请使用 tests/ 目录")
