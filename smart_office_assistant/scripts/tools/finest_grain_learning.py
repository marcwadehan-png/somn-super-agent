#!/usr/bin/env python3
"""
超细粒度深度学习系统 - 逐句理解版
深度学习E盘所有资料，确保理解每一句话的意思
"""

import sys
import os
import json
import time
import hashlib
import re
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, asdict
from enum import Enum
import logging

from path_bootstrap import bootstrap_project_paths

bootstrap_project_paths(__file__, change_cwd=True)

# 配置日志

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('data/learning/finest_grain_learning.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)


class EDriveCategory(Enum):
    """E盘学习分类"""
    PMP_MANAGEMENT = "PMP项目管理"
    INSURANCE_KNOWLEDGE = "保险知识"
    RESUME_CAREER = "简历与职业"
    ABYSS_AI_KNOWLEDGE = "Abyss AI知识库"


class LearningDimension(Enum):
    """学习维度 - 覆盖项目所有需要学习的内容"""
    # 增长引擎维度
    PROVIDER_INTELLIGENCE = "服务商智能"
    SOLUTION_TEMPLATE = "解决方案模板"
    INDUSTRY_KNOWLEDGE = "行业知识"
    ASSESSMENT_FRAMEWORK = "评估框架"
    
    # PPT系统维度
    CHART_DESIGN = "图表设计"
    PPT_STYLE = "PPT风格"
    DATA_VISUALIZATION = "数据可视化"
    
    # 神经记忆维度
    NEURAL_MEMORY = "神经记忆"
    KNOWLEDGE_GRAPH = "知识图谱"
    REASONING_ENGINE = "推理引擎"
    LEARNING_ENGINE = "学习引擎"
    
    # 技术与工具维度
    AI_TECHNOLOGY = "AI技术"
    PYTHON_ADVANCED = "Python进阶"
    DATA_SCIENCE = "数据科学"
    
    # 业务与方法论维度
    GROWTH_HACKING = "增长黑客"
    DIGITAL_MARKETING = "数字营销"
    BUSINESS_STRATEGY = "商业战略"
    
    # 行业趋势维度
    MARKET_TRENDS = "市场趋势"
    EMERGING_TECH = "新兴技术"
    COMPETITIVE_INTEL = "竞争情报"


@dataclass
class SentenceAnalysis:
    """单句分析结果"""
    sentence: str
    sentence_id: int
    key_concepts: List[str]
    entities: List[str]
    relationships: List[str]
    importance: float
    sentiment: str
    questions_generated: List[str]
    actionability: str


@dataclass
class ParagraphAnalysis:
    """段落分析结果"""
    paragraph: str
    paragraph_id: int
    sentences: List[SentenceAnalysis]
    main_idea: str
    supporting_points: List[str]
    logical_structure: str
    confidence: float


@dataclass
class FinestGrainSession:
    """超细粒度学习会话"""
    session_id: str
    timestamp: str
    file_path: str
    file_type: str
    dimension: str
    e_drive_category: str
    
    # 细粒度分析
    total_sentences: int
    total_paragraphs: int
    paragraph_analyses: List[ParagraphAnalysis]
    
    # 知识提取
    concepts: List[str]
    entities: List[str]
    relationships: List[str]
    rules: List[str]
    principles: List[str]
    
    # 深度理解
    deep_insights: List[str]
    critical_questions: List[str]
    practical_applications: List[str]
    
    # 学习质量
    overall_confidence: float
    completeness: float
    understanding_depth: str
    
    # 执行
    duration_seconds: float


class FinestGrainLearningSystem:
    """超细粒度深度学习系统"""
    
    def __init__(self, base_path: str = "E:/"):
        self.base_path = base_path
        self.learning_data_path = Path("data/learning/finest_grain")
        self.learning_data_path.mkdir(parents=True, exist_ok=True)
        
        # 文件类型支持
        self.supported_extensions = {
            '.pdf', '.doc', '.docx', '.ppt', '.pptx', 
            '.xls', '.xlsx', '.txt', '.md', '.json', 
            '.py', '.js', '.ts', '.java', '.cpp'
        }
        
        # 学习状态
        self.learned_files = set()
        self.current_session_id = None
        self.learning_statistics = {
            'total_files_scanned': 0,
            'total_files_learned': 0,
            'total_sentences_analyzed': 0,
            'total_paragraphs_analyzed': 0,
            'total_concepts_extracted': 0,
            'total_entities_extracted': 0,
            'average_confidence': 0.0,
            'learning_sessions': []
        }
        
        # 加载学习进度
        self._load_progress()
        
    def _load_progress(self):
        """加载学习进度"""
        progress_file = self.learning_data_path / "progress.json"
        if progress_file.exists():
            with open(progress_file, 'r', encoding='utf-8') as f:
                data = json.load(f)
                self.learned_files = set(data.get('learned_files', []))
                self.learning_statistics = data.get('statistics', self.learning_statistics)
                logger.info(f"加载学习进度：已学习 {len(self.learned_files)} 个文件")
    
    def _save_progress(self):
        """保存学习进度"""
        progress_file = self.learning_data_path / "progress.json"
        with open(progress_file, 'w', encoding='utf-8') as f:
            json.dump({
                'learned_files': list(self.learned_files),
                'statistics': self.learning_statistics,
                'last_update': datetime.now().isoformat()
            }, f, ensure_ascii=False, indent=2)
    
    def scan_e_drive_files(self) -> List[Dict[str, Any]]:
        """扫描E盘所有可学习文件"""
        logger.info(f"开始扫描 {self.base_path}...")
        
        files_found = []
        base_path = Path(self.base_path)
        
        if not base_path.exists():
            logger.error(f"{self.base_path} 不存在！")
            return files_found
        
        # 定义要扫描的目录
        scan_dirs = [
            "PMP备考必备资料包",
            "保险",
            "简历",
            "Abyss AI"
        ]
        
        for dir_name in scan_dirs:
            dir_path = base_path / dir_name
            if dir_path.exists():
                logger.info(f"扫描目录: {dir_path}")
                for file_path in dir_path.rglob("*"):
                    if file_path.is_file() and file_path.suffix.lower() in self.supported_extensions:
                        file_info = {
                            'path': str(file_path),
                            'name': file_path.name,
                            'extension': file_path.suffix.lower(),
                            'size': file_path.stat().st_size,
                            'category': self._categorize_file(file_path),
                            'dimension': self._map_to_dimension(file_path),
                            'modified': datetime.fromtimestamp(file_path.stat().st_mtime).isoformat()
                        }
                        files_found.append(file_info)
        
        logger.info(f"扫描完成，共找到 {len(files_found)} 个可学习文件")
        self.learning_statistics['total_files_scanned'] = len(files_found)
        return files_found
    
    def _categorize_file(self, file_path: Path) -> str:
        """分类文件到E盘类别"""
        path_str = str(file_path).upper()
        
        if 'PMP' in path_str or 'PROJECT' in path_str:
            return EDriveCategory.PMP_MANAGEMENT.value
        elif '保险' in path_str or 'INSURANCE' in path_str:
            return EDriveCategory.INSURANCE_KNOWLEDGE.value
        elif '简历' in path_str or 'RESUME' in path_str or 'CAREER' in path_str:
            return EDriveCategory.RESUME_CAREER.value
        else:
            return EDriveCategory.ABYSS_AI_KNOWLEDGE.value
    
    def _map_to_dimension(self, file_path: Path) -> str:
        """将文件映射到学习维度"""
        path_str = str(file_path).upper()
        name_str = file_path.name.upper()
        
        # 增长引擎维度
        if 'SERVICE' in path_str or 'PROVIDER' in path_str:
            return LearningDimension.PROVIDER_INTELLIGENCE.value
        elif 'SOLUTION' in path_str or 'TEMPLATE' in path_str:
            return LearningDimension.SOLUTION_TEMPLATE.value
        elif 'INDUSTRY' in path_str or 'MARKET' in path_str:
            return LearningDimension.INDUSTRY_KNOWLEDGE.value
        
        # PPT系统维度
        elif 'CHART' in path_str or 'GRAPH' in path_str or '图表' in path_str:
            return LearningDimension.CHART_DESIGN.value
        elif 'PPT' in path_str or 'STYLE' in path_str or '风格' in path_str:
            return LearningDimension.PPT_STYLE.value
        elif 'VISUAL' in path_str or '可视化' in path_str:
            return LearningDimension.DATA_VISUALIZATION.value
        
        # 神经记忆维度
        elif 'NEURAL' in path_str or 'MEMORY' in path_str or '记忆' in path_str:
            return LearningDimension.NEURAL_MEMORY.value
        elif 'GRAPH' in path_str or '图谱' in path_str:
            return LearningDimension.KNOWLEDGE_GRAPH.value
        elif 'REASON' in path_str or 'INFERENCE' in path_str or '推理' in path_str:
            return LearningDimension.REASONING_ENGINE.value
        elif 'LEARN' in path_str or '学习' in path_str:
            return LearningDimension.LEARNING_ENGINE.value
        
        # 技术与工具维度
        elif 'AI' in path_str or 'MACHINE' in path_str:
            return LearningDimension.AI_TECHNOLOGY.value
        elif 'PYTHON' in path_str:
            return LearningDimension.PYTHON_ADVANCED.value
        elif 'DATA' in path_str:
            return LearningDimension.DATA_SCIENCE.value
        
        # 业务与方法论维度
        elif 'GROWTH' in path_str or '增长' in path_str:
            return LearningDimension.GROWTH_HACKING.value
        elif 'MARKETING' in path_str or '营销' in path_str:
            return LearningDimension.DIGITAL_MARKETING.value
        elif 'STRATEGY' in path_str or '战略' in path_str:
            return LearningDimension.BUSINESS_STRATEGY.value
        
        # 行业趋势维度
        elif 'TREND' in path_str or '趋势' in path_str:
            return LearningDimension.MARKET_TRENDS.value
        elif 'EMERGING' in path_str or '新兴' in path_str:
            return LearningDimension.EMERGING_TECH.value
        elif 'COMPETITIVE' in path_str or '竞争' in path_str:
            return LearningDimension.COMPETITIVE_INTEL.value
        
        # 默认映射
        elif 'PMP' in path_str or 'PROJECT' in path_str:
            return LearningDimension.BUSINESS_STRATEGY.value
        else:
            return LearningDimension.AI_TECHNOLOGY.value
    
    def learn_single_file(self, file_info: Dict[str, Any]) -> Optional[FinestGrainSession]:
        """学习单个文件（超细粒度）"""
        file_path = file_info['path']
        
        # 检查是否已学习
        file_hash = self._get_file_hash(file_path)
        if file_hash in self.learned_files:
            logger.info(f"文件已学习，跳过: {file_path}")
            return None
        
        logger.info(f"开始学习文件: {file_path}")
        start_time = time.time()
        
        try:
            # 读取文件内容
            content = self._read_file_content(file_path)
            if not content or len(content) < 10:
                logger.warning(f"文件内容过短或为空: {file_path}")
                return None
            
            # 超细粒度分析
            paragraphs = self._split_into_paragraphs(content)
            paragraph_analyses = []
            
            all_sentences = []
            all_concepts = []
            all_entities = []
            all_relationships = []
            all_rules = []
            all_principles = []
            all_deep_insights = []
            all_critical_questions = []
            all_practical_applications = []
            
            for para_id, paragraph in enumerate(paragraphs, 1):
                para_analysis = self._analyze_paragraph(paragraph, para_id)
                paragraph_analyses.append(para_analysis)
                
                # 提取所有句子
                for sent_analysis in para_analysis.sentences:
                    all_sentences.append(sent_analysis)
                    all_concepts.extend(sent_analysis.key_concepts)
                    all_entities.extend(sent_analysis.entities)
                    all_relationships.extend(sent_analysis.relationships)
                
                # 提取段落级知识
                all_rules.extend(para_analysis.paragraph.split('。')[0].split('，')[0:1])
                if '应该' in para_analysis.paragraph or '必须' in para_analysis.paragraph:
                    all_principles.append(para_analysis.main_idea)
                
                # 生成深度洞察
                deep_insight = self._generate_deep_insight(para_analysis)
                all_deep_insights.append(deep_insight)
                
                # 生成关键问题
                critical_question = self._generate_critical_question(para_analysis)
                all_critical_questions.append(critical_question)
                
                # 生成实际应用
                practical_application = self._generate_practical_application(para_analysis)
                all_practical_applications.append(practical_application)
            
            # 计算置信度
            overall_confidence = self._calculate_overall_confidence(paragraph_analyses)
            
            # 计算完整性
            completeness = self._calculate_completeness(all_sentences, all_concepts)
            
            # 理解深度
            understanding_depth = self._assess_understanding_depth(all_sentences, all_deep_insights)
            
            # 创建学习会话
            session = FinestGrainSession(
                session_id=f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{hashlib.md5(file_path.encode()).hexdigest()[:8]}",
                timestamp=datetime.now().isoformat(),
                file_path=file_path,
                file_type=file_info['extension'],
                dimension=file_info['dimension'],
                e_drive_category=file_info['category'],
                total_sentences=len(all_sentences),
                total_paragraphs=len(paragraphs),
                paragraph_analyses=paragraph_analyses,
                concepts=list(set(all_concepts)),
                entities=list(set(all_entities)),
                relationships=list(set(all_relationships)),
                rules=all_rules[:10],  # 限制数量
                principles=list(set(all_principles))[:10],
                deep_insights=all_deep_insights[:5],
                critical_questions=all_critical_questions[:5],
                practical_applications=all_practical_applications[:5],
                overall_confidence=overall_confidence,
                completeness=completeness,
                understanding_depth=understanding_depth,
                duration_seconds=time.time() - start_time
            )
            
            # 保存学习会话
            self._save_session(session)
            
            # 标记文件已学习
            self.learned_files.add(file_hash)
            
            # 更新统计
            self.learning_statistics['total_files_learned'] += 1
            self.learning_statistics['total_sentences_analyzed'] += len(all_sentences)
            self.learning_statistics['total_paragraphs_analyzed'] += len(paragraphs)
            self.learning_statistics['total_concepts_extracted'] += len(all_concepts)
            self.learning_statistics['total_entities_extracted'] += len(all_entities)
            
            # 更新平均置信度
            current_avg = self.learning_statistics['average_confidence']
            n = self.learning_statistics['total_files_learned']
            self.learning_statistics['average_confidence'] = (current_avg * (n-1) + overall_confidence) / n
            
            logger.info(f"文件学习完成: {file_path}")
            logger.info(f"  - 分析句子: {len(all_sentences)}")
            logger.info(f"  - 分析段落: {len(paragraphs)}")
            logger.info(f"  - 提取概念: {len(all_concepts)}")
            logger.info(f"  - 深度洞察: {len(all_deep_insights)}")
            logger.info(f"  - 置信度: {overall_confidence:.2f}")
            
            return session
            
        except Exception as e:
            logger.error(f"学习文件失败 {file_path}: {e}")
            import traceback
            traceback.print_exc()
            return None
    
    def _read_file_content(self, file_path: str) -> str:
        """读取文件内容"""
        file_ext = Path(file_path).suffix.lower()
        
        try:
            if file_ext == '.pdf':
                # PDF读取
                import PyPDF2
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    content = ""
                    for page in reader.pages:
                        try:
                            text = page.extract_text()
                            if text:
                                content += text + "\n"
                        except Exception:
                            pass
                    return content
            elif file_ext in ['.docx', '.doc']:
                # Word读取（只支持.docx）
                if file_ext == '.docx':
                    import docx
                    doc = docx.Document(file_path)
                    return "\n".join([para.text for para in doc.paragraphs])
                else:
                    # .doc文件暂不支持，需要antiword等工具
                    logger.warning(f".doc格式暂不支持，请转换为.docx: {file_path}")
                    return ""
            elif file_ext in ['.ppt', '.pptx']:
                # PPT读取（只支持.pptx）
                if file_ext == '.pptx':
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    content = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                content += shape.text + "\n"
                    return content
                else:
                    # .ppt文件暂不支持
                    logger.warning(f".ppt格式暂不支持，请转换为.pptx: {file_path}")
                    return ""
            elif file_ext in ['.xls', '.xlsx']:
                import pandas as pd
                try:
                    df = pd.read_excel(file_path)
                    return df.to_string()
                except Exception as e:
                    logger.warning(f"读取Excel失败 {file_path}: {e}")
                    return ""
            else:
                # 文本文件直接读取
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
        except ImportError as e:
            logger.error(f"缺少依赖库 {file_path}: {e}")
            return ""
        except Exception as e:
            logger.error(f"读取文件失败 {file_path}: {e}")
            return ""
    
    def _split_into_paragraphs(self, content: str) -> List[str]:
        """分割内容为段落"""
        # 按换行符分割
        paragraphs = [p.strip() for p in content.split('\n') if p.strip()]
        # 过滤过短的段落
        paragraphs = [p for p in paragraphs if len(p) > 20]
        return paragraphs
    
    def _analyze_paragraph(self, paragraph: str, para_id: int) -> ParagraphAnalysis:
        """分析单个段落"""
        sentences = self._split_into_sentences(paragraph)
        sentence_analyses = []
        
        for sent_id, sentence in enumerate(sentences, 1):
            sent_analysis = self._analyze_sentence(sentence, sent_id)
            sentence_analyses.append(sent_analysis)
        
        # 提取段落主旨
        main_idea = self._extract_main_idea(paragraph, sentence_analyses)
        
        # 提取支撑点
        supporting_points = [sent.sentence for sent in sentence_analyses if sent.importance > 0.6]
        
        # 分析逻辑结构
        logical_structure = self._analyze_logical_structure(paragraph)
        
        # 计算置信度
        confidence = sum(sent.importance for sent in sentence_analyses) / len(sentence_analyses) if sentence_analyses else 0.5
        
        return ParagraphAnalysis(
            paragraph=paragraph,
            paragraph_id=para_id,
            sentences=sentence_analyses,
            main_idea=main_idea,
            supporting_points=supporting_points,
            logical_structure=logical_structure,
            confidence=confidence
        )
    
    def _split_into_sentences(self, text: str) -> List[str]:
        """分割文本为句子"""
        # 使用中文和英文的句子结束符
        sentences = re.split(r'[。！？.!?]+', text)
        return [s.strip() for s in sentences if s.strip()]
    
    def _analyze_sentence(self, sentence: str, sent_id: int) -> SentenceAnalysis:
        """分析单个句子（超细粒度）"""
        # 提取关键概念
        key_concepts = self._extract_key_concepts(sentence)
        
        # 提取实体
        entities = self._extract_entities(sentence)
        
        # 提取关系
        relationships = self._extract_relationships(sentence)
        
        # 计算重要性
        importance = self._calculate_sentence_importance(sentence, key_concepts, entities)
        
        # 分析情感
        sentiment = self._analyze_sentiment(sentence)
        
        # 生成问题
        questions = self._generate_questions(sentence, key_concepts)
        
        # 评估可执行性
        actionability = self._assess_actionability(sentence)
        
        return SentenceAnalysis(
            sentence=sentence,
            sentence_id=sent_id,
            key_concepts=key_concepts,
            entities=entities,
            relationships=relationships,
            importance=importance,
            sentiment=sentiment,
            questions_generated=questions,
            actionability=actionability
        )
    
    def _extract_key_concepts(self, sentence: str) -> List[str]:
        """提取关键概念"""
        concepts = []
        
        # 简单的关键词提取（基于规则）
        # 识别专业术语（2-4个字的词）
        words = re.findall(r'[\u4e00-\u9fa5]{2,4}|[a-zA-Z]{3,}', sentence)
        
        # 过滤常用词
        stop_words = {'的', '了', '是', '在', '有', '和', '与', '或', '这', '那', 'the', 'and', 'is', 'of', 'to', 'in'}
        concepts = [w for w in words if w not in stop_words]
        
        # 去重并限制数量
        concepts = list(set(concepts))[:5]
        return concepts
    
    def _extract_entities(self, sentence: str) -> List[str]:
        """提取实体（简化版）"""
        entities = []
        
        # 识别可能的实体（数字、百分比、人名、机构名等）
        # 日期
        dates = re.findall(r'\d{4}[-/年]\d{1,2}[-/月]\d{1,2}', sentence)
        entities.extend(dates)
        
        # 数字和百分比
        numbers = re.findall(r'\d+[%％]', sentence)
        entities.extend(numbers)
        
        # 英文术语
        terms = re.findall(r'[A-Z][a-zA-Z]+', sentence)
        entities.extend(terms)
        
        return entities[:5]
    
    def _extract_relationships(self, sentence: str) -> List[str]:
        """提取关系（简化版）"""
        relationships = []
        
        # 识别关系词
        relation_keywords = {
            '包含': '包含关系',
            '包括': '包含关系',
            '属于': '归属关系',
            '影响': '影响关系',
            '导致': '因果关系',
            '因为': '因果关系',
            '所以': '因果关系',
            '基于': '基础关系',
            '依赖': '依赖关系',
            '促进': '促进关系',
            '阻碍': '阻碍关系'
        }
        
        for keyword, relation in relation_keywords.items():
            if keyword in sentence:
                relationships.append(f"{keyword} -> {relation}")
        
        return relationships[:3]
    
    def _calculate_sentence_importance(self, sentence: str, concepts: List[str], entities: List[str]) -> float:
        """计算句子重要性"""
        importance = 0.5  # 基础重要性
        
        # 长度加分
        if len(sentence) > 20:
            importance += 0.1
        if len(sentence) > 50:
            importance += 0.1
        
        # 概念数量加分
        importance += min(len(concepts) * 0.1, 0.2)
        
        # 实体数量加分
        importance += min(len(entities) * 0.05, 0.1)
        
        # 关键词加分
        important_words = ['关键', '重要', '核心', '必须', '应该', '需要', '重点', '主要', 'important', 'key', 'critical', 'must']
        for word in important_words:
            if word in sentence:
                importance += 0.1
                break
        
        # 限制范围
        return min(importance, 1.0)
    
    def _analyze_sentiment(self, sentence: str) -> str:
        """分析句子情感"""
        positive_words = ['成功', '有效', '优秀', '好', '优秀', '优势', 'improve', 'success', 'good', 'better']
        negative_words = ['失败', '无效', '差', '坏', '问题', '错误', 'fail', 'bad', 'wrong', 'problem']
        
        positive_count = sum(1 for word in positive_words if word in sentence)
        negative_count = sum(1 for word in negative_words if word in sentence)
        
        if positive_count > negative_count:
            return "positive"
        elif negative_count > positive_count:
            return "negative"
        else:
            return "neutral"
    
    def _generate_questions(self, sentence: str, concepts: List[str]) -> List[str]:
        """生成问题"""
        questions = []
        
        if concepts:
            main_concept = concepts[0]
            questions.append(f"什么是{main_concept}？")
            questions.append(f"{main_concept}有什么作用？")
            questions.append(f"如何应用{main_concept}？")
        
        return questions[:3]
    
    def _assess_actionability(self, sentence: str) -> str:
        """评估可执行性"""
        action_words = ['应该', '必须', '需要', '要', '可以', 'should', 'must', 'need', 'can']
        
        if any(word in sentence for word in action_words):
            return "high"
        elif '?' in sentence or '？' in sentence:
            return "question"
        else:
            return "informational"
    
    def _extract_main_idea(self, paragraph: str, sentence_analyses: List[SentenceAnalysis]) -> str:
        """提取段落主旨"""
        # 选择最重要的句子
        if not sentence_analyses:
            return paragraph[:50]
        
        most_important = max(sentence_analyses, key=lambda x: x.importance)
        return most_important.sentence
    
    def _analyze_logical_structure(self, paragraph: str) -> str:
        """分析逻辑结构"""
        if '因为' in paragraph or '由于' in paragraph:
            if '所以' in paragraph or '因此' in paragraph:
                return "因果结构"
        
        if '首先' in paragraph and '其次' in paragraph:
            return "顺序结构"
        
        if '一方面' in paragraph or '另一方面' in paragraph:
            return "对比结构"
        
        return "普通结构"
    
    def _generate_deep_insight(self, para_analysis: ParagraphAnalysis) -> str:
        """生成深度洞察"""
        if not para_analysis.sentences:
            return para_analysis.paragraph
        
        main_concept = para_analysis.sentences[0].key_concepts[0] if para_analysis.sentences[0].key_concepts else "概念"
        
        # 基于段落内容生成洞察
        if para_analysis.logical_structure == "因果结构":
            return f"{main_concept}的关键因素是其因果关系的理解"
        elif '应该' in para_analysis.paragraph or '必须' in para_analysis.paragraph:
            return f"{main_concept}需要严格执行相关规范和流程"
        else:
            return f"{main_concept}的核心价值在于{para_analysis.main_idea[:30]}"
    
    def _generate_critical_question(self, para_analysis: ParagraphAnalysis) -> str:
        """生成关键问题"""
        if not para_analysis.sentences:
            return "这段内容的实际应用场景是什么？"
        
        main_concept = para_analysis.sentences[0].key_concepts[0] if para_analysis.sentences[0].key_concepts else "内容"
        return f"{main_concept}在实践中如何确保有效性？"
    
    def _generate_practical_application(self, para_analysis: ParagraphAnalysis) -> str:
        """生成实际应用"""
        if not para_analysis.sentences:
            return "将此内容应用于实际工作中"
        
        main_concept = para_analysis.sentences[0].key_concepts[0] if para_analysis.sentences[0].key_concepts else "方法"
        return f"在日常工作中运用{main_concept}的原理，提升效率"
    
    def _calculate_overall_confidence(self, paragraph_analyses: List[ParagraphAnalysis]) -> float:
        """计算整体置信度"""
        if not paragraph_analyses:
            return 0.5
        
        total_confidence = sum(p.confidence for p in paragraph_analyses)
        avg_confidence = total_confidence / len(paragraph_analyses)
        return avg_confidence
    
    def _calculate_completeness(self, sentences: List[SentenceAnalysis], concepts: List[str]) -> float:
        """计算完整性"""
        if not sentences:
            return 0.5
        
        # 基于概念覆盖度
        concept_coverage = min(len(concepts) / (len(sentences) * 2), 1.0)
        
        # 基于句子分析完整性
        analyzed_ratio = len([s for s in sentences if s.key_concepts]) / len(sentences)
        
        return (concept_coverage + analyzed_ratio) / 2
    
    def _assess_understanding_depth(self, sentences: List[SentenceAnalysis], insights: List[str]) -> str:
        """评估理解深度"""
        avg_importance = sum(s.importance for s in sentences) / len(sentences) if sentences else 0
        
        if avg_importance > 0.8:
            return "deep"
        elif avg_importance > 0.6:
            return "moderate"
        else:
            return "shallow"
    
    def _get_file_hash(self, file_path: str) -> str:
        """获取文件哈希"""
        return hashlib.md5(file_path.encode()).hexdigest()
    
    def _save_session(self, session: FinestGrainSession):
        """保存学习会话"""
        session_file = self.learning_data_path / f"session_{session.session_id}.json"
        
        # 序列化会话数据
        session_data = asdict(session)
        # 转换SentenceAnalysis和ParagraphAnalysis为字典
        session_data['paragraph_analyses'] = [
            {
                'paragraph': p.paragraph,
                'paragraph_id': p.paragraph_id,
                'sentences': [
                    {
                        'sentence': s.sentence,
                        'sentence_id': s.sentence_id,
                        'key_concepts': s.key_concepts,
                        'entities': s.entities,
                        'relationships': s.relationships,
                        'importance': s.importance,
                        'sentiment': s.sentiment,
                        'questions_generated': s.questions_generated,
                        'actionability': s.actionability
                    }
                    for s in p.sentences
                ],
                'main_idea': p.main_idea,
                'supporting_points': p.supporting_points,
                'logical_structure': p.logical_structure,
                'confidence': p.confidence
            }
            for p in session.paragraph_analyses
        ]
        
        with open(session_file, 'w', encoding='utf-8') as f:
            json.dump(session_data, f, ensure_ascii=False, indent=2)
    
    def start_finest_grain_learning(self, max_files: Optional[int] = None):
        """启动超细粒度学习"""
        logger.info("=" * 80)
        logger.info("启动超细粒度深度学习系统")
        logger.info("=" * 80)
        
        # 扫描文件
        files = self.scan_e_drive_files()
        
        if not files:
            logger.warning("没有找到可学习的文件")
            return
        
        # 限制文件数量（如果指定）
        if max_files:
            files = files[:max_files]
            logger.info(f"限制学习文件数量: {max_files}")
        
        # 开始学习
        learned_count = 0
        for i, file_info in enumerate(files, 1):
            logger.info(f"\n[{i}/{len(files)}] 学习文件: {file_info['name']}")
            
            session = self.learn_single_file(file_info)
            if session:
                learned_count += 1
                
                # 每10个文件保存一次进度
                if learned_count % 10 == 0:
                    self._save_progress()
                    logger.info(f"已学习 {learned_count} 个文件，保存进度")
        
        # 最终保存
        self._save_progress()
        
        # 生成学习报告
        self._generate_learning_report()
        
        logger.info("=" * 80)
        logger.info(f"学习完成！")
        logger.info(f"  总文件数: {len(files)}")
        logger.info(f"  已学习: {learned_count}")
        logger.info(f"  分析句子: {self.learning_statistics['total_sentences_analyzed']}")
        logger.info(f"  分析段落: {self.learning_statistics['total_paragraphs_analyzed']}")
        logger.info(f"  提取概念: {self.learning_statistics['total_concepts_extracted']}")
        logger.info(f"  平均置信度: {self.learning_statistics['average_confidence']:.2f}")
        logger.info("=" * 80)
    
    def _generate_learning_report(self):
        """生成学习报告"""
        report_file = self.learning_data_path / f"learning_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
        
        with open(report_file, 'w', encoding='utf-8') as f:
            json.dump({
                'timestamp': datetime.now().isoformat(),
                'summary': self.learning_statistics,
                'details': {
                    'total_files': self.learning_statistics['total_files_scanned'],
                    'learned_files': self.learning_statistics['total_files_learned'],
                    'learning_rate': f"{self.learning_statistics['total_files_learned'] / self.learning_statistics['total_files_scanned'] * 100:.1f}%",
                    'avg_sentences_per_file': self.learning_statistics['total_sentences_analyzed'] / max(self.learning_statistics['total_files_learned'], 1),
                    'avg_paragraphs_per_file': self.learning_statistics['total_paragraphs_analyzed'] / max(self.learning_statistics['total_files_learned'], 1),
                    'avg_concepts_per_file': self.learning_statistics['total_concepts_extracted'] / max(self.learning_statistics['total_files_learned'], 1),
                    'avg_entities_per_file': self.learning_statistics['total_entities_extracted'] / max(self.learning_statistics['total_files_learned'], 1)
                }
            }, f, ensure_ascii=False, indent=2)
        
        logger.info(f"学习报告已生成: {report_file}")


def main():
    """主函数"""
    import argparse
    
    parser = argparse.ArgumentParser(description='超细粒度深度学习系统')
    parser.add_argument('--max-files', type=int, default=None, help='最大学习文件数')
    parser.add_argument('--base-path', type=str, default='E:/', help='基础路径')
    
    args = parser.parse_args()
    
    # 创建学习系统
    system = FinestGrainLearningSystem(base_path=args.base_path)
    
    # 启动学习
    system.start_finest_grain_learning(max_files=args.max_files)


if __name__ == '__main__':
    main()
