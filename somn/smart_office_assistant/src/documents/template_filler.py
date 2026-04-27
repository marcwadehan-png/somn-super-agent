# -*- coding: utf-8 -*-
"""
智能模板填充系统 [v1.0 Phase 3]
Smart Template Fill System

功能:
- 支持Word/Excel/PPT模板的变量填充
- 支持列表/表格数据的批量填充
- 支持条件渲染和循环渲染
- 支持模板语法解析

版本: v1.0.0
日期: 2026-04-25
"""

from __future__ import annotations

import re
import json
from pathlib import Path
from typing import Dict, Any, List, Optional, Union, Callable
from dataclasses import dataclass, field
from enum import Enum, auto
import logging

logger = logging.getLogger(__name__)


class TemplateType(Enum):
    """模板类型"""
    WORD = auto()      # .docx 模板
    EXCEL = auto()     # .xlsx 模板
    POWERPOINT = auto() # .pptx 模板
    JSON = auto()      # JSON模板
    MARKDOWN = auto()  # Markdown模板
    HTML = auto()       # HTML模板
    UNKNOWN = auto()


@dataclass
class TemplateVariable:
    """模板变量定义"""
    name: str                    # 变量名 (如 {{title}})
    var_type: str = "string"     # 类型: string, number, date, list, table, conditional
    required: bool = True        # 是否必填
    default: Any = None          # 默认值
    description: str = ""        # 描述


@dataclass
class TemplateContext:
    """模板渲染上下文"""
    variables: Dict[str, Any] = field(default_factory=dict)  # 变量名 -> 值
    loops: Dict[str, List[Dict]] = field(default_factory=dict)  # 循环数据
    conditionals: Dict[str, bool] = field(default_factory=dict)  # 条件状态
    metadata: Dict[str, Any] = field(default_factory=dict)  # 元数据


@dataclass
class FillResult:
    """填充结果"""
    success: bool
    output_path: Optional[Path] = None
    data: Optional[Any] = None
    error: Optional[str] = None
    warnings: List[str] = field(default_factory=list)


class TemplateFiller:
    """
    智能模板填充器
    
    支持的模板语法:
    - {{variable}} - 简单变量
    - {{#loop name}}...{{/loop}} - 循环块
    - {{#if condition}}...{{/if}} - 条件块
    - {{@function}} - 函数调用
    
    示例:
    ```python
    filler = TemplateFiller()
    
    # 简单填充
    result = filler.fill_text(
        text="Hello, {{name}}! Today is {{date}}.",
        variables={"name": "World", "date": "2026-04-25"}
    )
    print(result)  # "Hello, World! Today is 2026-04-25."
    
    # 列表填充
    result = filler.fill_text(
        text="{{#loop items}}<li>{{name}}</li>{{/loop}}",
        variables={"items": [{"name": "Item1"}, {"name": "Item2"}]}
    )
    print(result)  # "<li>Item1</li><li>Item2</li>"
    ```
    """

    # 模板变量正则
    VAR_PATTERN = re.compile(r'\{\{([^}]+)\}\}')
    LOOP_START = re.compile(r'\{\{#loop\s+(\w+)\}\}')
    LOOP_END = re.compile(r'\{\{/loop\}\}')
    IF_START = re.compile(r'\{\{#if\s+(\w+)\}\}')
    IF_END = re.compile(r'\{\{/if\}\}')
    FUNC_PATTERN = re.compile(r'\{\{@(\w+)\}\}')
    
    # 内置函数
    BUILTIN_FUNCS: Dict[str, Callable] = {
        'upper': lambda x: str(x).upper(),
        'lower': lambda x: str(x).lower(),
        'title': lambda x: str(x).title(),
        'capitalize': lambda x: str(x).capitalize(),
        'len': lambda x: len(x) if hasattr(x, '__len__') else 1,
        'str': lambda x: str(x),
        'int': lambda x: int(x) if str(x).isdigit() else x,
        'float': lambda x: float(x) if str(x).replace('.', '').isdigit() else x,
        'date': lambda x: str(x) if hasattr(x, 'strftime') else x,
        'json': lambda x: json.dumps(x, ensure_ascii=False) if isinstance(x, (dict, list)) else x,
    }

    def __init__(self):
        """初始化模板填充器"""
        self.custom_funcs: Dict[str, Callable] = {}
        self._context_stack: List[TemplateContext] = []

    @property
    def all_funcs(self) -> Dict[str, Callable]:
        """获取所有函数（内置+自定义）"""
        return {**self.BUILTIN_FUNCS, **self.custom_funcs}

    def register_func(self, name: str, func: Callable) -> None:
        """注册自定义函数"""
        self.custom_funcs[name] = func

    def fill_text(
        self,
        text: str,
        variables: Optional[Dict[str, Any]] = None,
        context: Optional[TemplateContext] = None,
        strict: bool = False
    ) -> str:
        """
        填充文本模板
        
        Args:
            text: 模板文本
            variables: 变量字典
            context: 渲染上下文
            strict: 严格模式，未填充变量是否报错
        
        Returns:
            填充后的文本
        """
        variables = variables or {}
        context = context or TemplateContext(variables=variables)
        
        result = text
        
        # 处理循环块
        result = self._process_loops(result, variables)
        
        # 处理条件块
        result = self._process_conditionals(result, variables)
        
        # 处理变量替换
        result = self._process_variables(result, variables, strict)
        
        # 处理函数调用
        result = self._process_functions(result, variables)
        
        return result

    def _process_loops(self, text: str, variables: Dict[str, Any]) -> str:
        """处理循环块"""
        result = text
        
        while True:
            match = self.LOOP_START.search(result)
            if not match:
                break
            
            loop_name = match.group(1)
            start_pos = match.start()
            
            # 找到循环结束
            end_match = self.LOOP_END.search(result[start_pos:])
            if not end_match:
                break
            
            end_pos = start_pos + end_match.end()
            loop_content = result[start_pos + match.end():end_pos - len("{{/loop}}")]
            
            # 获取循环数据
            loop_data = variables.get(loop_name, [])
            if not isinstance(loop_data, list):
                loop_data = [loop_data]
            
            # 渲染循环内容
            rendered_items = []
            for item_data in loop_data:
                if isinstance(item_data, dict):
                    item_vars = {**variables, **item_data}
                else:
                    item_vars = {**variables, loop_name: item_data}
                rendered_items.append(self.fill_text(loop_content, item_vars))
            
            # 替换循环块
            result = result[:start_pos] + ''.join(rendered_items) + result[end_pos:]
        
        return result

    def _process_conditionals(self, text: str, variables: Dict[str, Any]) -> str:
        """处理条件块"""
        result = text
        
        while True:
            match = self.IF_START.search(result)
            if not match:
                break
            
            condition_name = match.group(1)
            start_pos = match.start()
            
            # 找到条件结束
            end_match = self.IF_END.search(result[start_pos:])
            if not end_match:
                break
            
            end_pos = start_pos + end_match.end()
            if_content = result[start_pos + match.end():end_pos - len("{{/if}}")]
            
            # 检查条件
            condition_value = variables.get(condition_name, False)
            if condition_value:
                rendered = self.fill_text(if_content, variables)
            else:
                rendered = ""
            
            # 替换条件块
            result = result[:start_pos] + rendered + result[end_pos:]
        
        return result

    def _process_variables(self, text: str, variables: Dict[str, Any], strict: bool) -> str:
        """处理变量替换"""
        def replace_var(match):
            var_name = match.group(1).strip()
            
            # 检查是否是循环变量
            if '.' in var_name:
                parts = var_name.split('.')
                value = variables
                for part in parts:
                    if isinstance(value, dict):
                        value = value.get(part)
                    else:
                        value = None
                        break
            else:
                value = variables.get(var_name)
            
            if value is not None:
                return str(value)
            elif strict:
                raise ValueError(f"未找到变量: {var_name}")
            else:
                return match.group(0)  # 保留原变量
        
        return self.VAR_PATTERN.sub(replace_var, text)

    def _process_functions(self, text: str, variables: Dict[str, Any]) -> str:
        """处理函数调用"""
        def replace_func(match):
            func_name = match.group(1).strip()
            
            if func_name in self.all_funcs:
                # 找到函数前面的变量值
                before = text[:match.start()]
                var_match = list(self.VAR_PATTERN.finditer(before))
                if var_match:
                    last_var = var_match[-1].group(1).strip()
                    value = variables.get(last_var, "")
                    try:
                        return str(self.all_funcs[func_name](value))
                    except Exception:
                        return match.group(0)
            
            return match.group(0)
        
        return self.FUNC_PATTERN.sub(replace_func, text)

    def fill_template_file(
        self,
        template_path: Union[str, Path],
        variables: Dict[str, Any],
        output_path: Optional[Union[str, Path]] = None,
        strict: bool = False
    ) -> FillResult:
        """
        填充模板文件
        
        Args:
            template_path: 模板文件路径
            variables: 变量字典
            output_path: 输出文件路径（默认覆盖原文件）
            strict: 严格模式
        
        Returns:
            FillResult: 填充结果
        """
        template_path = Path(template_path)
        
        if not template_path.exists():
            return FillResult(
                success=False,
                error=f"模板文件不存在: {template_path}"
            )
        
        try:
            # 根据文件类型处理
            suffix = template_path.suffix.lower()
            
            if suffix == '.txt':
                return self._fill_text_file(template_path, variables, output_path, strict)
            elif suffix == '.md':
                return self._fill_markdown_file(template_path, variables, output_path, strict)
            elif suffix == '.json':
                return self._fill_json_file(template_path, variables, output_path, strict)
            elif suffix == '.html':
                return self._fill_html_file(template_path, variables, output_path, strict)
            elif suffix == '.docx':
                return self._fill_docx_file(template_path, variables, output_path, strict)
            elif suffix == '.xlsx':
                return self._fill_xlsx_file(template_path, variables, output_path, strict)
            elif suffix == '.pptx':
                return self._fill_pptx_file(template_path, variables, output_path, strict)
            else:
                return FillResult(
                    success=False,
                    error=f"不支持的模板类型: {suffix}"
                )
        
        except Exception as e:
            logger.exception(f"模板填充失败: {template_path}")
            return FillResult(
                success=False,
                error=f"模板填充失败"
            )

    def _fill_text_file(
        self,
        template_path: Path,
        variables: Dict[str, Any],
        output_path: Optional[Path],
        strict: bool
    ) -> FillResult:
        """填充文本文件"""
        content = template_path.read_text(encoding='utf-8')
        result = self.fill_text(content, variables, strict=strict)
        
        output = Path(output_path) if output_path else template_path
        output.write_text(result, encoding='utf-8')
        
        return FillResult(success=True, output_path=output)

    def _fill_markdown_file(
        self,
        template_path: Path,
        variables: Dict[str, Any],
        output_path: Optional[Path],
        strict: bool
    ) -> FillResult:
        """填充Markdown文件"""
        return self._fill_text_file(template_path, variables, output_path, strict)

    def _fill_json_file(
        self,
        template_path: Path,
        variables: Dict[str, Any],
        output_path: Optional[Path],
        strict: bool
    ) -> FillResult:
        """填充JSON文件"""
        content = template_path.read_text(encoding='utf-8')
        
        # JSON模板需要特殊处理（变量语法可能冲突）
        # 使用 {{}} 包裹变量
        result = self.fill_text(content, variables, strict=strict)
        
        try:
            data = json.loads(result)
        except json.JSONDecodeError as e:
            return FillResult(success=False, error=f"JSON解析失败: {e}")
        
        output = Path(output_path) if output_path else template_path
        output.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding='utf-8')
        
        return FillResult(success=True, output_path=output, data=data)

    def _fill_html_file(
        self,
        template_path: Path,
        variables: Dict[str, Any],
        output_path: Optional[Path],
        strict: bool
    ) -> FillResult:
        """填充HTML文件"""
        return self._fill_text_file(template_path, variables, output_path, strict)

    def _fill_docx_file(
        self,
        template_path: Path,
        variables: Dict[str, Any],
        output_path: Optional[Path],
        strict: bool
    ) -> FillResult:
        """填充Word文档"""
        try:
            from docx import Document
        except ImportError:
            return FillResult(
                success=False,
                error="python-docx未安装，请运行: pip install python-docx"
            )
        
        doc = Document(template_path)
        warnings = []
        
        # 遍历所有段落
        for paragraph in doc.paragraphs:
            original = paragraph.text
            filled = self.fill_text(original, variables, strict=False)
            if filled != original:
                paragraph.text = filled
                # 检查未填充的变量
                unfilled = self.VAR_PATTERN.findall(original)
                if unfilled:
                    for var in unfilled:
                        if var not in variables:
                            warnings.append(f"段落未填充变量: {var}")
        
        # 遍历所有表格
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        original = paragraph.text
                        filled = self.fill_text(original, variables, strict=False)
                        if filled != original:
                            paragraph.text = filled
        
        output = Path(output_path) if output_path else template_path
        doc.save(output)
        
        return FillResult(success=True, output_path=output, warnings=warnings)

    def _fill_xlsx_file(
        self,
        template_path: Path,
        variables: Dict[str, Any],
        output_path: Optional[Path],
        strict: bool
    ) -> FillResult:
        """填充Excel文件"""
        try:
            from openpyxl import load_workbook
        except ImportError:
            return FillResult(
                success=False,
                error="openpyxl未安装，请运行: pip install openpyxl"
            )
        
        wb = load_workbook(template_path)
        warnings = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            
            # 遍历所有单元格
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value and isinstance(cell.value, str):
                        original = cell.value
                        filled = self.fill_text(original, variables, strict=False)
                        if filled != original:
                            cell.value = filled
        
        output = Path(output_path) if output_path else template_path
        wb.save(output)
        
        return FillResult(success=True, output_path=output, warnings=warnings)

    def _fill_pptx_file(
        self,
        template_path: Path,
        variables: Dict[str, Any],
        output_path: Optional[Path],
        strict: bool
    ) -> FillResult:
        """填充PowerPoint文件"""
        try:
            from pptx import Presentation
        except ImportError:
            return FillResult(
                success=False,
                error="python-pptx未安装，请运行: pip install python-pptx"
            )
        
        prs = Presentation(template_path)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text:
                                original = run.text
                                filled = self.fill_text(original, variables, strict=False)
                                if filled != original:
                                    run.text = filled
        
        output = Path(output_path) if output_path else template_path
        prs.save(output)
        
        return FillResult(success=True, output_path=output)


class BatchTemplateFiller:
    """
    批量模板填充器
    
    用于批量处理多个模板或同一模板的多组数据。
    
    示例:
    ```python
    batch = BatchTemplateFiller()
    
    # 添加多个填充任务
    batch.add_task(
        template="report.md",
        variables_list=[
            {"title": "年报2024", "year": 2024},
            {"title": "年报2025", "year": 2025},
        ],
        output_pattern="output/report_{year}.md"
    )
    
    results = batch.execute()
    ```
    """

    def __init__(self, filler: Optional[TemplateFiller] = None):
        self.filler = filler or TemplateFiller()
        self.tasks: List[Dict] = []
        self.results: List[FillResult] = []

    def add_task(
        self,
        template: Union[str, Path],
        variables_list: List[Dict[str, Any]],
        output_pattern: Optional[str] = None,
        output_dir: Optional[Union[str, Path]] = None
    ) -> None:
        """
        添加批量填充任务
        
        Args:
            template: 模板路径
            variables_list: 变量列表（每个元素一组变量）
            output_pattern: 输出文件命名模式，支持 {var_name} 格式
            output_dir: 输出目录
        """
        self.tasks.append({
            'template': template,
            'variables_list': variables_list,
            'output_pattern': output_pattern,
            'output_dir': Path(output_dir) if output_dir else None
        })

    def execute(self) -> List[FillResult]:
        """执行所有批量任务"""
        self.results = []
        
        for task in self.tasks:
            template = task['template']
            output_pattern = task['output_pattern']
            output_dir = task['output_dir']
            
            for variables in task['variables_list']:
                # 生成输出路径
                if output_pattern:
                    try:
                        output_path = output_pattern.format(**variables)
                        if output_dir:
                            output_path = output_dir / output_path
                        else:
                            output_path = Path(output_path)
                    except KeyError:
                        output_path = None
                else:
                    output_path = None
                
                # 执行填充
                result = self.filler.fill_template_file(
                    template,
                    variables,
                    output_path
                )
                self.results.append(result)
        
        return self.results

    def get_summary(self) -> Dict[str, Any]:
        """获取执行摘要"""
        total = len(self.results)
        success = sum(1 for r in self.results if r.success)
        failed = total - success
        
        return {
            'total': total,
            'success': success,
            'failed': failed,
            'success_rate': success / total if total > 0 else 0
        }


# 便捷函数
def quick_fill(
    template: str,
    **kwargs
) -> str:
    """
    快速填充文本模板
    
    示例:
    >>> quick_fill("Hello, {{name}}!", name="World")
    'Hello, World!'
    """
    filler = TemplateFiller()
    return filler.fill_text(template, kwargs)


def fill_file(
    template_path: Union[str, Path],
    variables: Dict[str, Any],
    output_path: Optional[Union[str, Path]] = None
) -> FillResult:
    """快速填充文件模板"""
    filler = TemplateFiller()
    return filler.fill_template_file(template_path, variables, output_path)
