"""
__all__ = [
    'add_visual_elements',
    'apply_alignment',
    'apply_color_scheme',
    'apply_fonts',
    'apply_layout',
    'apply_spacing',
    'auto_detect_layout',
    'beautify_presentation',
    'beautify_slide',
    'hex_to_rgb',
    'main',
]

PPT智能美化引擎 - 自动排版,配色,字体优化
集成高效配色矩阵generate器和品牌LOGO主色调提取
"""

from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, field
from enum import Enum
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import yaml
import logging
import random
import os

from .color_matrix_generator import ColorMatrixGenerator, Color, ColorScheme, LogoColorExtractor

logger = logging.getLogger(__name__)

class ColorScheme(Enum):
    """配色方案"""
    BUSINESS = "business"
    TECH = "tech"
    EDUCATION = "education"
    CREATIVE = "creative"
    MINIMAL = "minimal"

class LayoutType(Enum):
    """布局类型"""
    TITLE_CENTER = "title_center"          # 标题居中
    TWO_COLUMN = "two_column"              # 左右分栏
    LEFT_LIST_RIGHT_CONTENT = "left_list_right_content"  # 左列表右内容
    CHART_CENTER = "chart_center"          # 图表居中
    GRID = "grid"                          # 网格布局
    ASYMMETRIC = "asymmetric"              # 不对称布局
    WATERFALL = "waterfall"                # 瀑布流布局

@dataclass
class ColorPalette:
    """配色方案"""
    name: str
    primary: str           # 主色
    secondary: str         # 辅助色
    accent: str            # 强调色
    background: str        # 背景色
    text: str              # 文字色
    light_text: str        # 浅色文字(用于背景深色时)

@dataclass
class FontPair:
    """字体搭配"""
    name: str
    title_font: str
    body_font: str
    accent_font: Optional[str] = None
    code_font: Optional[str] = None

class PPTBeautifier:
    """PPT美化器"""

    def __init__(self, color_scheme: ColorScheme = ColorScheme.BUSINESS, 
                 brand_logo_path: Optional[str] = None):
        """
        initPPT美化器
        
        Args:
            color_scheme: 配色方案(如果未提供brand_logo_path则使用预设方案)
            brand_logo_path: 品牌LOGO图片路径(优先使用LOGO主色调)
        """
        self.brand_logo_path = brand_logo_path
        self.color_scheme = color_scheme
        
        # 如果提供了品牌LOGO,提取主色调并generate配色矩阵
        if brand_logo_path and os.path.exists(brand_logo_path):
            self.color_palette = self._generate_palette_from_logo(brand_logo_path, color_scheme)
        else:
            self.color_palette = self._load_color_palette(color_scheme)
        
        self.font_pair = self._load_font_pair(color_scheme)
        self.design_rules = self._load_design_rules()
        self.layout_patterns = self._load_layout_patterns()
    
    def _generate_palette_from_logo(self, logo_path: str, 
                                     fallback_scheme: ColorScheme) -> ColorPalette:
        """
        从品牌LOGOgenerate配色方案
        
        Args:
            logo_path: LOGO图片路径
            fallback_scheme: 如果提取失败使用的备用方案
            
        Returns:
            配色方案
        """
        try:
            # 提取LOGO主色调
            extractor = LogoColorExtractor(max_colors=3)
            dominant_colors = extractor.extract_from_image(logo_path)
            primary_color = dominant_colors[0]
            
            # 使用配色矩阵generate器generate方案
            generator = ColorMatrixGenerator(primary_color)
            
            # 根据fallback_scheme选择合适的配色方案
            if fallback_scheme == ColorScheme.BUSINESS:
                scheme = generator.generate_business_scheme()
            elif fallback_scheme == ColorScheme.TECH:
                scheme = generator.generate_tech_scheme()
            elif fallback_scheme == ColorScheme.EDUCATION:
                scheme = generator.generate_education_scheme()
            elif fallback_scheme == ColorScheme.CREATIVE:
                scheme = generator.generate_triadic_scheme()
            else:  # MINIMAL
                scheme = generator.generate_monochromatic_scheme()
            
            # 将配色方案转换为ColorPalette
            colors = scheme.colors
            name = f"品牌定制-{scheme.name}"
            
            return ColorPalette(
                name=name,
                primary=colors[0].hex_code,
                secondary=colors[1].hex_code,
                accent=colors[2].hex_code,
                background=colors[3].hex_code,
                text=colors[4].hex_code,
                light_text=colors[0].get_contrast_color().hex_code
            )
            
        except Exception as e:
            logger.warning(f"从LOGO提取配色失败,使用默认方案: {e}")
            return self._load_color_palette(fallback_scheme)

    def _load_color_palette(self, scheme: ColorScheme) -> ColorPalette:
        """加载配色方案"""
        palettes = {
            ColorScheme.BUSINESS: ColorPalette(
                name="Business",
                primary="#2C3E50",
                secondary="#3498DB",
                accent="#E74C3C",
                background="#FFFFFF",
                text="#2C3E50",
                light_text="#ECF0F1"
            ),
            ColorScheme.TECH: ColorPalette(
                name="Tech",
                primary="#1A1A1A",
                secondary="#00D4FF",
                accent="#FF006E",
                background="#0A0A0A",
                text="#FFFFFF",
                light_text="#B0B0B0"
            ),
            ColorScheme.EDUCATION: ColorPalette(
                name="Education",
                primary="#27AE60",
                secondary="#F39C12",
                accent="#E74C3C",
                background="#FDF6E3",
                text="#2C3E50",
                light_text="#7F8C8D"
            ),
            ColorScheme.CREATIVE: ColorPalette(
                name="Creative",
                primary="#8E44AD",
                secondary="#3498DB",
                accent="#F39C12",
                background="#F4F6F7",
                text="#2C3E50",
                light_text="#95A5A6"
            ),
            ColorScheme.MINIMAL: ColorPalette(
                name="Minimal",
                primary="#34495E",
                secondary="#7F8C8D",
                accent="#2980B9",
                background="#FFFFFF",
                text="#2C3E50",
                light_text="#BDC3C7"
            )
        }
        return palettes.get(scheme, palettes[ColorScheme.BUSINESS])

    def _load_font_pair(self, scheme: ColorScheme) -> FontPair:
        """加载字体搭配"""
        font_pairs = {
            ColorScheme.BUSINESS: FontPair(
                name="Business",
                title_font="Arial Bold",
                body_font="Arial",
                accent_font="Arial Italic"
            ),
            ColorScheme.TECH: FontPair(
                name="Tech",
                title_font="Roboto Mono Bold",
                body_font="Inter",
                code_font="Fira Code"
            ),
            ColorScheme.EDUCATION: FontPair(
                name="Education",
                title_font="Georgia Bold",
                body_font="Georgia",
                accent_font="Times New Roman Italic"
            ),
            ColorScheme.CREATIVE: FontPair(
                name="Creative",
                title_font="Impact",
                body_font="Arial",
                accent_font="Comic Sans MS"
            ),
            ColorScheme.MINIMAL: FontPair(
                name="Minimal",
                title_font="Helvetica Bold",
                body_font="Helvetica Regular",
                accent_font="Helvetica Italic"
            )
        }
        return font_pairs.get(scheme, font_pairs[ColorScheme.BUSINESS])

    def _load_design_rules(self) -> Dict[str, Any]:
        """加载设计规则"""
        return {
            "alignment": {
                "title": "left_aligned",
                "body": "left_aligned",
                "chart": "center_aligned",
                "quote": "centered"
            },
            "spacing": {
                "grid_system": 8,  # 8px网格系统
                "paragraph_spacing": 1.5,
                "section_padding": 2.5,
                "element_spacing": 1.0
            },
            "hierarchy": {
                "title_size": 32,
                "subtitle_size": 24,
                "body_size": 18,
                "caption_size": 14
            },
            "proportions": {
                "60_30_10": True,  # 60%主色, 30%辅助色, 10%强调色
                "max_colors_per_slide": 5,
                "whitespace_ratio": 0.2  # 至少20%留白
            }
        }

    def _load_layout_patterns(self) -> Dict[str, List[Dict[str, Any]]]:
        """加载布局模式"""
        return {
            "text_heavy": [
                {
                    "layout": LayoutType.TWO_COLUMN,
                    "condition": "bullet_points >= 6",
                    "optimization": "equal_width"
                },
                {
                    "layout": LayoutType.LEFT_LIST_RIGHT_CONTENT,
                    "condition": "bullet_points >= 4 and has_content",
                    "optimization": "30_70_split"
                }
            ],
            "chart_dominant": [
                {
                    "layout": LayoutType.CHART_CENTER,
                    "condition": "chart_area_ratio > 0.5",
                    "optimization": "maximize_chart"
                }
            ],
            "mixed_content": [
                {
                    "layout": LayoutType.ASYMMETRIC,
                    "condition": "multiple_content_types",
                    "optimization": "visual_balance"
                }
            ],
            "title_only": [
                {
                    "layout": LayoutType.TITLE_CENTER,
                    "condition": "is_title_slide",
                    "optimization": "center_align"
                }
            ]
        }

    def hex_to_rgb(self, hex_color: str) -> RGBColor:
        """将HEX颜色转换为RGB"""
        hex_color = hex_color.lstrip('#')
        rgb = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        return RGBColor(rgb[0], rgb[1], rgb[2])

    def apply_color_scheme(self, slide, is_title: bool = False):
        """
        应用配色方案到幻灯片
        """
        # 设置背景色
        if slide.background.fill.type == 1:  # No fill
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = self.hex_to_rgb(self.color_palette.background)

        # 遍历所有形状
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text_frame = shape.text_frame
            text_color = self.color_palette.light_text if self._is_dark_background() else self.color_palette.text

            # 应用到所有段落
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.color.type == 1:  # No fill
                        run.font.color.rgb = self.hex_to_rgb(text_color)

    def apply_fonts(self, slide, is_title: bool = False):
        """
        应用字体到幻灯片
        """
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text_frame = shape.text_frame

            # 标题使用标题字体
            if shape.has_text_frame and shape == slide.shapes.title and is_title:
                title_font_size = 44 if is_title else 32
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = self.font_pair.title_font
                        run.font.size = Pt(title_font_size)
                        run.font.bold = True

            # 正文使用正文字体
            else:
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name == "" or run.font.name == "Calibri":
                            run.font.name = self.font_pair.body_font
                            run.font.size = Pt(self.design_rules["hierarchy"]["body_size"])

    def apply_alignment(self, slide):
        """
        应用对齐规则
        """
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text_frame = shape.text_frame

            # judge是否是标题
            if shape == slide.shapes.title:
                alignment = PP_ALIGN.CENTER
            else:
                # 正文左对齐
                alignment = PP_ALIGN.LEFT

            # 应用对齐
            for paragraph in text_frame.paragraphs:
                paragraph.alignment = alignment

    def apply_spacing(self, slide):
        """
        应用间距规则
        """
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text_frame = shape.text_frame

            # 应用段落间距
            for paragraph in text_frame.paragraphs:
                paragraph.space_before = Pt(10)
                paragraph.space_after = Pt(10)
                paragraph.line_spacing = self.design_rules["spacing"]["paragraph_spacing"]

    def apply_layout(self, slide, layout_type: LayoutType):
        """
        应用布局
        """
        # 根据布局类型调整形状位置和大小
        if layout_type == LayoutType.TWO_COLUMN:
            self._apply_two_column_layout(slide)
        elif layout_type == LayoutType.LEFT_LIST_RIGHT_CONTENT:
            self._apply_left_list_right_content_layout(slide)
        elif layout_type == LayoutType.CHART_CENTER:
            self._apply_chart_center_layout(slide)
        elif layout_type == LayoutType.TITLE_CENTER:
            self._apply_title_center_layout(slide)

    def _apply_two_column_layout(self, slide):
        """应用左右分栏布局"""
        # get内容形状
        content_shapes = [s for s in slide.shapes if s.has_text_frame and s != slide.shapes.title]

        if len(content_shapes) >= 2:
            # 调整为左右分栏
            left_shape = content_shapes[0]
            right_shape = content_shapes[1]

            # 左栏
            left_shape.left = Inches(0.5)
            left_shape.top = Inches(2.5)
            left_shape.width = Inches(6)
            left_shape.height = Inches(4)

            # 右栏
            right_shape.left = Inches(6.8)
            right_shape.top = Inches(2.5)
            right_shape.width = Inches(6)
            right_shape.height = Inches(4)

    def _apply_left_list_right_content_layout(self, slide):
        """应用左列表右内容布局(30/70分)"""
        content_shapes = [s for s in slide.shapes if s.has_text_frame and s != slide.shapes.title]

        if len(content_shapes) >= 2:
            left_shape = content_shapes[0]
            right_shape = content_shapes[1]

            # 左栏30%
            left_shape.left = Inches(0.5)
            left_shape.top = Inches(2.5)
            left_shape.width = Inches(3.7)
            left_shape.height = Inches(4)

            # 右栏70%
            right_shape.left = Inches(4.3)
            right_shape.top = Inches(2.5)
            right_shape.width = Inches(8.5)
            right_shape.height = Inches(4)

    def _apply_chart_center_layout(self, slide):
        """应用图表居中布局"""
        content_shapes = [s for s in slide.shapes if s.has_text_frame and s != slide.shapes.title]

        if content_shapes:
            # 将内容居中
            shape = content_shapes[0]
            shape.left = Inches(2.5)
            shape.top = Inches(2.5)
            shape.width = Inches(8)
            shape.height = Inches(4)

    def _apply_title_center_layout(self, slide):
        """应用标题居中布局"""
        if slide.shapes.title:
            title = slide.shapes.title
            title.left = Inches(0)
            title.top = Inches(1)
            title.width = Inches(13.333)
            title.height = Inches(1.5)

    def add_visual_elements(self, slide):
        """
        添加视觉元素(形状,分隔线等)
        """
        # 可以添加装饰性形状,分隔线等
        pass

    def beautify_slide(self, slide, is_title: bool = False, layout_type: Optional[LayoutType] = None):
        """
        美化单个幻灯片
        """
        # 应用配色
        self.apply_color_scheme(slide, is_title)

        # 应用字体
        self.apply_fonts(slide, is_title)

        # 应用对齐
        self.apply_alignment(slide)

        # 应用间距
        self.apply_spacing(slide)

        # 应用布局
        if layout_type:
            self.apply_layout(slide, layout_type)

        # 添加视觉元素
        self.add_visual_elements(slide)

    def beautify_presentation(self, input_path: str, output_path: str):
        """
        美化整个演示文稿
        """
        prs = Presentation(input_path)

        for i, slide in enumerate(prs.slides):
            # judge是否是标题页(第一张通常是标题页)
            is_title = (i == 0)

            # 推断布局类型
            layout_type = self._infer_layout_type(slide)

            # 美化幻灯片
            self.beautify_slide(slide, is_title, layout_type)

            logger.info(f"已美化第 {i+1}/{len(prs.slides)} 张幻灯片")

        # 保存美化后的PPT
        prs.save(output_path)
        logger.info(f"PPT美化完成: {output_path}")

    def _infer_layout_type(self, slide) -> Optional[LayoutType]:
        """
        推断幻灯片布局类型
        """
        # 统计内容形状数量
        content_shapes = [s for s in slide.shapes if s.has_text_frame and s != slide.shapes.title]
        shape_count = len(content_shapes)

        # 第一张通常是标题页
        if slide.shapes.title.text == prs.slides[0].shapes.title.text:
            return LayoutType.TITLE_CENTER

        # 根据内容形状数量推断
        if shape_count >= 3:
            return LayoutType.TWO_COLUMN
        elif shape_count == 2:
            return LayoutType.LEFT_LIST_RIGHT_CONTENT
        elif shape_count == 1:
            return LayoutType.CHART_CENTER

        return None

    def _is_dark_background(self) -> bool:
        """judge是否是深色背景"""
        r = int(self.color_palette.background[1:3], 16)
        g = int(self.color_palette.background[3:5], 16)
        b = int(self.color_palette.background[5:7], 16)
        brightness = (r * 299 + g * 587 + b * 114) / 1000
        return brightness < 128

    def auto_detect_layout(self, slide) -> LayoutType:
        """
        自动检测最佳布局
        """
        # 统计内容
        bullet_points = 0
        has_content = False

        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        bullet_points += 1
                        has_content = True

        # 根据规则选择布局
        if bullet_points >= 6:
            return LayoutType.TWO_COLUMN
        elif bullet_points >= 4 and has_content:
            return LayoutType.LEFT_LIST_RIGHT_CONTENT
        else:
            return LayoutType.CHART_CENTER

def main():
    """主函数 - 示例"""
    from src.core.paths import PROJECT_ROOT
    beautifier = PPTBeautifier(color_scheme=ColorScheme.BUSINESS)

    # 美化PPT
    beautifier.beautify_presentation(
        input_path=str(PROJECT_ROOT / "outputs/demo_presentation.pptx"),
        output_path=str(PROJECT_ROOT / "outputs/demo_presentation_beautified.pptx")
    )

    print("PPT美化完成!")

# if __name__ == "__main__":
# #     raise RuntimeError("此入口已禁用 - 请使用 tests/ 目录")
