"""PPT智能generate服务 - unified服务接口(v1.0)
__all__ = [
    'apply_smart_design_rules',
    'beautify_ppt',
    'create_chart_dashboard',
    'embed_chart_to_ppt',
    'embed_multiple_charts',
    'generate_chart',
    'generate_color_schemes_from_logo',
    'generate_ppt',
    'get_chart_anti_patterns',
    'get_chart_best_practices',
    'get_chart_recommendations',
    'get_design_principles',
    'quick_beautify',
    'quick_generate',
    'recommend_color_by_mood',
    'recommend_style_by_content',
    'recommend_style_by_scene',
    'search_chart_knowledge',
]

集成图表自动generate和数据可视化
"""

from typing import Optional, Dict, Any, List, Tuple
from pathlib import Path
import logging
from datetime import datetime

from .ppt_generator import PPTContentGenerator, ContentFormat, SlideType
from .ppt_beautifier import PPTBeautifier, ColorScheme, LayoutType
from .ppt_memory_integration import PPTMemoryIntegrator
from .chart_generator import (
    ChartGenerator, ChartType, ChartLibrary, ChartConfig,
    DataAnalyzer, ChartRecommender, PlotlyChartGenerator
)
from .chart_learning import ChartLearningEngine
from .chart_embedder import ChartEmbedder
from .ppt_style_recommender import PPTStyleRecommender, SceneType
from .ppt_design_rules import SmartDesignEngine

logger = logging.getLogger(__name__)

class PPTService:
    """PPT智能generate服务(v1.0)"""

    def __init__(self, theme: str = "business", enable_learning: bool = True,
                 brand_logo_path: Optional[str] = None,
                 enable_charts: bool = True):
        """
        initPPT服务

        Args:
            theme: 主题 (business/tech/education/creative/minimal)
            enable_learning: 是否启用学习和记忆功能
            brand_logo_path: 品牌LOGO图片路径(用于generate品牌配色方案)
            enable_charts: 是否启用图表自动generate功能
        """
        self.theme = theme
        self.enable_learning = enable_learning
        self.brand_logo_path = brand_logo_path
        self.enable_charts = enable_charts

        # init各引擎
        self.generator = PPTContentGenerator(theme=theme)
        self.beautifier = PPTBeautifier(
            color_scheme=ColorScheme(theme),
            brand_logo_path=brand_logo_path
        )

        # init图表generate器
        if enable_charts:
            self.chart_generator = ChartGenerator(
                preferred_library=ChartLibrary.PLOTLY
            )
            self.chart_learning = ChartLearningEngine()
            self.chart_embedder = ChartEmbedder()
        else:
            self.chart_generator = None
            self.chart_learning = None
            self.chart_embedder = None

        # init记忆整合器
        if enable_learning:
            self.memory_integrator = PPTMemoryIntegrator()
        else:
            self.memory_integrator = None

        # initstyle推荐系统
        self.style_recommender = PPTStyleRecommender()
        self.design_engine = SmartDesignEngine()

        logger.info(f"PPT服务init完成 (主题: {theme}, 学习: {enable_learning}, "
                   f"品牌LOGO: {brand_logo_path}, 图表: {enable_charts})")

    def generate_ppt(
        self,
        content: str,
        format: ContentFormat = ContentFormat.MARKDOWN,
        output_path: str = None,
        beautify: bool = True,
        auto_charts: bool = True
    ) -> str:
        """
        generatePPT

        Args:
            content: 内容
            format: 内容格式
            output_path: 输出路径
            beautify: 是否美化
            auto_charts: 是否自动generate图表

        Returns:
            PPT文件路径
        """
        # generatePPT
        ppt_path = self.generator.generate_presentation(
            content=content,
            format=format,
            output_path=output_path
        )

        # 自动generate图表
        if auto_charts and self.chart_generator:
            self._generate_charts_from_content(content, ppt_path)

        # 美化PPT
        if beautify:
            beautified_path = ppt_path.replace('.pptx', '_beautified.pptx')
            self.beautifier.beautify_presentation(ppt_path, beautified_path)
            ppt_path = beautified_path

        logger.info(f"PPTgenerate完成: {ppt_path}")
        return ppt_path

    def generate_chart(
        self,
        data: Any,
        chart_type: Optional[ChartType] = None,
        library: Optional[ChartLibrary] = None,
        title: str = "",
        **kwargs
    ) -> tuple:
        """
        generate图表

        Args:
            data: 数据(字典,列表或DataFrame)
            chart_type: 图表类型(可选,自动推荐)
            library: 图表库(plotly/matplotlib)
            title: 图表标题
            **kwargs: 其他配置参数

        Returns:
            (图表文件路径, 推荐结果)
        """
        if not self.chart_generator:
            raise RuntimeError("图表generate功能未启用")

        return self.chart_generator.generate_chart(
            data=data,
            chart_type=chart_type,
            library=library,
            title=title,
            **kwargs
        )

    def get_chart_recommendations(self, data: Any) -> List[Dict]:
        """
        get图表推荐

        Args:
            data: 数据

        Returns:
            推荐列表
        """
        if not self.chart_generator:
            raise RuntimeError("图表generate功能未启用")

        recommendations = self.chart_generator.get_recommendations(data)

        return [
            {
                "chart_type": rec.chart_type.value,
                "library": rec.library.value,
                "confidence": rec.confidence,
                "reason": rec.reason
            }
            for rec in recommendations
        ]

    def beautify_ppt(self, input_path: str, output_path: str,
                      color_scheme: Optional[ColorScheme] = None,
                      brand_logo_path: Optional[str] = None) -> str:
        """
        美化现有PPT

        Args:
            input_path: 输入PPT路径
            output_path: 输出PPT路径
            color_scheme: 配色方案 (如果不指定则使用主题配色)
            brand_logo_path: 品牌LOGO图片路径(优先级高于color_scheme)

        Returns:
            美化后的PPT路径
        """
        if brand_logo_path:
            # 使用品牌LOGOgenerate配色
            beautifier = PPTBeautifier(
                color_scheme=color_scheme or ColorScheme(self.theme),
                brand_logo_path=brand_logo_path
            )
        elif color_scheme:
            # 使用指定配色方案
            beautifier = PPTBeautifier(color_scheme=color_scheme)
        else:
            # 使用默认美化器
            beautifier = self.beautifier

        beautifier.beautify_presentation(input_path, output_path)

        logger.info(f"PPT美化完成: {output_path}")
        return output_path

    def generate_color_schemes_from_logo(self, logo_path: str) -> Dict[str, Any]:
        """
        从品牌LOGOgenerate配色矩阵

        Args:
            logo_path: LOGO图片路径

        Returns:
            配色方案字典
        """
        from .color_matrix_generator import LogoColorExtractor, ColorMatrixGenerator

        try:
            # 提取LOGO主色调
            extractor = LogoColorExtractor(max_colors=3)
            dominant_colors = extractor.extract_from_image(logo_path)
            primary_color = dominant_colors[0]

            # generate配色矩阵
            generator = ColorMatrixGenerator(primary_color)
            matrix = generator.generate_color_matrix()

            # 转换为可序列化的格式
            result = {
                'logo_primary_color': primary_color.to_dict(),
                'all_dominant_colors': [c.to_dict() for c in dominant_colors],
                'color_matrix': {}
            }

            for category, schemes in matrix.items():
                result['color_matrix'][category] = [
                    scheme.to_dict() for scheme in schemes
                ]

            return result

        except Exception as e:
            logger.error(f"generate配色矩阵失败: {e}")
            return {'error': '配色矩阵生成失败'}

    def get_chart_best_practices(self, chart_type: str) -> List[str]:
        """
        get图表最佳实践

        Args:
            chart_type: 图表类型

        Returns:
            最佳实践列表
        """
        if not self.chart_learning:
            raise RuntimeError("图表学习功能未启用")

        return self.chart_learning.get_best_practices(chart_type)

    def get_chart_anti_patterns(self, chart_type: str) -> List[str]:
        """
        get图表反模式

        Args:
            chart_type: 图表类型

        Returns:
            反模式列表
        """
        if not self.chart_learning:
            raise RuntimeError("图表学习功能未启用")

        return self.chart_learning.get_anti_patterns(chart_type)

    def search_chart_knowledge(
        self,
        query: str,
        limit: int = 10
    ) -> List[Dict]:
        """
        搜索图表知识

        Args:
            query: 查询字符串
            limit: 返回结果限制

        Returns:
            搜索结果列表
        """
        if not self.chart_learning:
            raise RuntimeError("图表学习功能未启用")

        results = self.chart_learning.search_knowledge(query, limit=limit)

        return [
            {
                "title": result.knowledge.title,
                "content": result.knowledge.content,
                "category": result.knowledge.category.value,
                "confidence": result.knowledge.confidence,
                "relevance_score": result.relevance_score,
                "matched_fields": result.matched_fields
            }
            for result in results
        ]

    def embed_chart_to_ppt(
        self,
        ppt_path: str,
        chart_path: str,
        slide_index: int = -1,
        position: Optional[Tuple[float, float, float, float]] = None,
        title: Optional[str] = None,
        caption: Optional[str] = None
    ) -> str:
        """
        将图表嵌入到PPT中

        Args:
            ppt_path: PPT文件路径
            chart_path: 图表文件路径
            slide_index: 幻灯片索引(-1表示最后一页)
            position: 位置 (left, top, width, height)
            title: 图表标题
            caption: 图表说明

        Returns:
            嵌入后的PPT路径
        """
        if not self.chart_embedder:
            raise RuntimeError("图表嵌入功能未启用")

        return self.chart_embedder.embed_chart(
            ppt_path=ppt_path,
            chart_path=chart_path,
            slide_index=slide_index,
            position=position,
            title=title,
            caption=caption
        )

    def embed_multiple_charts(
        self,
        ppt_path: str,
        chart_configs: List[Dict[str, Any]]
    ) -> str:
        """
        批量嵌入多个图表

        Args:
            ppt_path: PPT文件路径
            chart_configs: 图表配置列表
                - chart_path: 图表路径
                - position: 位置 (left, top, width, height)
                - title: 标题
                - caption: 说明

        Returns:
            嵌入后的PPT路径
        """
        if not self.chart_embedder:
            raise RuntimeError("图表嵌入功能未启用")

        return self.chart_embedder.embed_multiple_charts(ppt_path, chart_configs)

    def create_chart_dashboard(
        self,
        chart_paths: List[str],
        output_path: str,
        layout: str = "grid",
        titles: List[str] = None
    ) -> str:
        """
        创建图表仪表板

        Args:
            chart_paths: 图表路径列表
            output_path: 输出PPT路径
            layout: 布局方式 (grid, horizontal, vertical)
            titles: 标题列表

        Returns:
            仪表板PPT路径
        """
        if not self.chart_embedder:
            raise RuntimeError("图表嵌入功能未启用")

        return self.chart_embedder.create_chart_dashboard(
            chart_paths=chart_paths,
            output_path=output_path,
            layout=layout,
            titles=titles
        )

    def _generate_charts_from_content(self, content: str, ppt_path: str) -> List[str]:
        """
        从内容中自动生成图表 [v1.0 已实现]

        Args:
            content: 内容
            ppt_path: PPT路径

        Returns:
            生成的图表路径列表
        """
        import re
        import json
        import pandas as pd
        from pathlib import Path

        chart_paths = []
        output_dir = Path(ppt_path).parent / "charts"
        output_dir.mkdir(parents=True, exist_ok=True)

        # 1. 提取表格数据 (Markdown表格格式)
        table_pattern = r'\|(.+?)\|\s*\n\|[-:\s|]+\|\s*\n((?:\|.+\|\s*\n)+)'
        tables = re.findall(table_pattern, content)

        for idx, (header, rows) in enumerate(tables):
            try:
                # 解析表头
                headers = [h.strip() for h in header.split('|') if h.strip()]
                # 解析数据行
                data_rows = []
                for row in rows.strip().split('\n'):
                    cells = [c.strip() for c in row.split('|') if c.strip()]
                    if cells:
                        data_rows.append(cells)

                if len(headers) >= 2 and len(data_rows) >= 2:
                    # 创建DataFrame
                    max_cols = max(len(headers), max(len(r) for r in data_rows))
                    headers = headers[:max_cols] + [f"Col{i}" for i in range(len(headers), max_cols)]
                    df = pd.DataFrame(data_rows, columns=headers[:max_cols])

                    # 尝试转换为数值列
                    for col in df.columns:
                        try:
                            df[col] = pd.to_numeric(df[col])
                        except Exception:
                            pass

                    # 分析数据类型并推荐图表
                    analyzer = DataAnalyzer()
                    recommender = ChartRecommender()

                    feature = analyzer.analyze_data(df)
                    recommendations = recommender.recommend(df)

                    if recommendations:
                        best_rec = recommendations[0]
                        chart_config = ChartConfig(
                            chart_type=best_rec.chart_type,
                            title=f"图表 {idx + 1}",
                            width=800,
                            height=500
                        )

                        # 尝试使用Plotly生成
                        try:
                            chart_gen = PlotlyChartGenerator()
                            chart_path = chart_gen.generate(df, chart_config, str(output_dir / f"chart_{idx}_{datetime.now().strftime('%Y%m%d%H%M%S')}.html"))
                            if chart_path and Path(chart_path).exists():
                                chart_paths.append(chart_path)
                                logger.info(f"从表格生成了图表: {chart_path}")
                        except Exception as e:
                            logger.warning(f"图表生成失败: {e}")

            except Exception as e:
                logger.warning(f"表格解析失败: {e}")

        # 2. 提取JSON数据
        json_pattern = r'```(?:json)?\s*(\{[\s\S]*?\})\s*```|`(\{[\s\S]*?\})`'
        json_matches = re.findall(json_pattern, content)
        for match in json_matches:
            json_str = match[0] or match[1]
            try:
                data = json.loads(json_str)
                if isinstance(data, dict):
                    # 尝试提取数值数组
                    for key, value in data.items():
                        if isinstance(value, list) and len(value) > 2:
                            try:
                                numeric_values = [float(v) for v in value if str(v).replace('.', '').replace('-', '').isdigit()]
                                if len(numeric_values) == len(value):
                                    df = pd.DataFrame({'指标': range(1, len(value) + 1), key: numeric_values})
                                    chart_gen = PlotlyChartGenerator()
                                    chart_path = chart_gen.generate(df, ChartConfig(chart_type=ChartType.BAR, title=key), str(output_dir / f"chart_{key}_{datetime.now().strftime('%Y%m%d%H%M%S')}.html"))
                                    if chart_path:
                                        chart_paths.append(chart_path)
                            except Exception:
                                pass
            except Exception:
                pass

        # 3. 提取键值对数据 (如 "销售额: 100万")
        kv_pattern = r'(\w+)\s*[:：]\s*([\d,.]+)\s*(?:万|亿|千|元|人|个|%)?'
        kv_matches = re.findall(kv_pattern, content)
        if len(kv_matches) >= 3:
            try:
                labels = [m[0] for m in kv_matches]
                values = [float(m[1].replace(',', '')) for m in kv_matches]
                df = pd.DataFrame({'指标': labels, '数值': values})

                chart_gen = PlotlyChartGenerator()
                chart_path = chart_gen.generate(df, ChartConfig(chart_type=ChartType.BAR, title="数据分析"), str(output_dir / f"chart_kv_{datetime.now().strftime('%Y%m%d%H%M%S')}.html"))
                if chart_path:
                    chart_paths.append(chart_path)
            except Exception as e:
                logger.warning(f"键值对图表生成失败: {e}")

        logger.info(f"从内容中自动生成了 {len(chart_paths)} 个图表")
        return chart_paths

    def recommend_style_by_scene(self, scene: str) -> Dict[str, Any]:
        """
        根据场景推荐style

        Args:
            scene: 场景描述文本

        Returns:
            style推荐字典
        """
        recommendation = self.style_recommender.recommend_by_scene(scene)

        return {
            "scene_type": recommendation.scene.value,
            "reference_platform": recommendation.platform,
            "color_scheme": recommendation.color_scheme.value,
            "color_palette": {
                "primary": recommendation.color_palette.primary,
                "secondary": recommendation.color_palette.secondary,
                "accent": recommendation.color_palette.accent,
                "background": recommendation.color_palette.background,
                "text": recommendation.color_palette.text
            },
            "layout": recommendation.layout.value,
            "fonts": {
                "chinese_title": recommendation.fonts.title,
                "chinese_body": recommendation.fonts.body,
                "english_title": recommendation.fonts.english_title,
                "english_body": recommendation.fonts.english_body
            },
            "design_rules": recommendation.design_rules,
            "confidence": recommendation.confidence
        }

    def recommend_style_by_content(self, content: str, audience: str = "") -> Dict[str, Any]:
        """
        根据内容和受众推荐style

        Args:
            content: 内容文本
            audience: 受众描述

        Returns:
            style推荐字典
        """
        recommendation = self.style_recommender.recommend_by_content(content, audience)

        return {
            "scene_type": recommendation.scene.value,
            "reference_platform": recommendation.platform,
            "color_scheme": recommendation.color_scheme.value,
            "color_palette": {
                "primary": recommendation.color_palette.primary,
                "secondary": recommendation.color_palette.secondary,
                "accent": recommendation.color_palette.accent,
                "background": recommendation.color_palette.background,
                "text": recommendation.color_palette.text
            },
            "layout": recommendation.layout.value,
            "fonts": {
                "chinese_title": recommendation.fonts.title,
                "chinese_body": recommendation.fonts.body,
                "english_title": recommendation.fonts.english_title,
                "english_body": recommendation.fonts.english_body
            },
            "design_rules": recommendation.design_rules,
            "confidence": recommendation.confidence
        }

    def recommend_color_by_mood(self, mood: str) -> Dict[str, str]:
        """
        根据情绪/氛围推荐配色

        Args:
            mood: 情绪描述

        Returns:
            配色方案字典
        """
        palette = self.style_recommender.recommend_color_by_mood(mood)

        return {
            "primary": palette.primary,
            "secondary": palette.secondary,
            "accent": palette.accent,
            "background": palette.background,
            "text": palette.text
        }

    def get_design_principles(self) -> Dict[str, str]:
        """
        get设计原则

        Returns:
            设计原则字典
        """
        return self.style_recommender.get_design_principles()

    def apply_smart_design_rules(self, ppt_path: str) -> str:
        """
        应用智能设计规则(基于Beautiful.ai原理)

        Args:
            ppt_path: PPT文件路径

        Returns:
            应用规则后的PPT路径
        """
        from pptx import Presentation

        try:
            # 加载PPT
            prs = Presentation(ppt_path)

            # 对每一页应用设计规则
            for slide in prs.slides:
                # 提取元素
                elements = self._extract_elements_from_slide(slide)

                # 自动布局
                layout = self.design_engine.auto_layout(
                    elements=elements,
                    page_width=prs.slide_width,
                    page_height=prs.slide_height
                )

                # 应用布局
                self._apply_layout_to_slide(slide, layout)

            # 保存
            output_path = ppt_path.replace(".pptx", "_smart_design.pptx")
            prs.save(output_path)

            logger.info(f"智能设计规则应用完成: {output_path}")
            return output_path

        except Exception as e:
            logger.error(f"应用智能设计规则失败: {e}")
            return ppt_path

    def _extract_elements_from_slide(self, slide) -> List:
        """
        从幻灯片中提取元素

        Args:
            slide: 幻灯片对象

        Returns:
            元素列表
        """
        from .ppt_design_rules import Element, HierarchyLevel

        elements = []

        # 提取文本框
        for shape in slide.shapes:
            if shape.has_text_frame:
                # 简化版:只提取位置和大小
                element = {
                    'type': 'text',
                    'x': shape.left,
                    'y': shape.top,
                    'width': shape.width,
                    'height': shape.height,
                    'content': shape.text
                }
                elements.append(element)

        return elements

    def _apply_layout_to_slide(self, slide, layout):
        """
        将布局应用到幻灯片

        Args:
            slide: 幻灯片对象
            layout: 布局对象
        """
        # 简化版:实际应用需要更复杂的逻辑
        # 这里只是示例

        elements = self._extract_elements_from_slide(slide)

        for i, element in enumerate(elements):
            if i < len(layout.elements):
                new_element = layout.elements[i]
                # 更新位置
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text == element['content']:
                        shape.left = new_element.x
                        shape.top = new_element.y
                        break

def quick_generate(
    content: str,
    theme: str = "business",
    output_path: str = None,
    auto_charts: bool = True
) -> str:
    """
    快速generatePPT

    Args:
        content: 内容
        theme: 主题
        output_path: 输出路径
        auto_charts: 是否自动generate图表

    Returns:
        PPT文件路径
    """
    service = PPTService(theme=theme, enable_charts=auto_charts)
    return service.generate_ppt(content=content, output_path=output_path)

def quick_beautify(
    input_path: str,
    output_path: str,
    theme: str = "business"
) -> str:
    """
    快速美化PPT

    Args:
        input_path: 输入路径
        output_path: 输出路径
        theme: 主题

    Returns:
        美化后的PPT路径
    """
    service = PPTService(theme=theme)
    return service.beautify_ppt(input_path, output_path)
