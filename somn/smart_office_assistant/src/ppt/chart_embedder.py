"""
__all__ = [
    'create_chart_dashboard',
    'embed_chart',
    'embed_multiple_charts',
]

图表嵌入PPT模块 - 智能将generate的图表嵌入到PPT中
"""

from typing import Dict, List, Any, Optional, Tuple
from pathlib import Path
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

class ChartEmbedder:
    """图表嵌入器"""

    def __init__(self):
        """init图表嵌入器"""
        pass

    def embed_chart(
        self,
        ppt_path: str,
        chart_path: str,
        slide_index: int = -1,
        position: Tuple[float, float, float, float] = None,
        title: str = None,
        caption: str = None
    ) -> str:
        """
        将图表嵌入到PPT中

        Args:
            ppt_path: PPT文件路径
            chart_path: 图表文件路径
            slide_index: 幻灯片索引(-1表示最后一页)
            position: 位置 (left, top, width, height)
            title: 图表标题
            caption: 图表说明

        Returns:
            嵌入后的PPT路径
        """
        try:
            # 加载PPT
            prs = Presentation(ppt_path)

            # 选择或创建幻灯片
            if slide_index == -1:
                slide = prs.slides[-1] if len(prs.slides) > 0 else prs.slides.add_slide(prs.slide_layouts[6])
            else:
                slide = prs.slides[slide_index]

            # 设置默认位置
            if position is None:
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(9.0)
                height = Inches(5.0)
            else:
                left, top, width, height = position
                left = Inches(left)
                top = Inches(top)
                width = Inches(width)
                height = Inches(height)

            # 添加图表
            slide.shapes.add_picture(chart_path, left, top, width, height)

            # 添加标题
            if title:
                title_box = slide.shapes.add_textbox(left, Inches(0.5), width, Inches(0.5))
                title_frame = title_box.text_frame
                title_frame.text = title
                title_para = title_frame.paragraphs[0]
                title_para.font.size = Pt(24)
                title_para.font.bold = True
                title_para.alignment = PP_ALIGN.CENTER

            # 添加说明
            if caption:
                caption_box = slide.shapes.add_textbox(left, top + height + Inches(0.2), width, Inches(0.5))
                caption_frame = caption_box.text_frame
                caption_frame.text = caption
                caption_para = caption_frame.paragraphs[0]
                caption_para.font.size = Pt(12)
                caption_para.alignment = PP_ALIGN.CENTER

            # 保存PPT
            output_path = ppt_path.replace('.pptx', '_with_chart.pptx')
            prs.save(output_path)

            logger.info(f"图表已嵌入PPT: {output_path}")
            return output_path

        except Exception as e:
            logger.error(f"嵌入图表失败: {e}")
            raise

    def embed_multiple_charts(
        self,
        ppt_path: str,
        chart_configs: List[Dict[str, Any]]
    ) -> str:
        """
        批量嵌入多个图表

        Args:
            ppt_path: PPT文件路径
            chart_configs: 图表配置列表
                - chart_path: 图表路径
                - position: 位置 (left, top, width, height)
                - title: 标题
                - caption: 说明

        Returns:
            嵌入后的PPT路径
        """
        try:
            prs = Presentation(ppt_path)

            for i, config in enumerate(chart_configs):
                # 创建新幻灯片
                slide = prs.slides.add_slide(prs.slide_layouts[6])

                # 设置位置
                position = config.get('position', (0.5, 1.5, 9.0, 5.0))
                left, top, width, height = position
                left = Inches(left)
                top = Inches(top)
                width = Inches(width)
                height = Inches(height)

                # 添加图表
                slide.shapes.add_picture(config['chart_path'], left, top, width, height)

                # 添加标题
                if config.get('title'):
                    title_box = slide.shapes.add_textbox(left, Inches(0.5), width, Inches(0.5))
                    title_frame = title_box.text_frame
                    title_frame.text = config['title']
                    title_para = title_frame.paragraphs[0]
                    title_para.font.size = Pt(24)
                    title_para.font.bold = True
                    title_para.alignment = PP_ALIGN.CENTER

                # 添加说明
                if config.get('caption'):
                    caption_box = slide.shapes.add_textbox(left, top + height + Inches(0.2), width, Inches(0.5))
                    caption_frame = caption_box.text_frame
                    caption_frame.text = config['caption']
                    caption_para = caption_frame.paragraphs[0]
                    caption_para.font.size = Pt(12)
                    caption_para.alignment = PP_ALIGN.CENTER

            # 保存PPT
            output_path = ppt_path.replace('.pptx', '_with_charts.pptx')
            prs.save(output_path)

            logger.info(f"多个图表已嵌入PPT: {output_path}")
            return output_path

        except Exception as e:
            logger.error(f"批量嵌入图表失败: {e}")
            raise

    def create_chart_dashboard(
        self,
        chart_paths: List[str],
        output_path: str,
        layout: str = "grid",
        titles: List[str] = None
    ) -> str:
        """
        创建图表仪表板

        Args:
            chart_paths: 图表路径列表
            output_path: 输出PPT路径
            layout: 布局方式 (grid, horizontal, vertical)
            titles: 标题列表

        Returns:
            仪表板PPT路径
        """
        try:
            prs = Presentation()

            if layout == "grid":
                self._create_grid_layout(prs, chart_paths, titles)
            elif layout == "horizontal":
                self._create_horizontal_layout(prs, chart_paths, titles)
            elif layout == "vertical":
                self._create_vertical_layout(prs, chart_paths, titles)
            else:
                raise ValueError(f"不支持的布局: {layout}")

            prs.save(output_path)
            logger.info(f"图表仪表板已创建: {output_path}")
            return output_path

        except Exception as e:
            logger.error(f"创建图表仪表板失败: {e}")
            raise

    def _create_grid_layout(
        self,
        prs: Presentation,
        chart_paths: List[str],
        titles: List[str] = None
    ):
        """创建网格布局"""
        charts_per_page = 4
        for i in range(0, len(chart_paths), charts_per_page):
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            page_charts = chart_paths[i:i + charts_per_page]
            page_titles = titles[i:i + charts_per_page] if titles else [None] * len(page_charts)

            for j, (chart_path, title) in enumerate(zip(page_charts, page_titles)):
                row = j // 2
                col = j % 2

                left = Inches(0.5 + col * 4.75)
                top = Inches(1.5 + row * 3.5)
                width = Inches(4.25)
                height = Inches(2.5)

                slide.shapes.add_picture(chart_path, left, top, width, height)

                if title:
                    title_box = slide.shapes.add_textbox(left, top - Inches(0.5), width, Inches(0.4))
                    title_frame = title_box.text_frame
                    title_frame.text = title
                    title_para = title_frame.paragraphs[0]
                    title_para.font.size = Pt(14)
                    title_para.font.bold = True

    def _create_horizontal_layout(
        self,
        prs: Presentation,
        chart_paths: List[str],
        titles: List[str] = None
    ):
        """创建水平布局"""
        charts_per_page = 3
        for i in range(0, len(chart_paths), charts_per_page):
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            page_charts = chart_paths[i:i + charts_per_page]
            page_titles = titles[i:i + charts_per_page] if titles else [None] * len(page_charts)

            for j, (chart_path, title) in enumerate(zip(page_charts, page_titles)):
                left = Inches(0.5 + j * 3.25)
                top = Inches(1.5)
                width = Inches(3.0)
                height = Inches(5.0)

                slide.shapes.add_picture(chart_path, left, top, width, height)

                if title:
                    title_box = slide.shapes.add_textbox(left, top - Inches(0.5), width, Inches(0.4))
                    title_frame = title_box.text_frame
                    title_frame.text = title
                    title_para = title_frame.paragraphs[0]
                    title_para.font.size = Pt(14)
                    title_para.font.bold = True

    def _create_vertical_layout(
        self,
        prs: Presentation,
        chart_paths: List[str],
        titles: List[str] = None
    ):
        """创建垂直布局"""
        charts_per_page = 2
        for i in range(0, len(chart_paths), charts_per_page):
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            page_charts = chart_paths[i:i + charts_per_page]
            page_titles = titles[i:i + charts_per_page] if titles else [None] * len(page_charts)

            for j, (chart_path, title) in enumerate(zip(page_charts, page_titles)):
                top = Inches(1.5 + j * 3.5)
                left = Inches(0.5)
                width = Inches(9.0)
                height = Inches(2.5)

                slide.shapes.add_picture(chart_path, left, top, width, height)

                if title:
                    title_box = slide.shapes.add_textbox(left, top - Inches(0.5), width, Inches(0.4))
                    title_frame = title_box.text_frame
                    title_frame.text = title
                    title_para = title_frame.paragraphs[0]
                    title_para.font.size = Pt(14)
                    title_para.font.bold = True
