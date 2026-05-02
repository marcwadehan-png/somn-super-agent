"""
================================================
    天枢 OutputEngine v1.0 — 多模态输出引擎
    Multi-Modal Output Engine for TianShu L7

    负责将八层管道的最终结果输出为多种格式：
    TEXT / MARKDOWN / HTML / IMAGE / PDF / PPTX / DOCX
================================================

Usage:
    from knowledge_cells.output_engine import (
        OutputEngine, OutputFormat, OutputArtifact,
    )

    engine = OutputEngine()
    artifact = engine.render(
        format=OutputFormat.HTML,
        final_answer="核心结论...",
        metadata={"demand_type": "分析研究", "grade": "deep"},
        sections={"结论": "...", "推理链路": "..."},
    )
    # artifact.file_path → "D:/output/xxx.html"
    # artifact.bytes → bytes of the file
"""

import os
import io
import time
import uuid
import json
import base64
from typing import Dict, List, Any, Optional, Tuple
from dataclasses import dataclass, field
from enum import Enum
from abc import ABC, abstractmethod


# ==================== 核心枚举 ====================

class OutputFormat(Enum):
    """输出格式枚举"""
    TEXT = "text"           # 纯文本
    MARKDOWN = "markdown"   # Markdown (带表格/标题)
    HTML = "html"           # HTML (可含图表/CSS)
    IMAGE = "image"         # PNG/SVG 图表
    PDF = "pdf"             # PDF 文档
    PPTX = "pptx"           # PPT 演示文稿
    DOCX = "docx"           # Word 文档


# ==================== 输出产物数据类 ====================

@dataclass
class OutputArtifact:
    """输出产物 — 统一的输出结果封装"""
    format: OutputFormat
    file_name: str = ""
    file_path: str = ""
    content: bytes = b""
    text_content: str = ""           # 纯文本回退
    size_bytes: int = 0
    duration_ms: float = 0.0
    metadata: Dict[str, Any] = field(default_factory=dict)

    def __post_init__(self):
        """自动计算 size_bytes"""
        if self.content and not self.size_bytes:
            self.size_bytes = len(self.content)

    @property
    def extension(self) -> str:
        """文件扩展名"""
        return {
            OutputFormat.TEXT: ".txt",
            OutputFormat.MARKDOWN: ".md",
            OutputFormat.HTML: ".html",
            OutputFormat.IMAGE: ".png",
            OutputFormat.PDF: ".pdf",
            OutputFormat.PPTX: ".pptx",
            OutputFormat.DOCX: ".docx",
        }[self.format]

    @property
    def mime_type(self) -> str:
        """MIME 类型"""
        return {
            OutputFormat.TEXT: "text/plain; charset=utf-8",
            OutputFormat.MARKDOWN: "text/markdown; charset=utf-8",
            OutputFormat.HTML: "text/html; charset=utf-8",
            OutputFormat.IMAGE: "image/png",
            OutputFormat.PDF: "application/pdf",
            OutputFormat.PPTX: "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            OutputFormat.DOCX: "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        }[self.format]


# ==================== 渲染上下文 ====================

@dataclass
class RenderContext:
    """渲染上下文 — 传入策略的完整数据"""
    final_answer: str = ""
    metadata: Dict[str, Any] = field(default_factory=dict)
    sections: Dict[str, Any] = field(default_factory=dict)
    confidence: float = 0.0
    grade: str = "basic"
    domain: str = "general"
    demand_type: str = "综合需求"
    decision_method: str = "标准流程"
    reasoning_chain: List[Dict] = field(default_factory=list)
    routing_path: List[str] = field(default_factory=list)
    argumentation_passed: bool = True
    argumentation_data: Dict[str, Any] = field(default_factory=dict)
    optimization_suggestions: List[str] = field(default_factory=list)
    request_id: str = ""
    total_duration_ms: float = 0.0
    output_dir: str = ""

    @classmethod
    def from_pipeline_result(cls, result: "PipelineResult", output_dir: str = "") -> "RenderContext":
        """从 PipelineResult 构建渲染上下文"""
        return cls(
            final_answer=result.final_answer,
            metadata=result.metadata,
            confidence=result.final_confidence,
            grade=result.grade.value if hasattr(result.grade, "value") else str(result.grade),
            domain=result.domain.value if hasattr(result.domain, "value") else str(result.domain),
            demand_type=result.metadata.get("demand_type", "综合需求"),
            reasoning_chain=result.reasoning_chain,
            routing_path=result.routing_path,
            argumentation_data=result.argumentation_result or {},
            optimization_suggestions=result.optimization_suggestions,
            request_id=result.request_id,
            total_duration_ms=result.total_duration_ms,
            output_dir=output_dir,
        )


# ==================== 输出策略基类 ====================

class OutputStrategy(ABC):
    """输出策略基类 — 策略模式"""

    format: OutputFormat = OutputFormat.TEXT
    display_name: str = "基础输出"

    @abstractmethod
    def render(self, ctx: RenderContext) -> bytes:
        """渲染为 bytes，子类必须实现"""
        ...

    def render_to_file(self, ctx: RenderContext) -> OutputArtifact:
        """渲染并写入文件"""
        start = time.perf_counter()
        content = self.render(ctx)
        elapsed = (time.perf_counter() - start) * 1000

        # 生成文件名
        file_id = uuid.uuid4().hex[:8]
        ts = time.strftime("%Y%m%d_%H%M%S")
        file_name = f"tianshu_{ts}_{file_id}"

        # 文本类型直接用 content
        if isinstance(content, str):
            content = content.encode("utf-8")

        file_path = ""
        if ctx.output_dir:
            os.makedirs(ctx.output_dir, exist_ok=True)
            ext = OutputArtifact(format=self.format).extension
            full_name = f"{file_name}{ext}"
            file_path = os.path.join(ctx.output_dir, full_name)
            with open(file_path, "wb") as f:
                f.write(content)

        return OutputArtifact(
            format=self.format,
            file_name=file_name,
            file_path=file_path,
            content=content,
            text_content=content.decode("utf-8", errors="replace") if len(content) < 100000 else "",
            size_bytes=len(content),
            duration_ms=round(elapsed, 2),
            metadata={"strategy": self.display_name},
        )


# ==================== 1. TEXT 纯文本策略 ====================

class TextOutputStrategy(OutputStrategy):
    """纯文本输出策略"""
    format = OutputFormat.TEXT
    display_name = "纯文本"

    def render(self, ctx: RenderContext) -> bytes:
        lines = self._build_text_lines(ctx)
        return "\n".join(lines).encode("utf-8")

    def _build_text_lines(self, ctx: RenderContext) -> List[str]:
        lines = [
            "=" * 60,
            f"  SageDispatch 八层处理管道 - {ctx.grade.upper()} 模式",
            "=" * 60,
            "",
            f"[需求类型] {ctx.demand_type}",
            f"[领域分类] {ctx.domain}",
            f"[处理等级] {ctx.grade}",
            f"[决策方法] {ctx.decision_method}",
            f"[置信度] {ctx.confidence:.1%}",
            "",
            "-" * 40,
            "核心结论",
            "-" * 40,
            ctx.final_answer or "处理完成。",
            "",
        ]

        if ctx.reasoning_chain:
            lines.append("-" * 40)
            lines.append("推理链路")
            lines.append("-" * 40)
            for i, step in enumerate(ctx.reasoning_chain, 1):
                step_name = step.get("step", f"步骤{i}")
                desc = step.get("description", "")
                lines.append(f"  {i}. {step_name}")
                if desc:
                    lines.append(f"     {desc}")
            lines.append("")

        if ctx.argumentation_data:
            lines.append("-" * 40)
            lines.append("论证结果")
            lines.append("-" * 40)
            r2 = ctx.argumentation_data.get("r2_argumentation", {})
            passed = ctx.argumentation_data.get("all_passed", True)
            if passed:
                lines.append(f"  ✓ 论证通过 (SD-R2: {r2.get('grade', 'N/A')})")
            else:
                lines.append(f"  ✗ 论证未通过 - {r2.get('recommendation', '需要修正')}")
            lines.append("")

        if ctx.optimization_suggestions:
            lines.append("-" * 40)
            lines.append("优化建议")
            lines.append("-" * 40)
            for sug in ctx.optimization_suggestions:
                lines.append(f"  • {sug}")
            lines.append("")

        lines.append("=" * 60)
        return lines


# ==================== 2. MARKDOWN 策略 ====================

class MarkdownOutputStrategy(OutputStrategy):
    """Markdown 输出策略"""
    format = OutputFormat.MARKDOWN
    display_name = "Markdown"

    def render(self, ctx: RenderContext) -> bytes:
        lines = [
            f"# SageDispatch 八层处理管道 — {ctx.grade.upper()} 模式",
            "",
            f"> 需求类型: **{ctx.demand_type}** | 领域: **{ctx.domain}** | 等级: **{ctx.grade}** | 置信度: **{ctx.confidence:.1%}**",
            "",
            "---",
            "",
            "## 核心结论",
            "",
            ctx.final_answer or "处理完成。",
            "",
        ]

        if ctx.reasoning_chain:
            lines.append("---")
            lines.append("")
            lines.append("## 推理链路")
            lines.append("")
            for i, step in enumerate(ctx.reasoning_chain, 1):
                step_name = step.get("step", f"步骤{i}")
                desc = step.get("description", "")
                lines.append(f"### {i}. {step_name}")
                if desc:
                    lines.append(f"> {desc}")
                lines.append("")

        if ctx.argumentation_data:
            lines.append("---")
            lines.append("")
            lines.append("## 论证结果")
            lines.append("")
            r2 = ctx.argumentation_data.get("r2_argumentation", {})
            passed = ctx.argumentation_data.get("all_passed", True)
            status = "✅ 通过" if passed else "❌ 未通过"
            lines.append(f"| 项目 | 结果 |")
            lines.append(f"|------|------|")
            lines.append(f"| SD-R2 谬误检测 | {r2.get('grade', 'N/A')} |")
            lines.append(f"| 论证状态 | {status} |")
            if not passed:
                lines.append(f"| 建议 | {r2.get('recommendation', '需要修正')} |")
            lines.append("")

        if ctx.optimization_suggestions:
            lines.append("---")
            lines.append("")
            lines.append("## 优化建议")
            lines.append("")
            for sug in ctx.optimization_suggestions:
                lines.append(f"- {sug}")
            lines.append("")

        if ctx.routing_path:
            lines.append("---")
            lines.append("")
            lines.append("## 调度路径")
            lines.append("")
            lines.append("`" + " → ".join(ctx.routing_path) + "`")
            lines.append("")

        return "\n".join(lines).encode("utf-8")


# ==================== 3. HTML 策略 ====================

class HtmlOutputStrategy(OutputStrategy):
    """HTML 输出策略 — 含样式和基础图表"""
    format = OutputFormat.HTML
    display_name = "HTML"

    def render(self, ctx: RenderContext) -> bytes:
        reasoning_html = self._render_reasoning(ctx.reasoning_chain)
        arg_html = self._render_argumentation(ctx.argumentation_data)
        opt_html = self._render_optimizations(ctx.optimization_suggestions)
        path_html = self._render_routing_path(ctx.routing_path)

        # 置信度颜色
        conf_color = "#22c55e" if ctx.confidence >= 0.8 else "#f59e0b" if ctx.confidence >= 0.5 else "#ef4444"
        conf_pct = f"{ctx.confidence:.1%}"

        html = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>SageDispatch - {ctx.demand_type}</title>
<style>
  :root {{ --primary: #6366f1; --bg: #0f172a; --card: #1e293b; --text: #e2e8f0; --border: #334155; }}
  * {{ margin: 0; padding: 0; box-sizing: border-box; }}
  body {{ font-family: -apple-system, 'Segoe UI', sans-serif; background: var(--bg); color: var(--text); padding: 2rem; }}
  .header {{ text-align: center; padding: 2rem 0; border-bottom: 1px solid var(--border); margin-bottom: 2rem; }}
  .header h1 {{ font-size: 1.8rem; color: var(--primary); }}
  .header .meta {{ margin-top: 0.5rem; color: #94a3b8; }}
  .meta-grid {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(180px, 1fr)); gap: 1rem; margin-bottom: 2rem; }}
  .meta-card {{ background: var(--card); border: 1px solid var(--border); border-radius: 8px; padding: 1rem; }}
  .meta-card .label {{ font-size: 0.75rem; color: #64748b; text-transform: uppercase; letter-spacing: 0.05em; }}
  .meta-card .value {{ font-size: 1.1rem; font-weight: 600; margin-top: 0.25rem; }}
  .conf-bar {{ height: 6px; background: var(--border); border-radius: 3px; margin-top: 0.5rem; }}
  .conf-fill {{ height: 100%; background: {conf_color}; border-radius: 3px; }}
  .section {{ background: var(--card); border: 1px solid var(--border); border-radius: 8px; padding: 1.5rem; margin-bottom: 1.5rem; }}
  .section h2 {{ font-size: 1.2rem; color: var(--primary); margin-bottom: 1rem; padding-bottom: 0.5rem; border-bottom: 1px solid var(--border); }}
  .section p {{ line-height: 1.7; }}
  .step {{ padding: 0.75rem 0; border-bottom: 1px solid var(--border); }}
  .step:last-child {{ border-bottom: none; }}
  .step-num {{ display: inline-block; width: 24px; height: 24px; background: var(--primary); color: white; border-radius: 50%; text-align: center; line-height: 24px; font-size: 0.75rem; margin-right: 0.5rem; }}
  .badge {{ display: inline-block; padding: 2px 8px; border-radius: 4px; font-size: 0.75rem; font-weight: 600; }}
  .badge-pass {{ background: #166534; color: #86efac; }}
  .badge-fail {{ background: #991b1b; color: #fca5a5; }}
  .opt-item {{ padding: 0.5rem 0; padding-left: 1rem; border-left: 2px solid var(--primary); margin-bottom: 0.5rem; }}
  .footer {{ text-align: center; color: #475569; font-size: 0.8rem; margin-top: 2rem; padding-top: 1rem; border-top: 1px solid var(--border); }}
</style>
</head>
<body>

<div class="header">
  <h1>SageDispatch 八层处理管道</h1>
  <div class="meta">{ctx.demand_type} | {ctx.domain} | {ctx.grade.upper()} 模式</div>
</div>

<div class="meta-grid">
  <div class="meta-card">
    <div class="label">置信度</div>
    <div class="value" style="color: {conf_color}">{conf_pct}</div>
    <div class="conf-bar"><div class="conf-fill" style="width: {ctx.confidence*100:.0f}%"></div></div>
  </div>
  <div class="meta-card">
    <div class="label">处理等级</div>
    <div class="value">{ctx.grade.upper()}</div>
  </div>
  <div class="meta-card">
    <div class="label">决策方法</div>
    <div class="value">{ctx.decision_method}</div>
  </div>
  <div class="meta-card">
    <div class="label">总耗时</div>
    <div class="value">{ctx.total_duration_ms:.1f} ms</div>
  </div>
</div>

<div class="section">
  <h2>核心结论</h2>
  <p>{ctx.final_answer or '处理完成。'}</p>
</div>

{reasoning_html}
{arg_html}
{path_html}
{opt_html}

<div class="footer">
  SageDispatch v6.1.1 | 天枢 TianShu | 请求 {ctx.request_id}
</div>

</body>
</html>"""
        return html.encode("utf-8")

    def _render_reasoning(self, chain: List[Dict]) -> str:
        if not chain:
            return ""
        steps = ""
        for i, step in enumerate(chain, 1):
            name = step.get("step", f"步骤{i}")
            desc = step.get("description", "")
            steps += f'<div class="step"><span class="step-num">{i}</span><strong>{name}</strong>'
            if desc:
                steps += f'<br><span style="color:#94a3b8;margin-left:2rem">{desc}</span>'
            steps += "</div>\n"
        return f'<div class="section"><h2>推理链路</h2>{steps}</div>'

    def _render_argumentation(self, arg_data: Dict) -> str:
        if not arg_data:
            return ""
        r2 = arg_data.get("r2_argumentation", {})
        passed = arg_data.get("all_passed", True)
        badge = "badge-pass" if passed else "badge-fail"
        status = "论证通过" if passed else "论证未通过"

        t2 = arg_data.get("t2_argumentation", {})
        t2_html = ""
        if t2:
            t2_passed = t2.get("passed", False)
            t2_badge = "badge-pass" if t2_passed else "badge-fail"
            t2_html = f'<p style="margin-top:0.5rem">RefuteCore T2: <span class="badge {t2_badge}">{"通过" if t2_passed else "未通过"}</span> '
            t2_html += f'({t2.get("strength_grade", "?")}级, {t2.get("risk_level", "?")}风险)</p>'

        r1 = arg_data.get("r1_supervision", {})
        r1_html = ""
        if r1:
            r1_html = f'<p style="margin-top:0.5rem">SD-R1 监管: {r1.get("report", "N/A")}</p>'

        return f"""<div class="section">
  <h2>论证结果</h2>
  <p>SD-R2: <span class="badge {badge}">{status}</span> ({r2.get('grade', 'N/A')})</p>
  {t2_html}{r1_html}
</div>"""

    def _render_optimizations(self, suggestions: List[str]) -> str:
        if not suggestions:
            return ""
        items = "".join(f'<div class="opt-item">{s}</div>' for s in suggestions)
        return f'<div class="section"><h2>优化建议</h2>{items}</div>'

    def _render_routing_path(self, path: List[str]) -> str:
        if not path:
            return ""
        path_str = " → ".join(f"<code>{p}</code>" for p in path)
        return f'<div class="section"><h2>调度路径</h2><p>{path_str}</p></div>'


# ==================== 4. IMAGE 图表策略 ====================

class ImageOutputStrategy(OutputStrategy):
    """图表输出策略 — 使用 matplotlib 生成信息图"""
    format = OutputFormat.IMAGE
    display_name = "图表"

    def render(self, ctx: RenderContext) -> bytes:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import matplotlib.patches as mpatches

        # 中文字体回退
        plt.rcParams["font.sans-serif"] = ["SimHei", "Microsoft YaHei", "PingFang SC", "DejaVu Sans"]
        plt.rcParams["axes.unicode_minus"] = False

        fig, ax = plt.subplots(figsize=(12, 8))
        fig.patch.set_facecolor("#0f172a")
        ax.set_facecolor("#1e293b")
        ax.axis("off")

        y = 0.95
        dy = 0.04

        # 标题
        ax.text(0.5, y, "SageDispatch 八层处理管道", fontsize=20, ha="center",
                color="#818cf8", fontweight="bold", transform=ax.transAxes)
        y -= dy * 1.8

        # 元信息行
        meta_items = [
            f"需求: {ctx.demand_type}",
            f"领域: {ctx.domain}",
            f"等级: {ctx.grade.upper()}",
            f"置信度: {ctx.confidence:.0%}",
        ]
        ax.text(0.5, y, "  |  ".join(meta_items), fontsize=10, ha="center",
                color="#94a3b8", transform=ax.transAxes)
        y -= dy * 1.5

        # 分隔线
        ax.plot([0.05, 0.95], [y + dy * 0.5, y + dy * 0.5], color="#334155",
                linewidth=0.5, transform=ax.transAxes)

        # 核心结论
        ax.text(0.05, y, "核心结论", fontsize=13, color="#818cf8", fontweight="bold",
                transform=ax.transAxes)
        y -= dy
        # 自动换行
        answer = ctx.final_answer or "处理完成。"
        if len(answer) > 120:
            answer = answer[:120] + "..."
        ax.text(0.05, y, answer, fontsize=10, color="#e2e8f0", transform=ax.transAxes,
                wrap=True)
        y -= dy * 2

        # 推理链路
        if ctx.reasoning_chain:
            ax.text(0.05, y, "推理链路", fontsize=13, color="#818cf8", fontweight="bold",
                    transform=ax.transAxes)
            y -= dy
            for i, step in enumerate(ctx.reasoning_chain[:6], 1):  # 最多6步
                name = step.get("step", f"步骤{i}")
                ax.text(0.07, y, f"{i}. {name}", fontsize=9, color="#cbd5e1",
                        transform=ax.transAxes)
                y -= dy * 0.8
            y -= dy * 0.5

        # 论证结果
        if ctx.argumentation_data:
            passed = ctx.argumentation_data.get("all_passed", True)
            color = "#22c55e" if passed else "#ef4444"
            status = "论证通过" if passed else "论证未通过"
            ax.text(0.05, y, f"论证: {status}", fontsize=10, color=color,
                    transform=ax.transAxes)
            y -= dy

        # 页脚
        ax.text(0.5, 0.02, f"SageDispatch v6.1.1 | {ctx.request_id}",
                fontsize=8, ha="center", color="#475569", transform=ax.transAxes)

        plt.tight_layout()
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                    facecolor=fig.get_facecolor())
        plt.close(fig)
        buf.seek(0)
        return buf.getvalue()


# ==================== 5. PDF 策略 ====================

class PdfOutputStrategy(OutputStrategy):
    """PDF 输出策略 — 使用 reportlab 生成"""
    format = OutputFormat.PDF
    display_name = "PDF报告"

    def render(self, ctx: RenderContext) -> bytes:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import mm
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
            HRFlowable, KeepTogether,
        )
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont

        buf = io.BytesIO()

        # 注册中文字体（尝试常见字体路径）
        _register_chinese_font()

        doc = SimpleDocTemplate(buf, pagesize=A4,
                                leftMargin=20 * mm, rightMargin=20 * mm,
                                topMargin=20 * mm, bottomMargin=20 * mm)

        styles = getSampleStyleSheet()

        # 自定义样式
        title_style = ParagraphStyle(
            "CustomTitle", parent=styles["Title"],
            fontSize=18, textColor=colors.HexColor("#4338ca"),
            spaceAfter=4 * mm, alignment=1,
        )
        subtitle_style = ParagraphStyle(
            "Subtitle", parent=styles["Normal"],
            fontSize=10, textColor=colors.HexColor("#6b7280"),
            alignment=1, spaceAfter=6 * mm,
        )
        heading_style = ParagraphStyle(
            "Heading", parent=styles["Heading2"],
            fontSize=14, textColor=colors.HexColor("#4338ca"),
            spaceBefore=4 * mm, spaceAfter=2 * mm,
        )
        body_style = ParagraphStyle(
            "Body", parent=styles["Normal"],
            fontSize=10, leading=16,
            spaceAfter=2 * mm,
        )
        meta_data = [
            ["需求类型", ctx.demand_type],
            ["领域分类", ctx.domain],
            ["处理等级", ctx.grade.upper()],
            ["决策方法", ctx.decision_method],
            ["置信度", f"{ctx.confidence:.1%}"],
            ["总耗时", f"{ctx.total_duration_ms:.1f} ms"],
        ]
        meta_table = Table(meta_data, colWidths=[35 * mm, 100 * mm])
        meta_table.setStyle(TableStyle([
            ("FONTNAME", (0, 0), (-1, -1), "Helvetica"),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("BACKGROUND", (0, 0), (0, -1), colors.HexColor("#f0f0ff")),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#d0d0d0")),
            ("TOPPADDING", (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ]))

        elements = [
            Paragraph("SageDispatch 八层处理管道", title_style),
            Paragraph(f"{ctx.demand_type} | {ctx.domain} | {ctx.grade.upper()} 模式", subtitle_style),
            HRFlowable(width="100%", thickness=1, color=colors.HexColor("#d0d0d0")),
            Spacer(1, 3 * mm),
            meta_table,
            Spacer(1, 6 * mm),
        ]

        # 核心结论
        elements.append(Paragraph("核心结论", heading_style))
        answer = ctx.final_answer or "处理完成。"
        elements.append(Paragraph(answer.replace("\n", "<br/>"), body_style))

        # 推理链路
        if ctx.reasoning_chain:
            elements.append(Spacer(1, 4 * mm))
            elements.append(Paragraph("推理链路", heading_style))
            chain_data = [["#", "步骤", "描述"]]
            for i, step in enumerate(ctx.reasoning_chain, 1):
                chain_data.append([
                    str(i),
                    step.get("step", f"步骤{i}"),
                    step.get("description", ""),
                ])
            chain_table = Table(chain_data, colWidths=[10 * mm, 35 * mm, 90 * mm])
            chain_table.setStyle(TableStyle([
                ("FONTNAME", (0, 0), (-1, -1), "Helvetica"),
                ("FONTSIZE", (0, 0), (-1, -1), 9),
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#4338ca")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#d0d0d0")),
                ("TOPPADDING", (0, 0), (-1, -1), 3),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
            ]))
            elements.append(chain_table)

        # 论证结果
        if ctx.argumentation_data:
            elements.append(Spacer(1, 4 * mm))
            elements.append(Paragraph("论证结果", heading_style))
            r2 = ctx.argumentation_data.get("r2_argumentation", {})
            passed = ctx.argumentation_data.get("all_passed", True)
            status_text = "通过" if passed else "未通过"
            status_color = colors.HexColor("#16a34a") if passed else colors.HexColor("#dc2626")
            elements.append(Paragraph(
                f'论证状态: <font color="{status_color.hexval()}">{status_text}</font> '
                f"(SD-R2: {r2.get('grade', 'N/A')})",
                body_style,
            ))

        # 优化建议
        if ctx.optimization_suggestions:
            elements.append(Spacer(1, 4 * mm))
            elements.append(Paragraph("优化建议", heading_style))
            for sug in ctx.optimization_suggestions:
                elements.append(Paragraph(f"• {sug}", body_style))

        # 页脚
        elements.append(Spacer(1, 10 * mm))
        elements.append(HRFlowable(width="100%", thickness=0.5, color=colors.HexColor("#d0d0d0")))
        elements.append(Paragraph(
            f"SageDispatch v6.1.1 | 请求 {ctx.request_id}",
            ParagraphStyle("Footer", parent=styles["Normal"], fontSize=8, textColor=colors.HexColor("#9ca3af")),
        ))

        doc.build(elements)
        buf.seek(0)
        return buf.getvalue()


# ==================== 6. PPTX 策略（开物 Kaiwu 增强） ====================

class PptxOutputStrategy(OutputStrategy):
    """
    PPT 演示文稿输出策略 [v1.1 — Kaiwu 开物增强]

    v1.1 增强：
    - 优先使用 KaiwuService（有风格学习 + Markdown 内容理解）
    - KaiwuService 不可用时回退到原生 python-pptx 实现
    - 自动将 RenderContext 转换为 Markdown 格式内容
    """
    format = OutputFormat.PPTX
    display_name = "PPT演示"

    def render(self, ctx: RenderContext) -> bytes:
        # === 尝试使用 KaiwuService（开物增强） ===
        kaiwu_error = None
        try:
            from smart_office_assistant.src.kaiwu.kaiwu_service import KaiwuService
            kaiwu = KaiwuService(theme="business", enable_learning=True)
            # 将 RenderContext 转换为 Markdown 内容
            md_content = self._render_context_to_markdown(ctx)
            output_path = kaiwu.generate_ppt(
                content=md_content,
                beautify=True,
                auto_charts=True,
            )
            with open(output_path, "rb") as f:
                return f.read()
        except ImportError:
            kaiwu_error = "KaiwuService 未安装（smart_office_assistant 未在路径中）"
        except Exception as e:
            kaiwu_error = f"KaiwuService 生成失败：{e}"

        # === 回退到原生 python-pptx 实现 ===
        if kaiwu_error:
            import logging
            logging.getLogger(__name__).warning(
                f"[OutputEngine] KaiwuService 不可用，回退到原生 PPTX：{kaiwu_error}"
            )
        return self._render_pptx_native(ctx)

    def _render_context_to_markdown(self, ctx: RenderContext) -> str:
        """将 RenderContext 渲染为 Markdown 格式（KaiwuService 输入格式）"""
        lines = [
            "# SageDispatch 八层处理管道",
            f"## {ctx.demand_type} | {ctx.domain} | {ctx.grade.upper()} 模式",
            "---",
            "## 核心结论",
            ctx.final_answer or "处理完成。",
        ]
        if ctx.reasoning_chain:
            lines += ["---", "## 推理链路"]
            for i, step in enumerate(ctx.reasoning_chain[:8], 1):
                name = step.get("step", f"步骤{i}")
                desc = step.get("description", "")
                lines.append(f"{i}. **{name}**")
                if desc:
                    lines.append(f"   {desc}")
        if ctx.argumentation_data:
            lines += ["---", "## 论证结果"]
            r2 = ctx.argumentation_data.get("r2_argumentation", {})
            passed = ctx.argumentation_data.get("all_passed", True)
            status = "✅ 通过" if passed else "❌ 未通过"
            lines.append(f"- 论证状态：{status}（{r2.get('grade', 'N/A')}）")
        if ctx.optimization_suggestions:
            lines += ["---", "## 优化建议"]
            for sug in ctx.optimization_suggestions:
                lines.append(f"- {sug}")
        return "\n\n".join(lines)

    def _render_pptx_native(self, ctx: RenderContext) -> bytes:
        """原生 python-pptx 实现（回退方案）"""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # 颜色常量
        BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
        PRIMARY = RGBColor(0x63, 0x66, 0xF1)
        TEXT_LIGHT = RGBColor(0xE2, 0xE8, 0xF0)
        TEXT_DIM = RGBColor(0x94, 0xA3, 0xB8)

        # === Slide 1: 封面 ===
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK

        # 标题
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "SageDispatch 八层处理管道"
        p.font.size = Pt(40)
        p.font.color.rgb = PRIMARY
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # 副标题
        p2 = tf.add_paragraph()
        p2.text = f"{ctx.demand_type} | {ctx.domain} | {ctx.grade.upper()} 模式"
        p2.font.size = Pt(18)
        p2.font.color.rgb = TEXT_DIM
        p2.alignment = PP_ALIGN.CENTER

        # === Slide 2: 核心结论 ===
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])
        bg2 = slide2.background.fill
        bg2.solid()
        bg2.fore_color.rgb = BG_DARK

        # 标题
        txBox2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf2 = txBox2.text_frame
        p2_title = tf2.paragraphs[0]
        p2_title.text = "核心结论"
        p2_title.font.size = Pt(28)
        p2_title.font.color.rgb = PRIMARY
        p2_title.font.bold = True

        # 结论内容
        answer = ctx.final_answer or "处理完成。"
        txBox2b = slide2.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
        tf2b = txBox2b.text_frame
        tf2b.word_wrap = True
        p2_body = tf2b.paragraphs[0]
        p2_body.text = answer
        p2_body.font.size = Pt(16)
        p2_body.font.color.rgb = TEXT_LIGHT
        p2_body.line_spacing = Pt(28)

        # === Slide 3: 分析详情 ===
        slide3 = prs.slides.add_slide(prs.slide_layouts[6])
        bg3 = slide3.background.fill
        bg3.solid()
        bg3.fore_color.rgb = BG_DARK

        txBox3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf3 = txBox3.text_frame
        p3_title = tf3.paragraphs[0]
        p3_title.text = "分析详情"
        p3_title.font.size = Pt(28)
        p3_title.font.color.rgb = PRIMARY
        p3_title.font.bold = True

        # 元信息表格
        from pptx.util import Inches as In
        table_shape = slide3.shapes.add_table(6, 2, Inches(0.8), Inches(1.3), Inches(5), Inches(3))
        table = table_shape.table
        meta_rows = [
            ("需求类型", ctx.demand_type),
            ("领域分类", ctx.domain),
            ("处理等级", ctx.grade.upper()),
            ("决策方法", ctx.decision_method),
            ("置信度", f"{ctx.confidence:.1%}"),
            ("总耗时", f"{ctx.total_duration_ms:.1f} ms"),
        ]
        for i, (label, value) in enumerate(meta_rows):
            cell_l = table.cell(i, 0)
            cell_r = table.cell(i, 1)
            cell_l.text = label
            cell_r.text = value
            for cell in (cell_l, cell_r):
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(12)
                    p.font.color.rgb = TEXT_LIGHT
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)

        # 论证结果
        if ctx.argumentation_data:
            passed = ctx.argumentation_data.get("all_passed", True)
            txBox3b = slide3.shapes.add_textbox(Inches(6.5), Inches(1.3), Inches(6), Inches(3))
            tf3b = txBox3b.text_frame
            tf3b.word_wrap = True
            p3_arg = tf3b.paragraphs[0]
            p3_arg.text = "论证结果"
            p3_arg.font.size = Pt(20)
            p3_arg.font.color.rgb = PRIMARY
            p3_arg.font.bold = True

            p3_status = tf3b.add_paragraph()
            p3_status.text = "通过" if passed else "未通过"
            p3_status.font.size = Pt(16)
            p3_status.font.color.rgb = RGBColor(0x22, 0xC5, 0x5E) if passed else RGBColor(0xEF, 0x44, 0x44)

        # === Slide 4: 推理链路 (如果有) ===
        if ctx.reasoning_chain:
            slide4 = prs.slides.add_slide(prs.slide_layouts[6])
            bg4 = slide4.background.fill
            bg4.solid()
            bg4.fore_color.rgb = BG_DARK

            txBox4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
            tf4 = txBox4.text_frame
            p4_title = tf4.paragraphs[0]
            p4_title.text = "推理链路"
            p4_title.font.size = Pt(28)
            p4_title.font.color.rgb = PRIMARY
            p4_title.font.bold = True

            for i, step in enumerate(ctx.reasoning_chain[:8], 1):
                y = Inches(1.2 + (i - 1) * 0.7)
                name = step.get("step", f"步骤{i}")
                desc = step.get("description", "")
                txStep = slide4.shapes.add_textbox(Inches(1), y, Inches(11), Inches(0.6))
                tfStep = txStep.text_frame
                pStep = tfStep.paragraphs[0]
                pStep.text = f"{i}. {name}"
                pStep.font.size = Pt(16)
                pStep.font.color.rgb = TEXT_LIGHT
                if desc:
                    pStep2 = tfStep.add_paragraph()
                    pStep2.text = f"   {desc}"
                    pStep2.font.size = Pt(12)
                    pStep2.font.color.rgb = TEXT_DIM

        # === Slide 5: 优化建议 (如果有) ===
        if ctx.optimization_suggestions:
            slide5 = prs.slides.add_slide(prs.slide_layouts[6])
            bg5 = slide5.background.fill
            bg5.solid()
            bg5.fore_color.rgb = BG_DARK

            txBox5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
            tf5 = txBox5.text_frame
            p5_title = tf5.paragraphs[0]
            p5_title.text = "优化建议"
            p5_title.font.size = Pt(28)
            p5_title.font.color.rgb = PRIMARY
            p5_title.font.bold = True

            txBox5b = slide5.shapes.add_textbox(Inches(1), Inches(1.3), Inches(11), Inches(5))
            tf5b = txBox5b.text_frame
            tf5b.word_wrap = True
            for j, sug in enumerate(ctx.optimization_suggestions):
                if j == 0:
                    p_sug = tf5b.paragraphs[0]
                else:
                    p_sug = tf5b.add_paragraph()
                p_sug.text = f"• {sug}"
                p_sug.font.size = Pt(16)
                p_sug.font.color.rgb = TEXT_LIGHT
                p_sug.space_after = Pt(12)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.getvalue()


# ==================== 7. DOCX 策略 ====================

class DocxOutputStrategy(OutputStrategy):
    """Word 文档输出策略"""
    format = OutputFormat.DOCX
    display_name = "Word文档"

    def render(self, ctx: RenderContext) -> bytes:
        from docx import Document
        from docx.shared import Inches as DocxInches, Pt as DocxPt, Cm
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.table import WD_TABLE_ALIGNMENT
        from docx.oxml.ns import qn, nsdecls
        from docx.oxml import parse_xml

        doc = Document()

        # 设置默认字体
        style = doc.styles["Normal"]
        font = style.font
        font.name = "Microsoft YaHei"
        font.size = DocxPt(11)
        # 中文字体回退
        style.element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")

        # 标题
        title = doc.add_heading("SageDispatch 八层处理管道", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # 副标题
        subtitle = doc.add_paragraph()
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = subtitle.add_run(f"{ctx.demand_type} | {ctx.domain} | {ctx.grade.upper()} 模式")
        run.font.size = DocxPt(12)
        run.font.color.rgb = None  # 默认颜色
        run.font.italic = True

        # 元信息表格
        doc.add_heading("基本信息", level=1)
        table = doc.add_table(rows=6, cols=2, style="Table Grid")
        table.alignment = WD_TABLE_ALIGNMENT.CENTER

        meta_rows = [
            ("需求类型", ctx.demand_type),
            ("领域分类", ctx.domain),
            ("处理等级", ctx.grade.upper()),
            ("决策方法", ctx.decision_method),
            ("置信度", f"{ctx.confidence:.1%}"),
            ("总耗时", f"{ctx.total_duration_ms:.1f} ms"),
        ]
        for i, (label, value) in enumerate(meta_rows):
            table.cell(i, 0).text = label
            table.cell(i, 1).text = value
            # 标签列加粗浅灰背景
            for p in table.cell(i, 0).paragraphs:
                for r in p.runs:
                    r.bold = True
            # 背景色
            shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="F0F0FF"/>')
            table.cell(i, 0)._tc.get_or_add_tcPr().append(shading)

        # 核心结论
        doc.add_heading("核心结论", level=1)
        doc.add_paragraph(ctx.final_answer or "处理完成。")

        # 推理链路
        if ctx.reasoning_chain:
            doc.add_heading("推理链路", level=1)
            chain_table = doc.add_table(rows=len(ctx.reasoning_chain) + 1, cols=3, style="Table Grid")
            chain_table.alignment = WD_TABLE_ALIGNMENT.CENTER
            # 表头
            headers = ["#", "步骤", "描述"]
            for j, h in enumerate(headers):
                chain_table.cell(0, j).text = h
                for p in chain_table.cell(0, j).paragraphs:
                    for r in p.runs:
                        r.bold = True
                shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="4338CA"/>')
                chain_table.cell(0, j)._tc.get_or_add_tcPr().append(shading)
                for p in chain_table.cell(0, j).paragraphs:
                    for r in p.runs:
                        r.font.color.rgb = None  # 白色字体

            for i, step in enumerate(ctx.reasoning_chain, 1):
                chain_table.cell(i, 0).text = str(i)
                chain_table.cell(i, 1).text = step.get("step", f"步骤{i}")
                chain_table.cell(i, 2).text = step.get("description", "")

        # 论证结果
        if ctx.argumentation_data:
            doc.add_heading("论证结果", level=1)
            r2 = ctx.argumentation_data.get("r2_argumentation", {})
            passed = ctx.argumentation_data.get("all_passed", True)
            status = "通过" if passed else "未通过"
            doc.add_paragraph(f"论证状态: {status} (SD-R2: {r2.get('grade', 'N/A')})")

            t2 = ctx.argumentation_data.get("t2_argumentation", {})
            if t2:
                t2_passed = t2.get("passed", False)
                doc.add_paragraph(
                    f"RefuteCore T2: {'通过' if t2_passed else '未通过'} "
                    f"({t2.get('strength_grade', '?')}级, {t2.get('risk_level', '?')}风险)"
                )

        # 优化建议
        if ctx.optimization_suggestions:
            doc.add_heading("优化建议", level=1)
            for sug in ctx.optimization_suggestions:
                doc.add_paragraph(sug, style="List Bullet")

        # 调度路径
        if ctx.routing_path:
            doc.add_heading("调度路径", level=1)
            doc.add_paragraph(" → ".join(ctx.routing_path))

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        return buf.getvalue()


# ==================== 格式检测器 ====================

class OutputFormatDetector:
    """
    格式检测器 — 根据需求类型和上下文自动选择最佳输出格式

    规则:
    1. 用户显式指定 → 直接使用
    2. 需求类型映射 → 按预设映射表
    3. 默认 → MARKDOWN
    """

    # 需求类型 → 推荐格式映射
    DEMAND_FORMAT_MAP: Dict[str, OutputFormat] = {
        "信息查询": OutputFormat.TEXT,
        "分析研究": OutputFormat.MARKDOWN,
        "战略规划": OutputFormat.DOCX,
        "决策选择": OutputFormat.HTML,
        "执行落地": OutputFormat.MARKDOWN,
        "创新突破": OutputFormat.PPTX,
        "论证反驳": OutputFormat.TEXT,
        "综合需求": OutputFormat.MARKDOWN,
    }

    # 备选格式（根据等级升级）
    GRADE_UPGRADE_MAP: Dict[str, Dict[str, OutputFormat]] = {
        "deep": {
            "分析研究": OutputFormat.HTML,
            "战略规划": OutputFormat.PDF,
            "综合需求": OutputFormat.HTML,
        },
        "super": {
            "分析研究": OutputFormat.PDF,
            "战略规划": OutputFormat.PPTX,
            "决策选择": OutputFormat.PPTX,
            "创新突破": OutputFormat.PDF,
            "综合需求": OutputFormat.PDF,
        },
    }

    @classmethod
    def detect(
        cls,
        demand_type: str = "",
        grade: str = "basic",
        explicit_format: Optional[OutputFormat] = None,
    ) -> OutputFormat:
        """检测最佳输出格式"""
        # 1. 显式指定优先
        if explicit_format is not None:
            return explicit_format

        # 2. 按需求类型映射
        base_format = cls.DEMAND_FORMAT_MAP.get(demand_type, OutputFormat.MARKDOWN)

        # 3. 按等级升级
        grade_upgrades = cls.GRADE_UPGRADE_MAP.get(grade, {})
        return grade_upgrades.get(demand_type, base_format)


# ==================== 输出引擎 ====================

class OutputEngine:
    """
    多模态输出引擎 — L7 输出层的核心执行器

    职责:
    - 格式检测与路由
    - 策略选择与渲染
    - 输出产物管理
    """

    # 格式 → 策略映射（延迟实例化）
    _STRATEGY_MAP: Dict[OutputFormat, type] = {
        OutputFormat.TEXT: TextOutputStrategy,
        OutputFormat.MARKDOWN: MarkdownOutputStrategy,
        OutputFormat.HTML: HtmlOutputStrategy,
        OutputFormat.IMAGE: ImageOutputStrategy,
        OutputFormat.PDF: PdfOutputStrategy,
        OutputFormat.PPTX: PptxOutputStrategy,
        OutputFormat.DOCX: DocxOutputStrategy,
    }

    def __init__(self, output_dir: str = ""):
        self._output_dir = output_dir
        self._strategies: Dict[OutputFormat, OutputStrategy] = {}
        self._stats = {
            "total_renders": 0,
            "by_format": {},
        }

    @property
    def output_dir(self) -> str:
        return self._output_dir

    @output_dir.setter
    def output_dir(self, value: str):
        self._output_dir = value

    def get_strategy(self, fmt: OutputFormat) -> OutputStrategy:
        """获取输出策略（单例缓存）"""
        if fmt not in self._strategies:
            strategy_cls = self._STRATEGY_MAP.get(fmt)
            if strategy_cls is None:
                raise ValueError(f"不支持的输出格式: {fmt}")
            self._strategies[fmt] = strategy_cls()
        return self._strategies[fmt]

    def detect_format(
        self,
        demand_type: str = "",
        grade: str = "basic",
        explicit_format: Optional[OutputFormat] = None,
    ) -> OutputFormat:
        """检测最佳输出格式（委托给 OutputFormatDetector）"""
        return OutputFormatDetector.detect(demand_type, grade, explicit_format)

    def render(
        self,
        fmt: OutputFormat,
        ctx: RenderContext,
    ) -> OutputArtifact:
        """
        渲染输出

        Args:
            fmt: 输出格式
            ctx: 渲染上下文

        Returns:
            OutputArtifact 输出产物
        """
        strategy = self.get_strategy(fmt)
        ctx.output_dir = self._output_dir
        artifact = strategy.render_to_file(ctx)

        # 更新统计
        self._stats["total_renders"] += 1
        fmt_name = fmt.value
        self._stats["by_format"][fmt_name] = self._stats["by_format"].get(fmt_name, 0) + 1

        return artifact

    def auto_render(self, ctx: RenderContext) -> OutputArtifact:
        """
        自动检测格式并渲染

        根据上下文中的 demand_type 和 grade 自动选择最佳输出格式
        """
        fmt = self.detect_format(
            demand_type=ctx.demand_type,
            grade=ctx.grade,
        )
        return self.render(fmt, ctx)

    def get_stats(self) -> Dict[str, Any]:
        """获取引擎统计"""
        return self._stats.copy()

    def supported_formats(self) -> List[str]:
        """列出支持的格式"""
        return [f.value for f in OutputFormat]

    def register_strategy(self, fmt: OutputFormat, strategy: OutputStrategy):
        """注册自定义输出策略"""
        self._strategies[fmt] = strategy


# ==================== 辅助函数 ====================

# PDF 中文字体注册缓存
_font_registered = False

def _register_chinese_font():
    """注册中文字体（仅执行一次）"""
    global _font_registered
    if _font_registered:
        return
    _font_registered = True

    try:
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont

        # Windows 常见中文字体路径
        font_paths = [
            ("SimHei", "C:/Windows/Fonts/simhei.ttf"),
            ("Microsoft YaHei", "C:/Windows/Fonts/msyh.ttc"),
            ("SimSun", "C:/Windows/Fonts/simsun.ttc"),
        ]
        for name, path in font_paths:
            if os.path.exists(path):
                pdfmetrics.registerFont(TTFont(name, path))
                break
    except Exception:
        pass  # 字体注册失败不影响 PDF 生成（英文回退）
