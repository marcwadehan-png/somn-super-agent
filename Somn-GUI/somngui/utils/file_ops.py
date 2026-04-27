# -*- coding: utf-8 -*-
"""Somn GUI - 文件操作工具

纯文件 I/O 工具函数，不依赖 Qt 或 Somn 核心。
提供多格式文件读取和内容提取能力。
"""

from __future__ import annotations

import hashlib
import os
import subprocess
import sys
from pathlib import Path
from typing import Optional

from loguru import logger

# 支持的文本扩展名
TEXT_EXTENSIONS = {
    '.txt', '.md', '.py', '.js', '.html', '.css',
    '.json', '.yaml', '.yml', '.xml',
}

# 支持的文档扩展名
DOC_EXTENSIONS = {
    '.docx', '.pdf', '.pptx', '.ppt', '.xlsx', '.xls', '.csv',
}

ALL_SUPPORTED = TEXT_EXTENSIONS | DOC_EXTENSIONS


def extract_file_content(file_path: Path | str) -> str:
    """提取文件内容为纯文本

    支持格式: txt, md, py, js, html, css, json, yaml, yml, xml,
              docx, pdf, pptx, xlsx, xls, csv

    Args:
        file_path: 文件路径

    Returns:
        提取的文本内容；若不支持则返回标记字符串
    """
    path = Path(file_path)
    suffix = path.suffix.lower()

    try:
        if suffix in TEXT_EXTENSIONS:
            return path.read_text(encoding='utf-8', errors='ignore')
        elif suffix == '.docx':
            return _extract_docx(path)
        elif suffix == '.pdf':
            return _extract_pdf(path)
        elif suffix in ('.pptx', '.ppt'):
            return _extract_pptx(path)
        elif suffix in ('.xlsx', '.xls', '.csv'):
            return _extract_excel(path, suffix)
        else:
            return f"[不支持的文件格式: {path.name}]"
    except Exception as e:
        logger.error(f"提取文件内容失败 {path}: {e}")
        return f"[无法读取: {path.name}]"


def _extract_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(path)
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            parts.append("\n[表格]")
            for row in table.rows:
                parts.append(" | ".join(cell.text for cell in row.cells))
        return '\n\n'.join(parts)
    except ImportError:
        return f"[Word文档: {path.name}，需安装 python-docx]"


def _extract_pdf(path: Path) -> str:
    try:
        import fitz
        doc = fitz.open(path)
        parts = []
        for i, page in enumerate(doc):
            parts.append(f"--- 第 {i + 1} 页 ---\n{page.get_text()}")
        doc.close()
        return '\n\n'.join(parts)
    except ImportError:
        return f"[PDF文档: {path.name}，需安装 PyMuPDF]"


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"=== 幻灯片 {i} ===")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
            parts.append("")
        return '\n'.join(parts)
    except ImportError:
        return f"[PPT: {path.name}，需安装 python-pptx]"


def _extract_excel(path: Path, suffix: str) -> str:
    try:
        import pandas as pd
        if suffix == '.csv':
            df = pd.read_csv(path)
            return df.to_string(index=False)
        else:
            xl = pd.ExcelFile(path)
            parts = []
            for sheet_name in xl.sheet_names:
                parts.append(f"=== 工作表: {sheet_name} ===")
                parts.append(pd.read_excel(path, sheet_name=sheet_name).to_string(index=False))
                parts.append("")
            return '\n'.join(parts)
    except ImportError:
        return f"[Excel: {path.name}，需安装 pandas]"


def open_with_system_default(file_path: Path | str) -> bool:
    """使用系统默认程序打开文件"""
    path = Path(file_path)
    try:
        if sys.platform == 'win32':
            os.startfile(path)
        elif sys.platform == 'darwin':
            subprocess.run(['open', str(path)], check=True)
        else:
            subprocess.run(['xdg-open', str(path)], check=True)
        return True
    except Exception as e:
        logger.error(f"系统打开失败: {e}")
        return False


def generate_file_id(file_path: Path | str) -> str:
    """生成文件唯一 ID（MD5 前 12 位）"""
    return hashlib.md5(str(file_path).encode()).hexdigest()[:12]


def scan_directory(dir_path: Path | str) -> dict:
    """扫描目录，返回文件统计"""
    path = Path(dir_path)
    if not path.exists():
        return {"total": 0, "files": 0, "dirs": 0, "by_ext": {}}

    files = 0
    dirs = 0
    by_ext: dict[str, int] = {}

    for item in path.rglob("*"):
        if item.is_file():
            files += 1
            ext = item.suffix.lower() or "(无扩展名)"
            by_ext[ext] = by_ext.get(ext, 0) + 1
        elif item.is_dir():
            dirs += 1

    return {
        "total": files + dirs,
        "files": files,
        "dirs": dirs,
        "by_ext": dict(sorted(by_ext.items(), key=lambda x: -x[1])),
    }
